import logging
import mimetypes

# Initialize logging early so it's available for optional dependencies
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

import time
import pandas as pd
import gc
import httpx
from pathlib import Path
import nltk.downloader
import pytesseract
import PyPDF2
from PyPDF2 import PdfReader
from PIL import Image, ImageOps, ImageFilter
from docx import Document as DOCXDocument
from pptx import Presentation
from project import settings
from pdfminer.high_level import extract_text as pdfminer_extract_text
try:
    import cv2
except ImportError:
    cv2 = None
    logger.warning("OpenCV (cv2) is not available; related features are disabled in this environment.")
import spacy
from langchain_openai import OpenAIEmbeddings # Use directly
# Import refactored providers and Langchain wrapper
from .llm_providers import (
    KNOWN_OLLAMA_EMBEDDING_DIMENSIONS,
    KNOWN_OPENAI_EMBEDDING_DIMENSIONS,
    OllamaProvider,
    OpenAIProvider,
    OllamaEmbeddings,
    BaseLLMProvider,
    LLMProviderError,
    LLMRateLimitError,
    LLMQuotaError,
    LLMAuthenticationError,
    LLMInvalidRequestError,
    LLMServiceUnavailableError,
    VectorStoreError,
    VectorStoreConnectionError,
)
from langchain.text_splitter import (
    RecursiveCharacterTextSplitter,
    MarkdownHeaderTextSplitter,
    NLTKTextSplitter,
)
from langchain_qdrant import QdrantVectorStore
from qdrant_client import QdrantClient
from qdrant_client.http.models import (
    Distance,
    VectorParams,
    MatchValue,
    MatchAny,
    Filter,
    FieldCondition,
    Datatype,
    PointStruct,
    PayloadSchemaType,
)
from langchain.schema import Document as LangChainDocument
import json
import xml.etree.ElementTree as ET
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough
from .search import HybridSearchRetriever
from concurrent.futures import ThreadPoolExecutor
from typing import List, Optional, Dict, Tuple, Any, Callable
from .models import (
    Document,
    LLMProviderConfig,
    DocumentAlert,
    Tenant,
    VectorStore,
    DocumentAccess,
    Assistant,
)
from django.contrib.auth import get_user_model
from knox.models import AuthToken
from django.shortcuts import get_object_or_404
import os
import subprocess
import platform
import speech_recognition as sr
from pydub import AudioSegment
import tempfile
from moviepy.editor import VideoFileClip
import whisper
import requests
try:
    from ddgs import DDGS  # DuckDuckGo Search client
except ImportError:  # pragma: no cover - optional dependency
    DDGS = None
from urllib.parse import urlparse
from functools import lru_cache
from collections import defaultdict
import re
import nltk
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords, wordnet
from nltk.stem import PorterStemmer
from nltk.sentiment import SentimentIntensityAnalyzer
import zipfile
import difflib
import numpy as np
import hashlib
import threading
from datetime import datetime
try:
    from langdetect import detect, LangDetectException
except Exception:  # pragma: no cover - fallback if package missing
    detect = lambda text: "unknown"
    class LangDetectException(Exception):
        pass
    logger.warning("langdetect not available; defaulting language metadata to 'unknown'.")
import tiktoken
from langchain.chains.llm import LLMChain
from .resilience import call_with_resilience, close_connections_before_io, CircuitBreakerOpenError
from bs4 import BeautifulSoup
import yaml
import configparser
from io import StringIO, BytesIO
from django.core.exceptions import ValidationError
from qdrant_client.http.exceptions import UnexpectedResponse
import openai
from contextlib import contextmanager
import base64
import io
import docx2txt
import chardet
import tarfile
import gzip
import rarfile
import py7zr
try:
    import magic
except ImportError:
    magic = None
    logger.warning("libmagic is not available; MIME type detection features are disabled in this environment.")
import uuid # For generating point IDs
import shutil
from django.utils import timezone
from langchain.docstore.document import Document as LangchainDocument
from .trace import trace_span

User = get_user_model()


# --- Provider Normalization Helpers ---

PROVIDER_EMBEDDING_DIMENSIONS = {
    "openai": 1536,
    "ollama": 1024,
}


def ensure_ollama_model_available(base_url: str, model: str, timeout: int = 30) -> bool:
    """Ensure a given Ollama model exists at base_url by checking /api/tags and pulling if missing.

    Returns True if the model is available (pre-existing or pulled), False otherwise.
    """
    import requests
    tags_url = f"{base_url.rstrip('/')}/api/tags"
    pull_url = f"{base_url.rstrip('/')}/api/pull"

    try:
        close_connections_before_io("Ollama tags request")
        resp = call_with_resilience(
            lambda: requests.get(tags_url, timeout=timeout),
            service="ollama_tags",
            exceptions=(requests.exceptions.RequestException,),
        )
        resp.raise_for_status()
        data = resp.json() if resp.content else {}
        models = {m.get("name") for m in data.get("models", []) if isinstance(m, dict)}
        if model in models:
            return True
    except Exception as exc:  # pragma: no cover - best-effort
        logger.warning("Failed to check Ollama tags at %s: %s", tags_url, exc)

    # Try to pull the model via Ollama HTTP API
    try:
        logger.info("Attempting to pull Ollama model '%s' from %s", model, pull_url)
        close_connections_before_io("Ollama pull request")
        pull_resp = call_with_resilience(
            lambda: requests.post(pull_url, json={"name": model}, timeout=timeout),
            service="ollama_pull",
            exceptions=(requests.exceptions.RequestException,),
        )
        pull_resp.raise_for_status()
        logger.info("Successfully pulled Ollama model '%s'", model)
        return True
    except Exception as exc:  # pragma: no cover - best-effort
        logger.warning("Failed to pull Ollama model '%s' via API: %s", model, exc)

    # As a fallback, attempt local CLI pull (may not be available on all hosts)
    try:
        logger.info("Attempting local 'ollama pull %s' as a fallback", model)
        subprocess.run(
            ["ollama", "pull", model],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            errors="ignore",
        )
        logger.info("Successfully pulled Ollama model '%s' via CLI fallback", model)
        return True
    except Exception as exc:  # pragma: no cover - best-effort
        logger.warning("CLI fallback failed to pull Ollama model '%s': %s", model, exc)

    return False

def normalize_provider_name(name: Optional[str]) -> str:
    """Return a canonical provider string for comparison and storage."""
    if not name:
        return "OpenAI"
    normalized = name.strip().lower()
    if normalized == "ollama":
        return "Ollama"
    if normalized == "claude":
        return "Claude"
    return "OpenAI"


def get_provider_embedding_dimension(provider: Optional[str]) -> Optional[int]:
    if not provider:
        return None
    return PROVIDER_EMBEDDING_DIMENSIONS.get(provider.lower())


def infer_provider_from_model(model: Optional[str], default: Optional[str] = None) -> str:
    """
    Best-effort provider inference from a model name. Falls back to the provided default
    (or OpenAI if not supplied) when the model name is ambiguous.
    """
    if not model:
        return normalize_provider_name(default)

    model_lc = model.strip().lower()

    # OpenAI model families (chat, completion, embeddings, audio, image)
    openai_prefixes = (
        "gpt-",
        "gpt",
        "gpt4",
        "gpt-4",
        "o1-",
        "text-",
        "dall-e",
        "whisper-",
        "text-embedding-",
        "embedding-",
    )
    openai_embeddings = {
        "text-embedding-3-small",
        "text-embedding-3-large",
        "text-embedding-ada-002",
    }
    if model_lc in openai_embeddings or model_lc.startswith(openai_prefixes):
        return "OpenAI"

    # Ollama models commonly use the `<model>:<tag>` form or known defaults
    ollama_prefixes = (
        "llama",
        "mistral",
        "phi",
        "codellama",
        "gemma",
        "nomic",
        "bge-",
        "all-minilm",
        "mxbai",
    )
    ollama_embeddings = {"nomic-embed-text", "bge-m3", "mxbai-embed-large", "bge-large-en"}
    if model_lc in ollama_embeddings:
        return "Ollama"

    if model_lc.startswith(ollama_prefixes):
        return "Ollama"

    if ":" in model or model_lc in {"gemma3", "gemma3:4b"}:
        return "Ollama"

    return normalize_provider_name(default)


def validate_model_availability(
    provider: Optional[str],
    model: Optional[str],
    api_key: Optional[str] = None,
    base_url: Optional[str] = None,
) -> None:
    """Validate that the requested model exists for the configured provider."""

    provider_name = normalize_provider_name(provider)
    if not model:
        raise ValueError("Model name is required for validation.")

    if provider_name == "OpenAI":
        effective_key = api_key or getattr(settings, "OPENAI_API_KEY", None)
        if not effective_key:
            raise ValueError("OpenAI API key is required to validate the model.")

        try:
            client = openai.OpenAI(api_key=effective_key)
            available_models = {getattr(m, "id", None) for m in getattr(client.models.list(), "data", [])}
            if model not in available_models:
                raise ValueError(
                    f"OpenAI model '{model}' not found or unavailable for the configured API key."
                )
        except openai.AuthenticationError as e:
            raise LLMAuthenticationError(f"OpenAI authentication error: {e.body.get('message') if e.body else str(e)}")
        except Exception as exc:
            raise LLMProviderError(f"Failed to validate OpenAI model '{model}': {exc}")

    elif provider_name == "Ollama":
        effective_base_url = base_url or getattr(
            settings, "OLLAMA_API_URL", getattr(settings, "OLLAMA_BASE_URL", "http://localhost:11434")
        )
        try:
            close_connections_before_io("Ollama tags availability check")
            response = call_with_resilience(
                lambda: requests.get(f"{effective_base_url.rstrip('/')}/api/tags", timeout=10),
                service="ollama_tags_check",
                exceptions=(requests.exceptions.RequestException,),
            )
            response.raise_for_status()
            models = response.json().get("models", [])
            if not any(model == entry.get("name") for entry in models):
                raise ValueError(f"Ollama model '{model}' is not downloaded or not found.")
        except requests.exceptions.HTTPError as e:
            if e.response.status_code in (401, 403):
                raise LLMAuthenticationError(f"Ollama authentication error during model validation: {e}")
            raise LLMProviderError(f"Ollama error {e.response.status_code} during model validation: {e}")
        except requests.exceptions.RequestException as exc:
            raise LLMServiceUnavailableError(f"Failed to connect to Ollama for model validation: {exc}")
        except Exception as exc:  # pragma: no cover - ensure consistent error surface
            raise LLMProviderError(f"Unexpected error during Ollama model validation: {exc}")

    else:
        raise ValueError(f"Unsupported provider for model validation: {provider}")


def get_default_embedding_dimension_for_provider(provider: Optional[str]) -> Optional[int]:
    """Return the expected embedding dimension for the given provider's default model."""
    normalized = normalize_provider_name(provider)
    if normalized == "Ollama":
        default_model = getattr(settings, "DEFAULT_OLLAMA_EMBEDDING_MODEL", "bge-m3")
        return KNOWN_OLLAMA_EMBEDDING_DIMENSIONS.get(default_model)
    if normalized == "OpenAI":
        default_model = getattr(settings, "DEFAULT_OPENAI_EMBEDDING_MODEL", "text-embedding-3-large")
        return KNOWN_OPENAI_EMBEDDING_DIMENSIONS.get(default_model)
    return None


def infer_provider_from_dimension(dimension: Optional[int]) -> Optional[str]:
    """Map a known embedding dimension back to a provider when possible."""
    if dimension is None:
        return None

    provider_dimensions = {
        get_default_embedding_dimension_for_provider("OpenAI"): "OpenAI",
        get_default_embedding_dimension_for_provider("Ollama"): "Ollama",
    }
    return provider_dimensions.get(dimension)


def ensure_default_assistant(user: Any) -> Optional[Assistant]:
    """Ensure the tenant has a default assistant aligned with the user's settings.

    This normalizes user preferences (provider, language), reuses an existing
    default assistant when present, and lazily creates one using the user's
    chosen provider and an available vector store when missing.
    """

    if not user or not getattr(user, "tenant", None):
        return None
    if not getattr(user, "llm_configured", False):
        return None

    existing_default = Assistant.objects.filter(tenant=user.tenant, is_default=True).first()
    if existing_default:
        return existing_default

    updates: dict[str, Any] = {}
    normalized_provider = normalize_provider_name(getattr(user, "selected_llm_provider", None))
    if getattr(user, "selected_llm_provider", None) != normalized_provider:
        updates["selected_llm_provider"] = normalized_provider

    default_language = getattr(settings, "DEFAULT_LANGUAGE", "en")
    if not getattr(user, "language", None):
        updates["language"] = default_language

    if updates:
        User.objects.filter(pk=user.pk).update(**updates)
        for field, value in updates.items():
            setattr(user, field, value)

    provider_default_model = (
        getattr(settings, "DEFAULT_OPENAI_MODEL", "gpt-4.1")
        if normalized_provider == "OpenAI"
        else getattr(settings, "DEFAULT_OLLAMA_MODEL", "gemma3:4b")
    )

    vector_store = (
        VectorStore.objects.filter(tenant=user.tenant, user=user).order_by("-created_at").first()
        or VectorStore.objects.filter(tenant=user.tenant).order_by("-created_at").first()
    )

    metadata: dict[str, Any] = {}
    if getattr(user, "language", None):
        metadata["language"] = user.language

    return Assistant.objects.create(
        tenant=user.tenant,
        name=f"Default Assistant ({normalized_provider})",
        vector_store=vector_store,
        instructions=None,
        model=provider_default_model,
        tools=[],
        metadata=metadata,
        is_default=True,
        creator=user,
    )


def detect_mime_type(file_path: str, default: str = "application/octet-stream") -> str:
    """Best-effort MIME type detection that tolerates missing libmagic."""
    if magic:
        try:
            mime = magic.from_file(file_path, mime=True)
            if mime:
                return mime
        except Exception as exc:  # pragma: no cover - defensive logging
            logger.warning("libmagic failed to detect MIME type for %s: %s", file_path, exc)

    guessed, _ = mimetypes.guess_type(file_path)
    return guessed or default


# Maximum bytes to read when we fall back to direct text ingestion for files that
# do not have a specific extractor. This keeps memory bounded while still
# retaining the original code/content structure instead of summarizing via LLM.
TEXT_FALLBACK_MAX_BYTES = 5 * 1024 * 1024  # 5MB


def _read_text_fallback(file_path: str, original_file_name: str) -> Optional[str]:
    """Best-effort direct text read for fallbackable types (code, logs, etc.).

    This is used before invoking LLM fallback to preserve the full code/content
    for text-like files, so chunking can operate on the source instead of an
    LLM-generated summary.
    """

    mime_type = detect_mime_type(file_path)
    ext = Path(file_path).suffix.lower()
    text_like_exts = {
        ".py",
        ".js",
        ".ts",
        ".java",
        ".c",
        ".cpp",
        ".cc",
        ".cs",
        ".go",
        ".rs",
        ".rb",
        ".php",
        ".sh",
        ".bash",
        ".zsh",
        ".ksh",
        ".ps1",
        ".psm1",
        ".r",
        ".pl",
        ".scala",
        ".swift",
        ".kt",
        ".kts",
        ".dart",
        ".lua",
        ".sql",
        ".html",
        ".htm",
        ".css",
        ".scss",
        ".sass",
        ".less",
        ".md",
        ".txt",
        ".log",
        ".json",
        ".yaml",
        ".yml",
        ".ini",
        ".cfg",
        ".toml",
        ".xml",
    }

    if mime_type.startswith("text/") or mime_type in {
        "application/json",
        "application/xml",
        "application/x-yaml",
        "application/javascript",
        "application/x-python",
    } or ext in text_like_exts:
        try:
            file_size = Path(file_path).stat().st_size
            if file_size > TEXT_FALLBACK_MAX_BYTES:
                logger.warning(
                    "Fallback text read for '%s' capped at %d bytes (size=%d)",
                    original_file_name,
                    TEXT_FALLBACK_MAX_BYTES,
                    file_size,
                )
            with open(file_path, "rb") as f:
                raw_bytes = f.read(TEXT_FALLBACK_MAX_BYTES)
            encoding = chardet.detect(raw_bytes).get("encoding") or "utf-8"
            text = raw_bytes.decode(encoding, errors="ignore")
            logger.info(
                "Fallback text ingestion succeeded for '%s' (len=%d, encoding=%s)",
                original_file_name,
                len(text),
                encoding,
            )
            return text
        except Exception as read_err:
            logger.warning(
                "Fallback direct text read failed for '%s': %s. Will continue with LLM fallback if available.",
                original_file_name,
                read_err,
                exc_info=True,
            )
    return None


# Configuration
MAX_FILE_SIZE = int(os.getenv('MAX_FILE_SIZE', 100 * 1024 * 1024))  # 100MB
CHUNK_SIZE = int(os.getenv('TEXT_CHUNK_SIZE', 1000))
CHUNK_OVERLAP = int(os.getenv('TEXT_CHUNK_OVERLAP', 200))
QDRANT_BATCH_SIZE = int(os.getenv('QDRANT_BATCH_SIZE', 64))
MAX_TOKEN_THRESHOLD = int(os.getenv('MAX_TOKEN_THRESHOLD', 4000))
HYBRID_SEARCH_TOP_K = int(os.getenv('HYBRID_SEARCH_TOP_K', 15))
MIN_SIMILARITY = float(os.getenv('MIN_SIMILARITY', 0.5))
ASK_TIMEOUT = float(os.getenv('ASK_TIMEOUT', 500.0))
MAX_QDRANT_ITERATIONS = int(os.getenv('MAX_QDRANT_ITERATIONS', 100))
DEFAULT_FALLBACK_DIMENSION = 1536 # Fallback dimension

# Constants for large file processing
LARGE_FILE_THRESHOLD = 10000  # rows for CSV/Excel
CHUNK_SIZE_LARGE = 1000  # rows per chunk for large files
MAX_SUMMARY_LENGTH = 2000  # characters per chunk summaryy


def log_document_processing_status(document_id: str, operation: str, details: dict = None):
    """Log detailed processing status for debugging missing rows"""
    timestamp = timezone.now().isoformat()
    log_entry = {
        "timestamp": timestamp,
        "document_id": document_id,
        "operation": operation,
        "details": details or {}
    }
    logger.info(f"PROCESSING_TRACKER: {json.dumps(log_entry)}")

# Lazy-loaded NLP models
_nlp = None
_sentiment_analyzer = None
_stemmer = None

def get_nlp():
    global _nlp
    if _nlp is None:
        try:
            _nlp = spacy.load("en_core_web_sm")
        except OSError:
            logger.info("Downloading spacy model en_core_web_sm...")
            spacy.cli.download("en_core_web_sm")
            _nlp = spacy.load("en_core_web_sm")
    return _nlp

def get_sentiment_analyzer():
    global _sentiment_analyzer
    if _sentiment_analyzer is None:
        try:
            # Ensure NLTK data is available before initializing
            download_nltk_data()
            _sentiment_analyzer = SentimentIntensityAnalyzer()
        except LookupError:
            logger.warning("NLTK VADER lexicon not available; sentiment analysis features are disabled.")
            _sentiment_analyzer = None
    return _sentiment_analyzer

def get_stemmer():
    global _stemmer
    if _stemmer is None:
        _stemmer = PorterStemmer()
    return _stemmer

# Download NLTK data
def download_nltk_data():
    nltk_data_path = os.path.join(nltk.data.path[0] if nltk.data.path else '.', 'nltk_data')
    os.makedirs(nltk_data_path, exist_ok=True)
    if nltk_data_path not in nltk.data.path:
        nltk.data.path.append(nltk_data_path)
    
    required_nltk_datasets = ['punkt', 'stopwords', 'wordnet', 'omw-1.4', 'vader_lexicon']
    for dataset in required_nltk_datasets:
        try:
            nltk.data.find(f'tokenizers/{dataset}' if dataset == 'punkt' else
                           f'corpora/{dataset}' if dataset in ['stopwords', 'wordnet', 'omw-1.4'] else
                           f'sentiment/{dataset}.zip')
        except LookupError:
            logger.info(f"Downloading NLTK dataset: {dataset}")
            nltk.download(dataset, download_dir=nltk_data_path, quiet=True)

# Timeout context manager
class TimeoutException(Exception):
    pass

@contextmanager
def timeout(seconds: float):
    if platform.system() == "Windows":
        # Windows doesn't support SIGALRM, use threading timer as fallback
        timer = threading.Timer(seconds, lambda: _raise_timeout())
        timer.start()
        try:
            yield
        finally:
            timer.cancel()
    else:
        import signal
        def signal_handler(signum, frame):
            raise TimeoutException("Operation timed out")
        current_handler = signal.getsignal(signal.SIGALRM)
        signal.signal(signal.SIGALRM, signal_handler)
        signal.alarm(int(seconds))
        try:
            yield
        finally:
            signal.alarm(0) # Disable the alarm
            signal.signal(signal.SIGALRM, current_handler) # Restore previous handler

def _raise_timeout():
    raise TimeoutException("Operation timed out via threading timer")

# --- Qdrant Client --- 
@lru_cache(maxsize=1)
def get_qdrant_client() -> QdrantClient:
    #Logger.info("Initializing Qdrant client...")
    close_connections_before_io("Qdrant client initialization")
    return QdrantClient(
        host=getattr(settings, 'QDRANT_HOST', 'localhost'),
        port=int(getattr(settings, 'QDRANT_PORT', 6333)), # Ensure port is int
        api_key=getattr(settings, 'QDRANT_API_KEY', None),
        timeout=float(getattr(settings, 'QDRANT_TIMEOUT', 20.0)), # Ensure timeout is float
        limits=httpx.Limits(
            max_connections=getattr(settings, 'QDRANT_MAX_CONNECTIONS', 50),
            max_keepalive_connections=getattr(settings, 'QDRANT_MAX_KEEPALIVE_CONNECTIONS', 20),
        ),
    )


def delete_qdrant_collection(collection_name: str) -> None:
    """Delete a Qdrant collection, ignoring missing collections."""

    logger.info("Deleting Qdrant collection '%s'", collection_name)
    close_connections_before_io("Qdrant collection deletion")
    client = get_qdrant_client()
    try:
        call_with_resilience(
            lambda: client.delete_collection(collection_name=collection_name),
            service="qdrant_delete_collection",
            exceptions=(Exception,),
        )
        logger.info("Deleted Qdrant collection '%s'", collection_name)
    except UnexpectedResponse as exc:
        if getattr(exc, "status_code", None) == 404:
            logger.info("Qdrant collection '%s' not found during deletion", collection_name)
            return
        logger.error("Qdrant failed to delete collection '%s': %s", collection_name, exc, exc_info=True)
        raise RuntimeError(
            f"Failed to delete Qdrant collection '{collection_name}': {exc.content.decode() if exc.content else str(exc)}"
        ) from exc
    except Exception as exc:
        logger.error("Unexpected error deleting Qdrant collection '%s': %s", collection_name, exc, exc_info=True)
        raise RuntimeError(f"Failed to delete Qdrant collection '{collection_name}': {exc}") from exc

# --- Authentication & LLM Provider Info (Refactored) --- 

@lru_cache(maxsize=256)
def get_llm_config(user_id: Optional[int] = None) -> Dict[str, Any]:
    """
    Determines the LLM configuration based on user settings or defaults.

    Returns a dictionary containing:
        - provider_name: 'OpenAI' or 'Ollama'
        - chat_model: Name of the chat model to use
        - embedding_model: Name of the embedding model to use
        - api_key: The API key (for OpenAI) or None
        - embedding_dimension: The dimension of the embedding model (fetched dynamically)
        - chat_config: Provider-specific overrides for chat generation (temperature, max_tokens, etc.)
    """
    config = {
        "provider_name": None,
        "chat_model": None,
        "embedding_model": None,
        "api_key": None,
        "embedding_dimension": None,
        "provider_instance": None, # Will hold OpenAIProvider or OllamaProvider instance
        "chat_config": {},
        "base_url": None,
    }

    user_key = None
    user = None
    selected_provider = normalize_provider_name(getattr(settings, "DEFAULT_LLM_PROVIDER", "OpenAI"))
    if user_id:
        try:
            user = User.objects.filter(id=user_id).first()
            if user and user.selected_llm_provider:
                selected_provider = normalize_provider_name(user.selected_llm_provider)

            user_key = LLMProviderConfig.objects.filter(
                user_id=user_id,
                is_active=True,
                is_valid=True,
                provider=selected_provider,
            ).first()
        except Exception as e:
            logger.error("Error fetching LLMProviderConfig for user %s: %s", user_id, e)
            #logger.error(f"Error fetching LLMProviderConfig for user {user_id}: {e}. Using defaults.")

    if user_key:
        #Logger.info(f"Using active configuration for user {user_id}: Provider={user_key.provider}, Model={user_key.model}")
        config["provider_name"] = user_key.provider
        config["chat_model"] = user_key.model # User's specified chat model
        config["chat_config"] = getattr(user_key, "chat_config", {}) or {}

        if config["provider_name"] == 'OpenAI':
            config["api_key"] = user_key.api_key
            # Check if user specified an embedding model (assuming field exists)
            config["embedding_model"] = getattr(user_key, 'embedding_model', None) or getattr(settings, 'DEFAULT_OPENAI_EMBEDDING_MODEL', 'text-embedding-3-large')
            if not config["api_key"]:
                #logger.error(f"User {user_id} selected OpenAI but provided no valid API key.")
                raise ValueError("OpenAI provider requires a valid API key.")
            config["provider_instance"] = OpenAIProvider()

        elif config["provider_name"] == 'Ollama':
            # Use user's specified chat model (already set)
            # Check if user specified an embedding model (assuming field exists)
            config["embedding_model"] = getattr(user_key, 'embedding_model', None) or getattr(settings, 'DEFAULT_OLLAMA_EMBEDDING_MODEL', 'bge-m3')
            config["api_key"] = None # Ollama doesn't use API key here
            config["base_url"] = getattr(user_key, "base_url", None) or getattr(settings, 'OLLAMA_API_URL', getattr(settings, 'OLLAMA_BASE_URL', 'http://localhost:11434'))
            config["provider_instance"] = OllamaProvider(base_url=config["base_url"])
        else:
            #logger.error(f"User {user_id} has invalid provider in key: {config['provider_name']}. Falling back to defaults.")
            user_key = None # Force fallback to defaults

    if not user_key: # Fallback to system defaults (Requirement 1)
        #Logger.info(f"No active user configuration found for user {user_id}. Using system defaults (OpenAI).")
        selected_provider = normalize_provider_name(selected_provider)
        preferred_provider = selected_provider or normalize_provider_name(getattr(settings, "DEFAULT_LLM_PROVIDER", "OpenAI"))
        config["provider_name"] = preferred_provider

        if config["provider_name"] == "Ollama":
            config["api_key"] = None
            config["chat_model"] = getattr(settings, 'DEFAULT_OLLAMA_MODEL', 'llama3.1:latest')
            config["embedding_model"] = getattr(settings, 'DEFAULT_OLLAMA_EMBEDDING_MODEL', 'bge-m3')
            config["base_url"] = getattr(settings, 'OLLAMA_API_URL', getattr(settings, 'OLLAMA_BASE_URL', 'http://localhost:11434'))
            config["provider_instance"] = OllamaProvider(base_url=config["base_url"])
        else:
            config["provider_name"] = "OpenAI" # Default provider is OpenAI
            config["api_key"] = getattr(settings, 'OPENAI_API_KEY', None)
            config["chat_model"] = getattr(settings, 'DEFAULT_OPENAI_MODEL', 'gpt-4.1')
            config["embedding_model"] = getattr(settings, 'DEFAULT_OPENAI_EMBEDDING_MODEL', 'text-embedding-3-large')
            config["chat_config"] = {}
            if not config["api_key"]:
                #logger.error("Defaulting to OpenAI, but no default OPENAI_API_KEY found in settings.")
                raise ValueError("Default OpenAI provider requires an OPENAI_API_KEY in settings.")
            config["provider_instance"] = OpenAIProvider()

    if config["provider_instance"]:
        # Provide user-specific chat overrides to the provider so they are merged with defaults per request
        config["provider_instance"].chat_config = config.get("chat_config", {}) or {}

    # Dynamically get embedding dimension AFTER provider and model are determined
    try:
        if not config["provider_instance"]:
            logger.error("Provider instance was not initialized before attempting to get embedding dimension.")
            raise ValueError("LLM provider instance could not be determined. Cannot fetch embedding dimension.")

        if config["provider_name"] == "Ollama":
            if not ensure_ollama_model_available(config["base_url"], config["chat_model"]):
                raise RuntimeError(
                    f"Ollama model '{config['chat_model']}' is unavailable at {config['base_url']}. "
                    "Pull the model on the target host and retry."
                )

        logger.info(f"Fetching embedding dimension for provider {config['provider_name']}, model: {config['embedding_model']}")
        embedding_dimension = config["provider_instance"].get_embedding_dimension(
            model=config["embedding_model"],
            api_key=config["api_key"]
        )
        if embedding_dimension is None: # get_embedding_dimension should raise error, but as a safeguard
            raise ValueError(f"Embedding dimension could not be determined for model {config['embedding_model']}.")

        config["embedding_dimension"] = embedding_dimension
        logger.info(f"Successfully determined embedding dimension: {embedding_dimension}")
        return config
    except Exception as e:
        logger.error(f"Failed to get embedding dimension: {e}")
        raise

# --- Deprecated Functions (kept for reference, remove later) ---
# @lru_cache(maxsize=128)
# def get_openai_api_key(user_id: int) -> Optional[str]: ... # Logic moved into get_llm_config
# @lru_cache(maxsize=128)
# def get_llm_provider_instance(user_id: Optional[int] = None) -> Tuple[BaseLLMProvider, str, Optional[str]]: ... # Logic moved into get_llm_config
# def get_embedding_provider_instance(user_id: Optional[int] = None) -> Tuple[BaseLLMProvider, str, Optional[str]]: ... # Logic moved into get_llm_config

# --- Authentication --- 
def get_authenticated_user(token: str) -> Any:
    #logger.debug(f"Authenticating token (last 4 chars): ...{token[-4:]}")
    try:
        # Use select_related to fetch user and tenant in one query
        auth_token = get_object_or_404(AuthToken.objects.select_related('user__tenant'), token_key=token)
        user = auth_token.user
        if not hasattr(user, 'tenant') or user.tenant is None:
            #logger.error(f"User {user.username} (ID: {user.id}) has no associated tenant.")
            raise ValidationError("User authentication successful, but user has no associated tenant.")
        #Logger.info(f"Authenticated user: {user.username} (ID: {user.id}), Tenant: {user.tenant.name}")
        return user
    except AuthToken.DoesNotExist:
        #logger.error(f"Authentication failed: Token not found (...{token[-4:]})")
        raise ValidationError("Authentication failed: Invalid token.")
    except Exception as e:
        #logger.error(f"Error authenticating user with token (...{token[-4:]}): {e}")
        raise ValidationError(f"Authentication failed: {str(e)}")


def sync_user_login_state(user: Any) -> Dict[str, Any]:
    """
    Normalize the user's provider, setup state, and collection flags during login.

    This aligns selected_llm_provider with the user's active configuration or the
    existing collection dimension to avoid mismatches (e.g., OpenAI selected while
    the tenant collection was created with an Ollama-sized embedding dimension).
    """
    if not user or not getattr(user, "tenant", None):
        return {"warnings": ["User has no tenant association; provider state not synced."]}

    warnings: list[str] = []
    updates: dict[str, Any] = {}

    active_collection = getattr(user, "active_collection", None)
    if not getattr(user, "llm_configured", False) and not active_collection:
        return {
            "warnings": ["LLM provider not configured; setup required before using APIs."],
            "collection_exists": False,
            "collection_dimension": None,
            "selected_provider": None,
        }

    active_config = (
        LLMProviderConfig.objects.filter(user=user, is_valid=True, is_active=True).first()
        or LLMProviderConfig.objects.filter(user=user, is_valid=True).order_by("-updated_at").first()
    )
    config_provider = normalize_provider_name(active_config.provider) if active_config else None

    collection_exists = False
    collection_dimension = None
    collection_provider = None
    if active_collection:
        collection_exists, collection_dimension = get_collection_status(
            active_collection.qdrant_collection_name
        )
        collection_provider = infer_provider_from_dimension(collection_dimension)
        if collection_exists and not getattr(user, "active_collection_ready", False):
            updates["active_collection_ready"] = True
            updates["active_collection"] = active_collection
        if collection_exists and not getattr(user, "llm_configured", False):
            updates["llm_configured"] = True

    desired_provider = config_provider or normalize_provider_name(
        getattr(user, "selected_llm_provider", getattr(settings, "DEFAULT_LLM_PROVIDER", "OpenAI"))
    )

    if collection_provider and desired_provider != collection_provider:
        warnings.append(
            f"Selected provider '{desired_provider}' realigned to '{collection_provider}' "
            f"based on existing collection dimension {collection_dimension}."
        )
        desired_provider = collection_provider

    if getattr(user, "selected_llm_provider", None) != desired_provider:
        updates["selected_llm_provider"] = desired_provider

    has_valid_config = LLMProviderConfig.objects.filter(
        user=user, provider=desired_provider, is_valid=True
    ).exists()
    if has_valid_config and not getattr(user, "is_setup", False):
        updates["is_setup"] = True
    elif not has_valid_config and getattr(user, "is_setup", False):
        warnings.append(
            f"No valid {desired_provider} configuration found; setup state left unchanged."
        )

    if updates:
        User.objects.filter(pk=user.pk).update(**updates)
        for field, value in updates.items():
            setattr(user, field, value)

    return {
        "collection_exists": collection_exists,
        "collection_dimension": collection_dimension,
        "selected_provider": desired_provider,
        "warnings": warnings,
    }


# --- Qdrant Collection and Vector Store Initialization (Updated) ---

DEFAULT_PAYLOAD_INDEXES = {
    "document_id": "keyword",
    "vector_store_id": "keyword",
    "tenant_id": "keyword",
    "user_id": "keyword",
    "language": "keyword",
    "page_number": "integer",
}


def ensure_payload_indexes(client: QdrantClient, collection_name: str, schema: Dict[str, str]) -> None:
    """Ensure payload indexes exist for specified fields."""
    try:
        info = client.get_collection(collection_name)
        existing = info.payload_schema or {}
    except Exception as e:
        logger.warning(
            f"Failed to fetch payload schema for collection '{collection_name}': {e}"
        )
        existing = {}

    for field, field_type in schema.items():
        if field in existing:
            continue
        try:
            client.create_payload_index(
                collection_name=collection_name,
                field_name=field,
                field_schema=PayloadSchemaType(field_type),
            )
            logger.info(
                f"Created payload index for field '{field}' ({field_type}) in collection '{collection_name}'."
            )
        except Exception as idx_e:
            logger.warning(
                f"Failed to create payload index for field '{field}': {idx_e}"
            )


def extract_collection_dimension(collection_info: Any) -> Optional[int]:
    """Safely extract the vector dimension from a Qdrant collection response."""
    try:
        if (
            hasattr(collection_info, "config")
            and collection_info.config
            and hasattr(collection_info.config, "params")
            and collection_info.config.params
        ):
            vectors = collection_info.config.params.vectors
        else:
            return None

        if isinstance(vectors, dict):
            if "size" in vectors:
                return vectors["size"]
            default_vector = vectors.get("")
            if default_vector and hasattr(default_vector, "size"):
                return default_vector.size
        elif hasattr(vectors, "size"):
            return vectors.size
    except Exception as exc:  # pragma: no cover - defensive logging
        logger.warning("Failed to extract dimension from collection info: %s", exc)
    return None


def get_collection_status(collection_name: str) -> Tuple[bool, Optional[int]]:
    """Return whether the collection exists and its dimension (if readable)."""
    close_connections_before_io("Qdrant collection status check")
    client = get_qdrant_client()
    try:
        info = call_with_resilience(
            lambda: client.get_collection(collection_name=collection_name),
            service="qdrant_get_collection_status",
            exceptions=(Exception,),
        )
        return True, extract_collection_dimension(info)
    except UnexpectedResponse as exc:
        if getattr(exc, "status_code", None) == 404:
            return False, None
        logger.warning(
            "Unexpected Qdrant response for collection '%s': %s",
            collection_name,
            exc,
        )
    except Exception as exc:  # pragma: no cover - defensive logging
        logger.warning(
            "Failed to check collection '%s' status: %s", collection_name, exc, exc_info=True
        )
    return False, None


def initialize_qdrant_collection(collection_name: str, vector_dimension: int) -> QdrantClient:
    close_connections_before_io("Qdrant collection initialization")
    client = get_qdrant_client()
    logger.info(f"Initializing Qdrant collection: '{collection_name}' with expected dimension: {vector_dimension}")
    try:
        logger.debug(f"Checking if collection '{collection_name}' exists...")
        collection_info = call_with_resilience(
            lambda: client.get_collection(collection_name=collection_name),
            service="qdrant_get_collection",
            exceptions=(Exception,),
        )
        
        # Validate collection parameters, especially vector dimension
        if (hasattr(collection_info, 'config') and collection_info.config and
            hasattr(collection_info.config, 'params') and collection_info.config.params and
            hasattr(collection_info.config.params, 'vectors') and
            isinstance(collection_info.config.params.vectors, dict) and # If single vector config
            'size' in collection_info.config.params.vectors):
            existing_dimension = collection_info.config.params.vectors['size']
        elif (hasattr(collection_info, 'config') and collection_info.config and
              hasattr(collection_info.config, 'params') and collection_info.config.params and
              hasattr(collection_info.config.params, 'vectors') and # if multiple named vectors
              isinstance(collection_info.config.params.vectors, dict) and
              "" in collection_info.config.params.vectors and # Check for default unnamed vector
              hasattr(collection_info.config.params.vectors[""], 'size')):
            existing_dimension = collection_info.config.params.vectors[""].size
        elif (hasattr(collection_info, 'config') and collection_info.config and
              hasattr(collection_info.config, 'params') and collection_info.config.params and
              hasattr(collection_info.config.params, 'vectors') and # if single vector config (older client?)
              hasattr(collection_info.config.params.vectors, 'size')):
             existing_dimension = collection_info.config.params.vectors.size
        else:
            logger.error(f"Collection '{collection_name}' exists but its vector configuration is in an unrecognized format or missing. Params: {collection_info.config.params if hasattr(collection_info.config, 'params') else 'N/A'}")
            # This case is problematic. Recreating might be an option but is risky.
            # For now, raise an error indicating manual check is needed.
            raise RuntimeError(f"Collection '{collection_name}' has an unreadable or incomplete vector configuration. Manual review required.")

        if existing_dimension != vector_dimension:
            error_message = (
                f"Critical Dimension Mismatch for Qdrant collection '{collection_name}'. "
                f"Existing dimension: {existing_dimension}, Expected dimension: {vector_dimension}. "
                "Updating the dimension of an existing collection with data is not supported directly by Qdrant "
                "and can lead to data corruption or errors. Manual intervention is required. "
                "You may need to re-index data or create a new collection with the correct dimension."
            )
            logger.error(error_message)
            # Option 1 (Safer): Raise a ValueError
            raise ValueError(error_message)
            # Option 2 (More Aggressive - NOT IMPLEMENTED as per instruction):
            # logger.warning(f"Attempting to delete and recreate collection '{collection_name}' due to dimension mismatch. THIS IS A DESTRUCTIVE OPERATION.")
            # client.delete_collection(collection_name=collection_name)
            # logger.info(f"Collection '{collection_name}' deleted.")
            # client.create_collection(
            #     collection_name=collection_name,
            #     vectors_config=VectorParams(size=vector_dimension, distance=Distance.COSINE)
            # )
            # logger.info(f"Collection '{collection_name}' recreated with new dimension {vector_dimension}.")
        else:
            logger.info(f"Collection '{collection_name}' already exists with correct dimension {vector_dimension}.")
            ensure_payload_indexes(client, collection_name, DEFAULT_PAYLOAD_INDEXES)

    except UnexpectedResponse as e:
        if hasattr(e, 'status_code') and e.status_code == 404:
            logger.info(f"Collection '{collection_name}' not found. Creating with dimension {vector_dimension}.")
            try:
                call_with_resilience(
                    lambda: client.create_collection(
                        collection_name=collection_name,
                        vectors_config=VectorParams(size=vector_dimension, distance=Distance.COSINE),
                    ),
                    service="qdrant_create_collection",
                    exceptions=(Exception,),
                )
                logger.info(f"Collection '{collection_name}' created successfully with dimension {vector_dimension}.")
                ensure_payload_indexes(client, collection_name, DEFAULT_PAYLOAD_INDEXES)
            except UnexpectedResponse as create_e_qdrant:
                logger.error(f"Qdrant specific error during creation of collection '{collection_name}': {create_e_qdrant.content.decode() if create_e_qdrant.content else str(create_e_qdrant)}")
                raise RuntimeError(f"Failed to create Qdrant collection '{collection_name}' due to Qdrant error: {create_e_qdrant.content.decode() if create_e_qdrant.content else str(create_e_qdrant)}")
            except Exception as create_e_generic:
                logger.error(f"Generic error during creation of collection '{collection_name}': {create_e_generic}")
                raise RuntimeError(f"Failed to create Qdrant collection '{collection_name}': {str(create_e_generic)}")
        else:
            logger.error(f"Unexpected Qdrant response while checking collection '{collection_name}': {e.status_code} - {e.content.decode() if e.content else str(e)}")
            raise RuntimeError(f"Failed to check Qdrant collection '{collection_name}' status: {e.status_code} - {e.content.decode() if e.content else str(e)}")
    except ValueError as ve: # Catch ValueError from dimension mismatch
        raise ve # Re-raise it to be caught by the caller
    except Exception as e:
        logger.error(f"An unexpected error occurred while initializing Qdrant collection '{collection_name}': {e}", exc_info=True)
        # Check for connection errors string matching since generic Exception catches everything
        msg = str(e).lower()
        if "connection refused" in msg or "active refused" in msg or "failed to connect" in msg or "10061" in msg:
             host = getattr(settings, 'QDRANT_HOST', 'localhost')
             port = getattr(settings, 'QDRANT_PORT', 6333)
             raise VectorStoreConnectionError(f"Could not connect to Qdrant vector store at {host}:{port}. Is the service running?")
        
        raise RuntimeError(f"Failed to initialize Qdrant collection '{collection_name}': {str(e)}")
    return client

def get_qdrant_vector_store(user: Any, collection_name: str) -> QdrantVectorStore:
    logger.info(f"Getting Qdrant vector store for user {user.id if user else 'system'}, collection '{collection_name}'")
    try:
        close_connections_before_io("Qdrant vector store initialization")
        # 1. Get the unified LLM config (this also determines vector_dimension)
        llm_config = get_llm_config(user.id if user else None)
        provider_name = llm_config["provider_name"]
        embedding_model = llm_config["embedding_model"]
        api_key = llm_config["api_key"] # For OpenAI
        vector_dimension = llm_config["embedding_dimension"]

        if not vector_dimension:
             # This error should ideally be caught by get_llm_config, but as a safeguard:
             logger.error(f"Vector dimension could not be determined for collection '{collection_name}'.")
             raise ValueError("Failed to determine vector dimension, cannot initialize Qdrant vector store.")

        # 2. Initialize Qdrant collection, ensuring it exists with the correct dimension.
        # This will raise an error if there's an unrecoverable issue (e.g., dimension mismatch).
        client = initialize_qdrant_collection(collection_name, vector_dimension)

        # 3. Instantiate the appropriate Langchain Embeddings class
        logger.debug(f"Instantiating embeddings for provider: {provider_name}, model: {embedding_model}")
        if provider_name == 'OpenAI':
            embeddings = OpenAIEmbeddings(
                api_key=api_key,
                model=embedding_model,
                dimensions=vector_dimension
            )
        elif provider_name == 'Ollama':
            embeddings = OllamaEmbeddings(model=embedding_model)
        else:
            raise ValueError(f"Unsupported provider name for embeddings: {provider_name}")

        # 4. Create the Langchain QdrantVectorStore wrapper
        vector_store = QdrantVectorStore(
            client=client,
            collection_name=collection_name,
            embedding=embeddings
        )
        #Logger.info(f"QdrantVectorStore initialized successfully for collection '{collection_name}' using {provider_name} embeddings ({embedding_model}).")
        return vector_store

    except Exception as e:
        logger.error(f"Error initializing Qdrant vector store for user {user.id if user else 'system'}, collection '{collection_name}': {e}", exc_info=True)
        # Wrap exception for clarity, or re-raise if already specific enough
        if isinstance(e, (ValueError, RuntimeError, ValidationError, LLMProviderError, VectorStoreError)):
            raise e
        
        # Check for connection errors string matching here too as a fallback for other parts
        msg = str(e).lower()
        if "connection refused" in msg or "active refused" in msg or "10061" in msg:
             raise VectorStoreConnectionError(f"Connection to vector store failed: {str(e)}")

        raise LLMProviderError(str(e))
    

def is_large_file_by_row_count(file_path: str, ext: str) -> bool:
    """Check if CSV/Excel file exceeds threshold rows"""
    try:
        if ext == '.csv':
            # Quick row count for CSV
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                row_count = sum(1 for _ in f)
                # Subtract 1 for header if you want to be precise, but for threshold check it's fine
        elif ext in ['.xlsx', '.xls']:
            # Quick row count for Excel
            import pandas as pd
            xl = pd.ExcelFile(file_path)
            row_count = 0
            for sheet_name in xl.sheet_names:
                try:
                    # Read just first few rows to get shape estimate or nrows=1
                    sheet_info = pd.read_excel(file_path, sheet_name=sheet_name, nrows=1)
                    # This reads 1 row but gives us the column structure.
                    # Getting exact count without loading full sheet is tricky.
                    # Let's use a more robust approach for Excel.
                    # Alternative: Use openpyxl/xlrd for row counting without loading data.
                    # For now, a simple approach: read first 1000 rows to check if it's large enough
                    # This is a compromise between speed and accuracy.
                    # A better way might be to use the engine directly.
                    # Let's try a different approach:
                    if ext == '.xlsx':
                        import openpyxl
                        wb = openpyxl.load_workbook(file_path, read_only=True)
                        sheet_row_count = wb[sheet_name].max_row
                        wb.close()
                    else: # .xls
                        import xlrd
                        wb = xlrd.open_workbook(file_path)
                        sheet_row_count = wb.sheet_by_name(sheet_name).nrows
                    row_count += sheet_row_count
                except Exception as e:
                    logger.warning(f"Could not get row count for sheet '{sheet_name}' in '{file_path}': {e}")
                    # Continue checking other sheets
                    continue
        else:
            return False

        result = row_count > LARGE_FILE_THRESHOLD
        logger.info(f"File '{file_path}' row count: {row_count}, Threshold: {LARGE_FILE_THRESHOLD}, Is Large: {result}")
        return result
    except Exception as e:
        logger.warning(f"Could not determine row count for large file check for '{file_path}': {e}")
        return False # Default to not large if check fails


def process_large_file_streaming(file_path: str, ext: str, original_file_name: str, user: Any = None) -> tuple:
    """
    Process large CSV/Excel files in chunks to avoid memory issues
    """
    try:
        logger.info(f"Starting streaming processing for large file: {original_file_name}")
        close_connections_before_io("large file streaming")
        
        if ext == '.csv':
            return process_large_csv_streaming(file_path, original_file_name, user)
        elif ext in ['.xlsx', '.xls']:
            return process_large_excel_streaming(file_path, original_file_name, user)
        else:
            # Fallback to regular processing
            logger.warning(f"Large file processing called for unsupported extension '{ext}', falling back to standard extraction.")
            extracted_data = extract_text_from_file(file_path, original_file_name, user)
            summary = generate_summary_from_text(extracted_data, user)
            return extracted_data, summary
            
    except Exception as e:
        logger.error(f"Error in streaming processing for '{original_file_name}': {e}", exc_info=True)
        raise RuntimeError(f"Failed to process large file '{original_file_name}': {e}")

def _append_preview(parts: List[str], chunk_text: str, max_chars: int) -> None:
    """Append chunk text up to max_chars total to avoid unbounded memory growth."""
    if not chunk_text or max_chars <= 0:
        return
    current_len = sum(len(part) for part in parts)
    if current_len >= max_chars:
        return
    remaining = max_chars - current_len
    parts.append(chunk_text[:remaining])


def _format_rows_to_text(sheet_name: str, rows: List[Any]) -> str:
    header = f"Sheet: {sheet_name}\n"
    row_lines = []
    for row in rows:
        row_lines.append(" | ".join("" if cell is None else str(cell) for cell in row))
    return header + "\n".join(row_lines)


def _summarize_chunk_if_needed(
    summaries: List[str],
    chunk_text: str,
    user: Any,
    sheet_name: str,
    chunk_index: int,
) -> List[str]:
    if len(chunk_text.strip()) > 100 and len(summaries) < 10:
        try:
            chunk_summary = generate_summary_from_text(chunk_text[:2000], user)
            summaries.append(chunk_summary)
            logger.debug(
                "Generated summary for sheet '%s' chunk %s, length: %s",
                sheet_name,
                chunk_index,
                len(chunk_summary),
            )
        except Exception as summary_e:
            logger.warning(
                "Failed to generate summary for Excel sheet '%s' chunk %s: %s",
                sheet_name,
                chunk_index,
                summary_e,
            )
            summaries.append(f"[Summary generation failed for sheet '{sheet_name}' chunk {chunk_index}]")
    return summaries

def process_large_csv_streaming(file_path: str, original_file_name: str, user: Any = None) -> tuple:
    """
    Process large CSV files in chunks with hierarchical summarization
    """
    summaries = []
    content_preview_parts = []
    chunk_count = 0
    total_rows = 0
    max_preview_chars = getattr(settings, "MAX_LARGE_FILE_CONTENT_CHARS", 200000)
    
    try:
        logger.info(f"Starting CSV streaming processing for: {original_file_name}")
        # Process CSV in chunks
        for chunk in pd.read_csv(file_path, chunksize=CHUNK_SIZE_LARGE):
            chunk_count += 1
            chunk_rows = len(chunk)
            total_rows += chunk_rows
            logger.info(f"Processing CSV chunk {chunk_count} with {chunk_rows} rows (Total so far: {total_rows})")
            
            # Convert chunk to text (store only a bounded preview to avoid memory growth)
            chunk_text = chunk.to_string(index=False, max_colwidth=50)
            _append_preview(content_preview_parts, chunk_text, max_preview_chars)
            
            # Generate summary for this chunk (limit to 10 chunks for performance/cost)
            if len(chunk_text.strip()) > 100 and len(summaries) < 10:
                try:
                    chunk_summary = generate_summary_from_text(chunk_text[:2000], user)  # Limit size for LLM call
                    summaries.append(chunk_summary)
                    logger.debug(f"Generated summary for chunk {chunk_count}, length: {len(chunk_summary)}")
                except Exception as summary_e:
                    logger.warning(f"Failed to generate summary for CSV chunk {chunk_count}: {summary_e}")
                    # Add a placeholder summary
                    summaries.append(f"[Summary generation failed for chunk {chunk_count}]")
            
            # Memory management - force garbage collection
            if chunk_count % 20 == 0:  # More frequent cleanup
                gc.collect()
                logger.debug(f"Performed garbage collection after chunk {chunk_count}")
            del chunk

        # Combine preview content for storage/summary purposes only
        full_content = "".join(content_preview_parts)
        logger.info(f"Completed processing {chunk_count} CSV chunks, total rows: {total_rows}")
        
        # Generate final summary using hierarchical approach
        final_summary = generate_hierarchical_summary(summaries, user)
        logger.info(f"Generated final hierarchical summary for CSV file '{original_file_name}'")

        return full_content, final_summary
    
    except Exception as e:
        logger.error(f"Error processing large CSV '{original_file_name}' after {chunk_count} chunks: {e}", exc_info=True)
        # Fallback to basic extraction if streaming fails
        logger.info("Falling back to standard CSV extraction due to streaming error.")
        extracted_data = extract_text_from_file(file_path, original_file_name, user)
        summary = generate_summary_from_text(extracted_data[:4000], user)  # Limit size
        return extracted_data, summary
    
def process_large_excel_streaming(file_path: str, original_file_name: str, user: Any = None) -> tuple:
    """
    Process large Excel files in chunks with hierarchical summarization
    """
    summaries = []
    content_preview_parts = []
    total_chunks = 0
    total_rows = 0
    max_preview_chars = getattr(settings, "MAX_LARGE_FILE_CONTENT_CHARS", 200000)
    
    try:
        logger.info(f"Starting Excel streaming processing for: {original_file_name}")
        if file_path.lower().endswith(".xlsx"):
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            sheet_names = wb.sheetnames
            for sheet_name in sheet_names:
                sheet = wb[sheet_name]
                logger.info("Processing Excel sheet: %s", sheet_name)
                chunk_rows = []
                sheet_chunk_count = 0
                for row in sheet.iter_rows(values_only=True):
                    chunk_rows.append(row)
                    if len(chunk_rows) >= CHUNK_SIZE_LARGE:
                        sheet_chunk_count += 1
                        total_chunks += 1
                        total_rows += len(chunk_rows)
                        chunk_text = _format_rows_to_text(sheet_name, chunk_rows)
                        _append_preview(content_preview_parts, chunk_text, max_preview_chars)
                        summaries = _summarize_chunk_if_needed(
                            summaries, chunk_text, user, sheet_name, sheet_chunk_count
                        )
                        chunk_rows.clear()
                        if total_chunks % 15 == 0:
                            gc.collect()
                            logger.debug("Performed garbage collection after chunk %s", total_chunks)
                if chunk_rows:
                    sheet_chunk_count += 1
                    total_chunks += 1
                    total_rows += len(chunk_rows)
                    chunk_text = _format_rows_to_text(sheet_name, chunk_rows)
                    _append_preview(content_preview_parts, chunk_text, max_preview_chars)
                    summaries = _summarize_chunk_if_needed(
                        summaries, chunk_text, user, sheet_name, sheet_chunk_count
                    )
                    chunk_rows.clear()
            wb.close()
        else:
            import xlrd
            wb = xlrd.open_workbook(file_path, on_demand=True)
            for sheet in wb.sheets():
                sheet_name = sheet.name
                logger.info("Processing Excel sheet: %s", sheet_name)
                sheet_chunk_count = 0
                for start_idx in range(0, sheet.nrows, CHUNK_SIZE_LARGE):
                    end_idx = min(start_idx + CHUNK_SIZE_LARGE, sheet.nrows)
                    chunk_rows = [sheet.row_values(i) for i in range(start_idx, end_idx)]
                    sheet_chunk_count += 1
                    total_chunks += 1
                    total_rows += len(chunk_rows)
                    chunk_text = _format_rows_to_text(sheet_name, chunk_rows)
                    _append_preview(content_preview_parts, chunk_text, max_preview_chars)
                    summaries = _summarize_chunk_if_needed(
                        summaries, chunk_text, user, sheet_name, sheet_chunk_count
                    )
                    if total_chunks % 15 == 0:
                        gc.collect()
                        logger.debug("Performed garbage collection after chunk %s", total_chunks)
            wb.release_resources()

        full_content = "".join(content_preview_parts)
        logger.info(
            "Completed processing Excel file '%s' with %s chunks, total rows: %s",
            original_file_name,
            total_chunks,
            total_rows,
        )

        final_summary = generate_hierarchical_summary(summaries, user)
        logger.info(f"Generated final hierarchical summary for Excel file '{original_file_name}'")

        return full_content, final_summary
    
    except Exception as e:
        logger.error(f"Error processing large Excel '{original_file_name}' after {total_chunks} chunks: {e}", exc_info=True)
        logger.info("Falling back to standard Excel extraction due to streaming error.")
        extracted_data = extract_text_from_file(file_path, original_file_name, user)
        summary = generate_summary_from_text(extracted_data[:4000], user)
        return extracted_data, summary

            
def generate_hierarchical_summary(summaries: list, user: Any = None) -> str:
    """
    Generate a final summary from multiple chunk summaries using hierarchical approach
    """
    if not summaries:
        logger.info("No summaries provided for hierarchical processing, returning default.")
        return "No content to summarize."
    
    if len(summaries) == 1:
        result = summaries[0][:MAX_SUMMARY_LENGTH] if len(summaries[0]) > MAX_SUMMARY_LENGTH else summaries[0]
        logger.debug("Only one summary, returning it directly.")
        return summaries[0][:MAX_SUMMARY_LENGTH]
    
    # If we have many summaries, consolidate them in batches
    if len(summaries) > 20:
        logger.info(f"Hierarchical summary: Processing {len(summaries)} summaries in batches of 10.")
        # First level consolidation: group summaries into batches
        batch_summaries = []
        batch_size = 10
        
        for i in range(0, len(summaries), batch_size):
            batch = summaries[i:i + batch_size]
            combined_batch = "\n\n".join(batch)
            logger.debug(f"Combining batch {i//batch_size + 1} with {len(batch)} summaries.")
            try:
                batch_summary = generate_summary_from_text(combined_batch, user, max_length=500)
                batch_summaries.append(batch_summary)
                logger.debug(f"Generated batch summary {i//batch_size + 1}, length: {len(batch_summary)}")
            except Exception as batch_e:
                logger.warning(f"Failed to generate summary for batch {i//batch_size + 1}: {batch_e}")
                batch_summaries.append(f"[Batch summary failed for batch {i//batch_size + 1}]")
        
        # Second level: consolidate batch summaries
        final_combined = "\n\n".join(batch_summaries)
        logger.info(f"Final consolidation of {len(batch_summaries)} batch summaries.")
        try:
            result = generate_summary_from_text(final_combined, user, max_length=MAX_SUMMARY_LENGTH)
            logger.info(f"Generated final hierarchical summary, length: {len(result)}")
            return result
        except Exception as final_e:
            logger.error(f"Failed to generate final summary from batch summaries: {final_e}")
            # Fallback to truncating the combined batch summaries
            fallback_result = final_combined[:MAX_SUMMARY_LENGTH] + "..." if len(final_combined) > MAX_SUMMARY_LENGTH else final_combined
            logger.info(f"Fallback final summary generated (truncated), length: {len(fallback_result)}")
            return fallback_result

    else:
        logger.info(f"Hierarchical summary: Direct consolidation of {len(summaries)} summaries.")
        # Direct consolidation for smaller number of summaries
        combined = "\n\n".join(summaries)
        try:
            result = generate_summary_from_text(combined, user, max_length=MAX_SUMMARY_LENGTH)
            logger.info(f"Generated direct hierarchical summary, length: {len(result)}")
            return result
        except Exception as direct_e:
            logger.error(f"Failed to generate direct summary: {direct_e}")
            # Fallback to truncation
            fallback_result = combined[:MAX_SUMMARY_LENGTH] + "..." if len(combined) > MAX_SUMMARY_LENGTH else combined
            logger.info(f"Fallback direct summary generated (truncated), length: {len(fallback_result)}")
            return fallback_result



def generate_summary_from_text(text: str, user: Any = None, max_length: int = MAX_SUMMARY_LENGTH) -> str:
    """
    Generate a summary from text using the user's configured LLM provider.
    Falls back to intelligent truncation if LLM summarization fails.
    """
    if not text or not text.strip():
        return "No content to summarize."

    # Limit the input text to prevent exceeding LLM context windows or causing timeouts
    input_text = text[:MAX_TOKEN_THRESHOLD] if len(text) > MAX_TOKEN_THRESHOLD else text
    logger.debug(f"Generating summary for text of length {len(input_text)} (original: {len(text)}) for user {user.id if user else 'N/A'}")

    try:
        if not user:
            logger.warning("No user provided for LLM summary generation, attempting with default settings.")
            # Return truncated version if no user context
            return _generate_fallback_summary(input_text, max_length)

        # Use your existing LLM configuration system
        llm_config = get_llm_config(user.id)
        provider_instance = llm_config.get("provider_instance")
        chat_model = llm_config.get("chat_model")
        api_key = llm_config.get("api_key")
        provider_name = llm_config.get("provider_name")

        if not provider_instance or not chat_model:
            logger.warning(f"LLM provider or model not configured for user {user.id}. Falling back to truncation.")
            return _generate_fallback_summary(input_text, max_length)

        logger.debug(f"Using LLM Provider: {provider_name}, Model: {chat_model} for summary generation.")

        # Construct the prompt for summarization
        requested_words = min(max_length // 4, 100)  # Rough heuristic: 4 chars per word, cap at 100 words
        summary_prompt = (
            f"Provide a concise summary (maximum {requested_words} words) "
            f"of the following document content:\n\n{input_text}"
        )
        logger.debug(f"Summary prompt (first 200 chars): {summary_prompt[:200]}...")

        # Call the LLM using your provider's method
        messages = [{'role': 'user', 'content': summary_prompt}]
        llm_response = provider_instance.get_chat_completion(messages, model=chat_model, api_key=api_key)

        # Extract the summary text from the response
        summary_text = _extract_answer_text(llm_response)

        summary_text = summary_text.strip()

        # Ensure the final summary respects the max_length constraint
        if len(summary_text) > max_length:
            logger.warning(f"LLM summary for user {user.id if user else 'N/A'} exceeded max_length ({len(summary_text)} > {max_length}), truncating.")
            summary_text = summary_text[:max_length].rsplit(' ', 1)[0] + "..."

        logger.info(f"Successfully generated LLM summary of length {len(summary_text)} for user {user.id if user else 'N/A'}")
        return summary_text

    except Exception as e:
        logger.error(f"LLM summary generation failed for user {user.id if user else 'N/A'}: {e}. Input text length: {len(input_text)}. Falling back to truncation.", exc_info=True)
        return _generate_fallback_summary(input_text, max_length)


def _generate_fallback_summary(text: str, max_length: int) -> str:
    """Generate a fallback summary using intelligent truncation"""
    try:
        # Provide a meaningful truncated version if LLM fails
        words = text.split()
        target_word_count = min(100, len(words))  # Aim for up to 100 words or total words if less
        
        if target_word_count > 50:
            # Take first 30% and last 40% of the target word count, with an indicator
            front_count = int(target_word_count * 0.3)
            back_count = int(target_word_count * 0.4)
            middle_indicator = ["[... CONTENT SUMMARIZED ...]"]
            if front_count + back_count < target_word_count:
                # Adjust if rounding caused a small difference
                back_count += target_word_count - (front_count + back_count)
            truncated_words = words[:front_count] + middle_indicator + (words[-back_count:] if back_count > 0 else [])
        else:
            # For very short texts or low target, just take the beginning
            truncated_words = words[:target_word_count]

        truncated_text = " ".join(truncated_words)

        # Ensure final truncated text also respects max_length
        if len(truncated_text) > max_length:
            truncated_text = truncated_text[:max_length].rsplit(' ', 1)[0] + "..."
        elif len(truncated_text) < len(text):
            truncated_text += "..."  # Indicate truncation happened

        logger.info(f"Fallback truncation summary generated (length: {len(truncated_text)})")
        return truncated_text
        
    except Exception as fallback_e:
        logger.critical(f"Fallback truncation also failed: {fallback_e}. Returning error message.")
        return f"(Summary generation failed: {str(fallback_e)})"


def generate_caption_from_text(text: str, user: Any = None, max_length: int = 300) -> str:
    """Generate a concise caption for image-based content."""
    try:
        return generate_summary_from_text(text, user, max_length=max_length)
    except Exception:
        return text[:max_length]

def debug_large_file_processing(file_path: str, document_id: str, collection_name: str):
    """Debug function to track large file processing steps"""
    debug_info = {
        "file_path": file_path,
        "document_id": document_id,
        "collection_name": collection_name,
        "timestamp": timezone.now(),
        "steps": []
    }
    
    try:
        # Check file size and row count
        file_size = os.path.getsize(file_path)
        debug_info["file_size_mb"] = file_size / (1024 * 1024)
        
        if file_path.endswith(('.xlsx', '.xls')):
            xl_file = pd.ExcelFile(file_path)
            total_rows = 0
            for sheet in xl_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet, nrows=1)
                sheet_rows = len(pd.read_excel(file_path, sheet_name=sheet))
                total_rows += sheet_rows
                debug_info["steps"].append(f"Sheet {sheet}: {sheet_rows} rows")
            
            debug_info["total_rows"] = total_rows
            debug_info["estimated_chunks"] = total_rows // getattr(settings, 'ROWS_PER_CHUNK', 100)
            xl_file.close()
        
        elif file_path.endswith('.csv'):
            # Quick row count for CSV
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                total_rows = sum(1 for _ in f)
            debug_info["total_rows"] = total_rows
            debug_info["estimated_chunks"] = total_rows // getattr(settings, 'ROWS_PER_CHUNK', 100)
        
        # Check if large file processing will be triggered
        is_large = total_rows > getattr(settings, 'LARGE_FILE_THRESHOLD', 10000)
        debug_info["is_large_file"] = is_large
        debug_info["processing_method"] = "streaming" if is_large else "standard"
        
        # Log memory usage
        try:
            import psutil
            process = psutil.Process()
            debug_info["memory_usage_mb"] = process.memory_info().rss / (1024 * 1024)
        except ImportError:
            debug_info["memory_usage_mb"] = "psutil not available"
        
        # Add settings information
        debug_info["settings"] = {
            "LARGE_FILE_THRESHOLD": getattr(settings, 'LARGE_FILE_THRESHOLD', 10000),
            "ROWS_PER_CHUNK": getattr(settings, 'ROWS_PER_CHUNK', 100),
            "EXCEL_MAX_CHUNK_SIZE": getattr(settings, 'EXCEL_MAX_CHUNK_SIZE', 3000),
            "QDRANT_BATCH_SIZE": getattr(settings, 'QDRANT_BATCH_SIZE', 50)
        }
        
        logger.info(f"LARGE FILE DEBUG INFO: {json.dumps(debug_info, indent=2, default=str)}")
        
        return debug_info
        
    except Exception as e:
        logger.error(f"Error in debug_large_file_processing: {e}")
        debug_info["error"] = str(e)
        return debug_info
    


def monitor_memory_usage(operation_name: str, document_id: str):
    """Monitor memory usage during large file operations"""
    try:
        import psutil
        import gc
        
        process = psutil.Process()
        memory_info = process.memory_info()
        memory_mb = memory_info.rss / (1024 * 1024)
        
        logger.info(f"MEMORY MONITOR [{operation_name}] Document {document_id}: {memory_mb:.2f}MB used")
        
        # Force garbage collection if memory usage is high
        if memory_mb > 1024:  # 1GB threshold
            logger.warning(f"High memory usage detected ({memory_mb:.2f}MB), forcing garbage collection")
            gc.collect()
            
            # Log memory after cleanup
            post_gc_memory = psutil.Process().memory_info().rss / (1024 * 1024)
            logger.info(f"Memory after GC: {post_gc_memory:.2f}MB (freed: {memory_mb - post_gc_memory:.2f}MB)")
            
        return memory_mb
        
    except ImportError:
        logger.warning("psutil not available for memory monitoring")
        return 0
    except Exception as e:
        logger.warning(f"Could not monitor memory usage: {e}")
        return 0



# --- File Processing & Text Extraction ---

def process_file(file: Any = None, s3_file_url: Optional[str] = None, file_name: Optional[str] = None, checksum: Optional[str] = None) -> Tuple[str, str]:
    """
    Processes an uploaded file or downloads it from an S3 URL, saves it temporarily,
    and optionally validates its checksum.
    Returns the temporary file path and the actual file name.
    """
    logger.debug(f"Processing file. Provided file object: {bool(file)}, S3 URL: {s3_file_url}, Original filename: {file_name}")
    tmp_path: Optional[str] = None
    actual_file_name: str = ""

    try:
        if not file and not s3_file_url:
            raise ValidationError("Either a file object or an S3 file URL must be provided.")

        # Determine suffix for temp file
        suffix = ".tmp"
        if file_name:
            suffix = Path(file_name).suffix
        elif s3_file_url:
            suffix = Path(urlparse(s3_file_url).path).suffix
        elif file and getattr(file, 'name', None):
            suffix = Path(getattr(file, 'name')).suffix or ".tmp"

        # Create a temporary file to store the content
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp_path = tmp.name  # Store the path immediately for cleanup in finally

            if file: # Handle uploaded file
                actual_file_name = file_name or getattr(file, 'name', 'uploaded_file')
                logger.info(f"Processing uploaded file: '{actual_file_name}' (Size: {getattr(file, 'size', 'unknown')} bytes)")
                if getattr(file, 'size', 0) > MAX_FILE_SIZE:
                    raise ValidationError(f"File size ({file.size / 1024 / 1024:.2f}MB) exceeds limit of {MAX_FILE_SIZE / 1024 / 1024}MB.")

                buffer_size = 65536  # 64KB buffer
                for chunk in file.chunks(chunk_size=buffer_size):
                    tmp.write(chunk)
                logger.debug(f"Uploaded file '{actual_file_name}' saved to temporary path: {tmp_path}")

            elif s3_file_url: # Handle S3 file URL
                actual_file_name = file_name or Path(urlparse(s3_file_url).path).name
                logger.info(f"Downloading S3 file: '{actual_file_name}' from URL: {s3_file_url}")
                try:
                    close_connections_before_io("S3 file download")
                    response = call_with_resilience(
                        lambda: requests.get(s3_file_url, stream=True, timeout=60),
                        service="s3_file_download",
                        exceptions=(requests.exceptions.RequestException,),
                    )
                    response.raise_for_status()

                    content_length_header = response.headers.get('content-length')
                    if content_length_header and int(content_length_header) > MAX_FILE_SIZE:
                        raise ValidationError(f"S3 file size ({int(content_length_header) / 1024 / 1024:.2f}MB based on 'content-length' header) exceeds limit of {MAX_FILE_SIZE / 1024 / 1024}MB.")

                    downloaded_size = 0
                    buffer_size = 65536  # 64KB buffer
                    for chunk in response.iter_content(chunk_size=buffer_size):
                        tmp.write(chunk)
                        downloaded_size += len(chunk)
                        if downloaded_size > MAX_FILE_SIZE:
                            raise ValidationError(f"S3 file download exceeded size limit of {MAX_FILE_SIZE / 1024 / 1024}MB during streaming.")
                    logger.debug(f"S3 file '{actual_file_name}' downloaded and saved to temporary path: {tmp_path}")
                except requests.exceptions.Timeout:
                    logger.error(f"Timeout while downloading S3 file {s3_file_url}", exc_info=True)
                    raise RuntimeError(f"Failed to download file from S3 due to timeout: {s3_file_url}")
                except requests.exceptions.HTTPError as http_err:
                    logger.error(f"HTTP error {http_err.response.status_code} while downloading S3 file {s3_file_url}", exc_info=True)
                    raise RuntimeError(f"Failed to download file from S3: HTTP {http_err.response.status_code} for {s3_file_url}")
                except requests.exceptions.RequestException as req_e:
                    logger.error(f"Failed to download S3 file {s3_file_url}: {req_e}", exc_info=True)
                    raise RuntimeError(f"Failed to download file from S3: {str(req_e)}")

        # tmp file is now closed, proceed with checksum if needed
        if checksum and tmp_path: # tmp_path must exist
            logger.info(f"Validating checksum for '{actual_file_name}'...")
            sha = hashlib.sha256()
            buffer_size = 65536  # 64KB buffer
            try:
                with open(tmp_path, 'rb') as f_checksum:
                    while True:
                        chunk = f_checksum.read(buffer_size)
                        if not chunk:
                            break
                        sha.update(chunk)
            except IOError as ioe:
                logger.error(f"IOError during checksum calculation for {tmp_path}: {ioe}", exc_info=True)
                raise RuntimeError(f"Could not read temporary file for checksum: {str(ioe)}")

            calculated_checksum = sha.hexdigest()
            if calculated_checksum != checksum:
                logger.error(f"Checksum mismatch for '{actual_file_name}'. Expected: {checksum}, Calculated: {calculated_checksum}")
                raise ValidationError("Checksum validation failed. File may be corrupted or altered.")
            logger.info(f"Checksum validated successfully for '{actual_file_name}'.")

        if not tmp_path: # Should not happen if logic is correct, but as a safeguard
            raise RuntimeError("Temporary file path was not set during file processing.")

        logger.info(f"Successfully processed file: '{actual_file_name}', temporary path: {tmp_path}")
        return tmp_path, actual_file_name

    except ValidationError: # Re-raise ValidationError directly
        raise
    except Exception as e: # Catch other unexpected errors
        logger.error(f"Unexpected error processing file '{file_name or s3_file_url}': {e}", exc_info=True)
        # Wrap non-ValidationError exceptions in a generic one for consistent error handling upstream
        raise RuntimeError(f"File processing failed due to an unexpected error: {str(e)}")
    finally:
        # Ensure cleanup happens if tmp_path was created, even if an error occurred after file creation
        # but before the function could return normally (e.g. during checksum).
        # If an error occurs *during* tmp.write, NamedTemporaryFile handles its own cleanup if delete=True.
        # Since we use delete=False, we must manually clean up.
        # Note: If an error occurs *before* tmp_path is assigned, this finally block won't do anything, which is fine.
        # If an error occurred *after* returning tmp_path, the caller is responsible for cleanup.
        # This 'finally' is for errors within this function *after* tmp_path is set.
        # However, the current structure returns tmp_path for the caller to manage,
        # so this 'finally' might only be relevant if we change that pattern.
        # For now, if an error occurs and tmp_path is set, it implies the caller won't get it, so we should try to clean.
        # Let's consider if an error like Checksum validation failure should lead to cleanup here or if caller still needs the file.
        # If process_file raises an exception, the caller (IngestAPIView) *does* try to clean up tmp_path.
        # So, this specific finally block might be redundant for cleanup if errors are always propagated.
        # Keeping it minimal for now.
        pass

# --- Text Extraction Methods (Individual file types) ---
# (Keep existing extraction methods like extract_text_from_pdf, _docx, etc.)
# ... (Assume existing methods for pdf, docx, pptx, txt, json, xml, html, md, yaml, ini are here)

def extract_text_from_pdf(file_path: str, user: Any = None) -> str:
    logger.debug(f"Attempting to extract text from PDF: '{file_path}'")
    basic_text = ""
    try:
        # Try PyPDF2 first
        logger.debug(f"Using PyPDF2 for PDF: '{file_path}'")
        text_parts = []
        with open(file_path, 'rb') as f:
            reader = PdfReader(f)
            if reader.is_encrypted:
                logger.warning(f"PDF file '{file_path}' is encrypted. Cannot extract text without decryption key.")
                raise ValidationError(f"Cannot extract text from encrypted PDF: {Path(file_path).name}")
            for i, page in enumerate(reader.pages):
                try:
                    page_text = page.extract_text()
                    text_parts.append(page_text or "")
                except Exception as page_e: # Catch error per page
                    logger.warning(f"Error extracting text from page {i+1} in PDF '{file_path}' using PyPDF2: {page_e}", exc_info=True)
                    text_parts.append("") # Append empty string on page error to not lose other pages
        basic_text = "\n".join(filter(None, text_parts)).strip()
        logger.info(f"PyPDF2 extracted {len(basic_text)} characters from '{file_path}'.")

        # If PyPDF2 yields little text (e.g., scanned PDF), try pdfminer.high_level
        if len(basic_text) < 100: # Arbitrary threshold, adjust as needed
            logger.info(f"PyPDF2 extracted minimal text from '{file_path}'. Attempting fallback with pdfminer.")
            try:
                miner_text = pdfminer_extract_text(file_path).strip()
                if len(miner_text) > len(basic_text):
                    logger.info(f"pdfminer extracted more text ({len(miner_text)} chars) from '{file_path}' than PyPDF2. Using pdfminer result.")
                    return miner_text
                else:
                    logger.info(f"pdfminer did not extract significantly more text from '{file_path}'. Sticking with PyPDF2 result.")
            except Exception as miner_e:
                logger.warning(f"pdfminer extraction failed for '{file_path}': {miner_e}. Using PyPDF2's result if available.", exc_info=True)

        if not basic_text:
            logger.warning(f"No text extracted from PDF '{file_path}' using PyPDF2 or pdfminer. File might be image-based or empty.")
            # OCR fallback is handled by extract_text_from_file if this returns empty
            return ""

        return basic_text

    except FileNotFoundError:
        logger.error(f"PDF file not found: '{file_path}'", exc_info=True)
        raise RuntimeError(f"PDF file not found: {Path(file_path).name}")
    except (PdfReader.PdfReadError, PyPDF2.errors.PdfReadError) as pdf_err: # PyPDF2.errors for older versions
        logger.error(f"Invalid or corrupted PDF file '{file_path}': {pdf_err}", exc_info=True)
        raise ValidationError(f"Invalid or corrupted PDF file: {Path(file_path).name}. Error: {pdf_err}")
    except Exception as e:
        logger.error(f"Unexpected error extracting text from PDF '{file_path}': {e}", exc_info=True)
        if isinstance(e, ValidationError): # Re-raise if it's already a ValidationError
             raise e
        raise RuntimeError(f"PDF extraction failed for {Path(file_path).name}: {str(e)}")

def extract_text_from_docx(file_path: str, user: Any = None) -> str:
    logger.debug(f"Attempting to extract text from DOCX: '{file_path}'")
    text = ""
    try:
        # Attempt with docx2txt first
        logger.debug(f"Using docx2txt for DOCX: '{file_path}'")
        text = docx2txt.process(file_path).strip()
        if text:
            logger.info(f"docx2txt successfully extracted {len(text)} characters from '{file_path}'.")
            return text

        # Fallback to python-docx if docx2txt yields no text
        logger.warning(f"docx2txt extracted no text from DOCX '{file_path}'. Attempting fallback with python-docx.")
        doc = DOCXDocument(file_path)
        text_parts = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
        # Optionally extract from tables, headers, footers if needed (more complex)
        # for table in doc.tables:
        #     for row in table.rows:
        #         for cell in row.cells:
        #             text_parts.append(cell.text.strip())
        text = '\n'.join(text_parts).strip()
        if text:
            logger.info(f"python-docx successfully extracted {len(text)} characters from '{file_path}' as fallback.")
        else:
            logger.warning(f"Both docx2txt and python-docx extracted no text from '{file_path}'.")
        return text

    except FileNotFoundError:
        logger.error(f"DOCX file not found: '{file_path}'", exc_info=True)
        raise RuntimeError(f"DOCX file not found: {Path(file_path).name}")
    except Exception as e: # Catches errors from both libraries if they occur e.g. docx.opc.exceptions.PackageNotFoundError
        logger.error(f"Error extracting text from DOCX '{file_path}': {e}", exc_info=True)
        raise RuntimeError(f"DOCX extraction failed for {Path(file_path).name}: {str(e)}")


def _preprocess_for_ocr(img: Image.Image) -> Image.Image:
    """
    Generic preprocessing that works reasonably well for most images:
    - ensure RGB
    - upscale small images
    - grayscale + autocontrast
    - mild sharpening
    (no aggressive thresholding so natural photos still work)
    """
    if img.mode != "RGB":
        img = img.convert("RGB")

    # Upscale small images so OCR has more pixels to work with
    min_side = min(img.size)
    if min_side < 1000:
        scale = 2
        img = img.resize((img.width * scale, img.height * scale), Image.LANCZOS)

    gray = img.convert("L")
    gray = ImageOps.autocontrast(gray)
    # light sharpen to make text edges crisper without destroying photos
    gray = gray.filter(ImageFilter.UnsharpMask(radius=1, percent=150, threshold=3))

    return gray


def _run_ocr_passes(img: Image.Image) -> str:
    """
    Run OCR with a couple of sensible configs and merge results.
    This helps across different types of layouts (dense text, sparse text, etc.).
    """
    results = []

    # Pass 1: regular block of text (good default)
    try:
        txt1 = pytesseract.image_to_string(
            img,
            lang="eng",
            config=r"--oem 3 --psm 6",
        )
        if txt1:
            results.append(txt1)
    except Exception as e:
        logger.warning(f"OCR pass1 failed: {e}", exc_info=True)

    # Pass 2: sparse text / irregular layout (e.g., diagrams, posters)
    try:
        txt2 = pytesseract.image_to_string(
            img,
            lang="eng",
            config=r"--oem 3 --psm 11",
        )
        if txt2:
            results.append(txt2)
    except Exception as e:
        logger.warning(f"OCR pass2 failed: {e}", exc_info=True)

    # Merge, dedupe lines
    merged = "\n".join(results)
    lines = [l.strip() for l in merged.splitlines()]
    # drop empty & near-duplicate lines
    seen = set()
    cleaned_lines = []
    for l in lines:
        if not l:
            continue
        key = l.lower()
        if key not in seen:
            seen.add(key)
            cleaned_lines.append(l)

    return "\n".join(cleaned_lines).strip()


def _ocr_text_quality_score(text: str) -> float:
    """
    Heuristic score (0–1) for 'how text-like' the OCR output is.
    Used to decide if we should fall back to LLM description.
    """
    if not text:
        return 0.0

    length = len(text)
    if length < 10:
        return 0.1

    # ratio of alphanumeric characters
    alpha_num = sum(c.isalnum() for c in text)
    ratio_alpha = alpha_num / max(length, 1)

    # non-texty symbols that often show up when OCR is garbage
    noisy_chars = "|=/\\[]{}@#%^*~_"
    noisy = sum(1 for c in text if c in noisy_chars)
    ratio_noisy = noisy / max(length, 1)

    score = ratio_alpha - ratio_noisy
    # clamp to [0,1]
    return max(0.0, min(1.0, score))


def extract_text_from_image(file_path: str, user: Any = None) -> str:
    """
    Generic image → text/description extractor.

    Behavior:
    1. Try OCR with preprocessing and multiple passes.
       - If OCR looks good: return cleaned text.
    2. If OCR fails or looks like noise:
       - If user & multimodal model configured:
           * Ask LLM to:
             - Extract any visible text (if present)
             - AND give a short description when there is no/very little text
       - Else: return a best-effort OCR result or a fallback string.
    """
    logger.debug(f"Attempting to extract text/description from Image: '{file_path}'")
    ocr_text = ""
    try:
        logger.info(f"Attempting OCR on image: '{file_path}'")
        img = Image.open(file_path)
        preprocessed = _preprocess_for_ocr(img)
        ocr_text = _run_ocr_passes(preprocessed)

        quality = _ocr_text_quality_score(ocr_text)
        logger.info(
            f"OCR finished for '{file_path}'. Length={len(ocr_text)}, quality_score={quality:.3f}"
        )

        # If OCR result looks reasonably text-like, use it directly
        if quality >= 0.4:
            logger.info(
                f"OCR result accepted for '{file_path}' with quality {quality:.3f}."
            )
            return ocr_text

        logger.info(
            f"OCR result for '{file_path}' seems poor (quality {quality:.3f}). "
            f"Will consider LLM fallback if available."
        )

    except pytesseract.TesseractNotFoundError:
        logger.error(
            "Tesseract is not installed or not found in PATH. OCR is unavailable for image processing."
        )
        # allow fallback to LLM if user is provided
        ocr_text = ""
    except FileNotFoundError:
        logger.error(f"Image file not found: '{file_path}'", exc_info=True)
        raise RuntimeError(f"Image file not found: {Path(file_path).name}")
    except Exception as ocr_e:
        logger.warning(
            f"OCR failed for image '{file_path}': {ocr_e}. Proceeding to LLM description if applicable.",
            exc_info=True,
        )
        ocr_text = ""

    # 2. Fallback to LLM (if user context is available)
    if not user:
        logger.warning(
            f"No user context provided for image '{file_path}', and OCR was empty/low-quality. "
            "Returning OCR output (which may be empty or noisy)."
        )
        return ocr_text  # may be empty; best we can do without LLM

    logger.info(f"Using LLM to enhance/describe image: '{file_path}' (User: {user.id})")
    try:
        llm_config = get_llm_config(user.id)
        provider_instance = llm_config.get("provider_instance")
        chat_model = llm_config.get("chat_model")
        api_key = llm_config.get("api_key")

        if not provider_instance or not chat_model:
            logger.warning(
                f"LLM provider or chat model not configured for user {user.id}. "
                f"Cannot describe image '{file_path}'."
            )
            # fall back to whatever OCR we have
            return ocr_text or "Image content (OCR failed, LLM description unavailable due to configuration)"

        if not provider_instance.is_multimodal(model=chat_model, api_key=api_key):
            logger.warning(
                f"LLM provider '{llm_config.get('provider_name')}' model '{chat_model}' "
                f"is not available or not multimodal. Cannot describe image '{file_path}'."
            )
            return ocr_text or "Image content (OCR failed, LLM description unavailable as model is not multimodal)"

        with open(file_path, "rb") as image_file:
            image_data = base64.b64encode(image_file.read()).decode("utf-8")

        mime_type = detect_mime_type(file_path, default="image/jpeg")
        logger.debug(
            f"Image '{file_path}' has MIME type: {mime_type} for LLM description."
        )

        # Prompt is generic: works for charts, nature, humans, quotes, etc.
        #  - If there is text, we want it captured cleanly.
        #  - If there is no text, we still want a useful description.
        base_prompt = (
            "You are processing an arbitrary image uploaded by a user.\n"
            "1) If there is any legible text in the image (signs, labels, captions, quotes, UI text, etc.), "
            "extract it as accurately as possible.\n"
            "2) If there is little or no text, say 'NO_CLEAR_TEXT_FOUND' on the first line.\n"
            "3) On the next line, starting with 'DESCRIPTION:', provide a concise but informative description "
            "of what is shown in the image (objects, layout, relationships, etc.).\n\n"
            "Examples:\n"
            "Text-heavy image:\n"
            "HELLO WORLD\n"
            "DESCRIPTION: A white poster with the words 'HELLO WORLD' in large black letters.\n\n"
            "No-text image:\n"
            "NO_CLEAR_TEXT_FOUND\n"
            "DESCRIPTION: A landscape photo of mountains and a lake at sunset.\n\n"
            "Now process this image."
        )

        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": base_prompt},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:{mime_type};base64,{image_data}",
                        },
                    },
                ],
            }
        ]

        response = provider_instance.get_chat_completion(
            messages, model=chat_model, api_key=api_key
        )

        description = _extract_answer_text(response)

        if not description or not description.strip():
            logger.warning(f"LLM description for image '{file_path}' was empty.")
            # If we still have some OCR text, at least return that
            return ocr_text or "Image content (OCR failed, LLM description was empty)"

        final_text = description.strip()

        # If OCR was somewhat OK, you may optionally prepend it for traceability:
        # (Uncomment if you want BOTH)
        #
        # if ocr_text and _ocr_text_quality_score(ocr_text) >= 0.2:
        #     final_text = f"OCR_TEXT:\n{ocr_text.strip()}\n\nLLM_ENHANCED:\n{final_text}"

        logger.info(
            f"LLM description successful for image '{file_path}'. Length: {len(final_text)}"
        )
        return final_text

    except FileNotFoundError:
        logger.error(
            f"Image file not found during LLM description phase: '{file_path}'",
            exc_info=True,
        )
        raise RuntimeError(f"Image file not found: {Path(file_path).name}")
    except Exception as llm_e:
        logger.error(
            f"Error getting LLM description for image '{file_path}': {llm_e}",
            exc_info=True,
        )
        # Return OCR text if it was available, otherwise indicate combined failure
        return ocr_text or f"Image content (OCR/LLM description failed: {str(llm_e)})"

def extract_text_from_pptx(file_path: str, user: Any = None) -> str:
    logger.debug(f"Attempting to extract text from PPTX: '{file_path}'")
    try:
        prs = Presentation(file_path)
        text_parts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    logger.debug(f"Extracting notes from slide {i+1} in '{file_path}'")
                    slide_texts.append(f"--- Notes Slide {i+1} ---\n{notes_text}")

            if slide_texts:
                 text_parts.append(f"--- Slide {i+1} ---\n" + '\n'.join(slide_texts))

        full_text = '\n\n'.join(text_parts).strip()
        logger.info(f"Successfully extracted {len(full_text)} characters from PPTX '{file_path}'.")
        return full_text
    except FileNotFoundError:
        logger.error(f"PPTX file not found: '{file_path}'", exc_info=True)
        raise RuntimeError(f"PPTX file not found: {Path(file_path).name}")
    except Exception as e: # Catch errors like pptx.exc.PackageNotFoundError
        logger.error(f"Error extracting text from PPTX '{file_path}': {e}", exc_info=True)
        raise RuntimeError(f"PPTX extraction failed for {Path(file_path).name}: {str(e)}")


def extract_structured_from_excel_or_csv(file_path: str, original_file_name: str | None = None) -> str:
    logger.debug("Attempting to extract structured data from Excel/CSV: '%s'", file_path)
    file_path_obj = Path(file_path)
    ext = file_path_obj.suffix.lower()
    if ext not in [".xls", ".xlsx", ".csv"] and original_file_name:
        ext = Path(original_file_name).suffix.lower()
    output = ""

    try:
        if ext in ['.xls', '.xlsx']:
            logger.debug(f"Processing Excel file: '{file_path}'")
            # Check if it's a large file first
            if is_large_file_by_row_count(file_path, ext):
                logger.info(f"Large Excel file detected: '{file_path}'. Using streaming approach.")
                return process_large_excel_streaming(file_path, original_file_name or file_path_obj.name)[0]

            # For regular-sized files, use enhanced processing
            excel_data = pd.read_excel(file_path, sheet_name=None) # Reads all sheets
            if not excel_data:
                logger.warning(f"No data found in Excel file: '{file_path}' (no sheets or all empty).")
                return "(Empty Excel File)"
            
            total_rows_processed = 0

            for sheet_name, df in excel_data.items():
                logger.info(f"Processing sheet '{sheet_name}' with {len(df)} rows and {len(df.columns)} columns")
                output += f"## Sheet: {sheet_name}\n\n"

                if df.empty:
                    output += "(Empty Sheet)\n\n"
                    continue

                # Enhanced processing to prevent row loss
                try:
                    # Method 1: Use to_string for better reliability
                    df_text = df.to_string(index=False, max_colwidth=100, max_rows=None)
                    output += df_text + "\n\n"
                    total_rows_processed += len(df)
                    logger.debug(f"Sheet '{sheet_name}': Processed {len(df)} rows using to_string method")
                    
                except Exception as string_e:
                    logger.warning(f"to_string failed for sheet '{sheet_name}': {string_e}. Trying row-by-row processing.")
                    
                    # Method 2: Row-by-row processing as fallback
                    try:
                        # Add column headers
                        headers = " | ".join(str(col) for col in df.columns)
                        output += headers + "\n"
                        output += " | ".join(["---"] * len(df.columns)) + "\n"
                        
                        # Process each row individually
                        rows_added = 0
                        for idx, row in df.iterrows():
                            try:
                                row_text = " | ".join(str(val) if pd.notna(val) else "" for val in row.values)
                                output += row_text + "\n"
                                rows_added += 1
                            except Exception as row_e:
                                logger.warning(f"Failed to process row {idx} in sheet '{sheet_name}': {row_e}")
                                # Add placeholder for failed row
                                output += f"[Row {idx}: Processing failed]\n"
                                rows_added += 1
                        
                        output += "\n"
                        total_rows_processed += rows_added
                        logger.info(f"Sheet '{sheet_name}': Processed {rows_added}/{len(df)} rows using row-by-row method")
                        
                        if rows_added < len(df):
                            logger.warning(f"Sheet '{sheet_name}': Lost {len(df) - rows_added} rows during processing")
                    
                    except Exception as row_e:
                        logger.error(f"Row-by-row processing failed for sheet '{sheet_name}': {row_e}")
                        # Last resort: basic info
                        output += f"[Sheet processing failed - {len(df)} rows, {len(df.columns)} columns]\n\n"

            logger.info(f"Excel processing complete: {total_rows_processed} total rows processed from {len(excel_data)} sheets")

        elif ext == '.csv':
            logger.debug(f"Processing CSV file: '{file_path}'")
            
            # Check if it's a large file first
            if is_large_file_by_row_count(file_path, ext):
                logger.info(f"Large CSV file detected: '{file_path}'. Using streaming approach.")
                return process_large_csv_streaming(file_path, original_file_name or file_path_obj.name)[0]
            
            # Detect encoding for CSV
            with open(file_path, 'rb') as f_enc:
                raw_data = f_enc.read(50000)
                encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
            logger.debug(f"Detected CSV encoding for '{file_path}': {encoding}")

            df = pd.read_csv(file_path, encoding=encoding)
            logger.info(f"CSV loaded: {len(df)} rows, {len(df.columns)} columns")
            
            output = "## CSV Data\n\n"
            if df.empty:
                output += "(Empty CSV)"
            else:
                try:
                    # Try to_string first (more reliable than to_markdown)
                    csv_text = df.to_string(index=False, max_colwidth=100, max_rows=None)   
                    output += csv_text
                    logger.info(f"CSV processed: {len(df)} rows using to_string method")
                    
                except Exception as string_e:
                    logger.warning(f"to_string failed for CSV: {string_e}. Trying row-by-row processing.")
                    
                    # Fallback to row-by-row
                    try:
                        headers = " | ".join(str(col) for col in df.columns)
                        output += headers + "\n"
                        output += " | ".join(["---"] * len(df.columns)) + "\n"
                        
                        rows_added = 0
                        for idx, row in df.iterrows():
                            try:
                                row_text = " | ".join(str(val) if pd.notna(val) else "" for val in row.values)
                                output += row_text + "\n"
                                rows_added += 1
                            except Exception as row_e:
                                logger.warning(f"Failed to process CSV row {idx}: {row_e}")
                                output += f"[Row {idx}: Processing failed]\n"
                                rows_added += 1
                        
                        logger.info(f"CSV processed: {rows_added}/{len(df)} rows using row-by-row method")
                        
                        if rows_added < len(df):
                            logger.warning(f"CSV: Lost {len(df) - rows_added} rows during processing")
                    
                    except Exception as row_e:
                        logger.error(f"Row-by-row CSV processing failed: {row_e}")
                        output += f"[CSV processing failed - {len(df)} rows, {len(df.columns)} columns]"

            logger.info(f"Successfully extracted data from CSV file '{file_path_obj.name}'.")
        else:
            logger.warning(f"Unsupported file type for structured extraction: {ext} for file '{file_path}'")
            raise ValueError(f"Unsupported file type for structured extraction: {ext}")
        
        # Add row verification before returning
        verification_result = verify_row_extraction(file_path, output, ext)
        if "error" not in verification_result:
            logger.info(f"Row extraction verification for '{file_path_obj.name}': "
                       f"Actual rows: {verification_result['actual_rows']}, "
                       f"Extracted lines: {verification_result['extracted_lines']}, "
                       f"Extraction rate: {verification_result['extraction_rate']:.1f}%, "
                       f"Potential loss: {verification_result['potential_loss']} rows")
            
            if verification_result['potential_loss'] > 0:
                logger.warning(f"Potential row loss detected in '{file_path_obj.name}': "
                             f"{verification_result['potential_loss']} rows may be missing")
        else:
            logger.warning(f"Row verification failed for '{file_path_obj.name}': {verification_result['error']}")

        return output.strip()
        
    except FileNotFoundError:
        logger.error(f"Excel/CSV file not found: '{file_path}'", exc_info=True)
        raise RuntimeError(f"Excel/CSV file not found: {file_path_obj.name}")
    except pd.errors.EmptyDataError:
        logger.warning(f"Excel/CSV file '{file_path}' is empty or contains no data.", exc_info=True)
        return f"(Empty {ext.upper()} File)"
    except Exception as e:
        logger.error(f"Error extracting structured data from '{file_path_obj.name}': {e}", exc_info=True)
        raise RuntimeError(f"Failed to extract data from {ext.upper()} file '{file_path_obj.name}': {str(e)}")

# Add this function after the extract_structured_from_excel_or_csv function
def verify_row_extraction(file_path: str, extracted_text: str, ext: str) -> Dict[str, Any]:
    """Verify that all rows were extracted from Excel/CSV files"""
    try:
        if ext == '.csv':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                actual_rows = sum(1 for _ in f) - 1  # Subtract header
        elif ext in ['.xlsx', '.xls']:
            xl_file = pd.ExcelFile(file_path)
            actual_rows = 0
            try:
                for sheet_name in xl_file.sheet_names:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    actual_rows += len(df)
            finally:
                # Safely close the ExcelFile object
                try:
                    xl_file.close()
                except (AttributeError, Exception):
                    pass  # Some ExcelFile objects don't have close() method
        else:
            return {"error": "Unsupported file type for row verification"}
        
        # Count rows in extracted text (rough estimate)
        text_lines = len([line for line in extracted_text.split('\n') if line.strip() and not line.startswith('#')])
        
        return {
            "actual_rows": actual_rows,
            "extracted_lines": text_lines,
            "extraction_rate": (text_lines / actual_rows * 100) if actual_rows > 0 else 0,
            "potential_loss": max(0, actual_rows - text_lines)
        }
    except Exception as e:
        return {"error": str(e)}


def extract_text_from_txt(file_path: str) -> str:
    logger.debug(f"Attempting to extract text from TXT file: '{file_path}'")
    try:
        with open(file_path, 'rb') as f_enc: # Read as bytes first for encoding detection
            raw_data = f_enc.read()
            encoding = chardet.detect(raw_data)['encoding'] or 'utf-8' # Default to UTF-8
        logger.debug(f"Detected encoding for TXT file '{file_path}': {encoding}")

        with open(file_path, 'r', encoding=encoding, errors='replace') as f_text: # Open with detected encoding
            text = f_text.read()

        stripped_text = text.strip()
        logger.info(f"Successfully extracted {len(stripped_text)} characters from TXT file '{file_path}'.")
        return stripped_text
    except FileNotFoundError:
        logger.error(f"TXT file not found: '{file_path}'", exc_info=True)
        raise RuntimeError(f"TXT file not found: {Path(file_path).name}")
    except IOError as ioe:
        logger.error(f"IOError reading TXT file '{file_path}': {ioe}", exc_info=True)
        raise RuntimeError(f"TXT file reading failed for {Path(file_path).name}: {str(ioe)}")
    except Exception as e:
        logger.error(f"Unexpected error extracting text from TXT file '{file_path}': {e}", exc_info=True)
        raise RuntimeError(f"TXT extraction failed for {Path(file_path).name}: {str(e)}")

def extract_structured_from_json(file_path: str) -> str:
    logger.debug(f"Attempting to extract structured data from JSON file: '{file_path}'")
    try:
        with open(file_path, 'r', encoding='utf-8') as f: # JSON is typically UTF-8
            data = json.load(f)
        output = json.dumps(data, indent=2) # Pretty print for better readability if stored/logged
        logger.info(f"Successfully extracted and formatted JSON data from '{file_path}'. Length: {len(output)}")
        return output
    except FileNotFoundError:
        logger.error(f"JSON file not found: '{file_path}'", exc_info=True)
        raise RuntimeError(f"JSON file not found: {Path(file_path).name}")
    except json.JSONDecodeError as json_e:
        logger.error(f"Invalid JSON format in file '{file_path}': {json_e}", exc_info=True)
        raise ValidationError(f"Invalid JSON format in {Path(file_path).name}: {str(json_e)}")
    except IOError as ioe:
        logger.error(f"IOError reading JSON file '{file_path}': {ioe}", exc_info=True)
        raise RuntimeError(f"JSON file reading failed for {Path(file_path).name}: {str(ioe)}")
    except Exception as e:
        logger.error(f"Unexpected error extracting data from JSON file '{file_path}': {e}", exc_info=True)
        raise RuntimeError(f"JSON extraction failed for {Path(file_path).name}: {str(e)}")

def extract_text_from_xml(file_path: str) -> str:
    logger.debug(f"Attempting to extract text from XML file: '{file_path}'")
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        # Efficiently extract text content from all elements
        text_parts = [elem.text.strip() for elem in root.iter() if elem.text and elem.text.strip()]
        text_content = ' '.join(text_parts)
        # Alternative: return the formatted XML string itself if that's more useful:
        # output = ET.tostring(root, encoding='unicode', method='xml')
        logger.info(f"Successfully extracted {len(text_content)} characters of text from XML file '{file_path}'.")
        return text_content
    except FileNotFoundError:
        logger.error(f"XML file not found: '{file_path}'", exc_info=True)
        raise RuntimeError(f"XML file not found: {Path(file_path).name}")
    except ET.ParseError as xml_e:
        logger.error(f"Invalid XML format in file '{file_path}': {xml_e}", exc_info=True)
        raise ValidationError(f"Invalid XML format in {Path(file_path).name}: {str(xml_e)}")
    except IOError as ioe:
        logger.error(f"IOError reading XML file '{file_path}': {ioe}", exc_info=True)
        raise RuntimeError(f"XML file reading failed for {Path(file_path).name}: {str(ioe)}")
    except Exception as e:
        logger.error(f"Unexpected error extracting text from XML file '{file_path}': {e}", exc_info=True)
        raise RuntimeError(f"XML extraction failed for {Path(file_path).name}: {str(e)}")

def extract_text_from_html(file_path: str) -> str:
    logger.debug(f"Attempting to extract text from HTML file: '{file_path}'")
    try:
        # Detect encoding first
        with open(file_path, 'rb') as f_enc:
            raw_bytes = f_enc.read(20000) # Read more bytes for better encoding detection
            encoding = chardet.detect(raw_bytes)['encoding'] or 'utf-8'
        logger.debug(f"Detected encoding for HTML file '{file_path}': {encoding}")

        with open(file_path, 'r', encoding=encoding, errors='replace') as f_html:
            soup = BeautifulSoup(f_html, 'html.parser')

        # Remove script, style, and other non-visible elements
        for element_type in ['script', 'style', 'noscript', 'head', 'meta', 'link', 'comment']:
            for element in soup.find_all(element_type):
                element.decompose()

        # Get text, trying to preserve some structure with newlines
        text = soup.get_text(separator='\n', strip=True)
        logger.info(f"Successfully extracted {len(text)} characters from HTML file '{file_path}'.")
        return text
    except FileNotFoundError:
        logger.error(f"HTML file not found: '{file_path}'", exc_info=True)
        raise RuntimeError(f"HTML file not found: {Path(file_path).name}")
    except IOError as ioe:
        logger.error(f"IOError reading HTML file '{file_path}': {ioe}", exc_info=True)
        raise RuntimeError(f"HTML file reading failed for {Path(file_path).name}: {str(ioe)}")
    except Exception as e: # Catch potential BeautifulSoup parsing errors
        logger.error(f"Error extracting text from HTML file '{file_path}': {e}", exc_info=True)
        raise RuntimeError(f"HTML extraction failed for {Path(file_path).name}: {str(e)}")

def extract_text_from_md(file_path: str) -> str:
    logger.debug(f"Extracting text from Markdown file: '{file_path}' (delegating to TXT extractor)")
    # For Markdown, reading as plain text is usually sufficient.
    # If specific Markdown parsing to HTML then text is needed, a library like 'markdown' or 'mistune' could be used here.
    return extract_text_from_txt(file_path)

def extract_text_from_yaml(file_path: str) -> str:
    logger.debug(f"Attempting to extract structured data from YAML file: '{file_path}'")
    try:
        with open(file_path, 'r', encoding='utf-8') as f: # YAML often UTF-8
            # Use safe_load to prevent arbitrary code execution from malicious YAML
            data = yaml.safe_load(f)

        # Convert to formatted JSON string for consistent output, similar to JSON/INI
        output = json.dumps(data, indent=2)
        logger.info(f"Successfully extracted and formatted YAML data from '{file_path}'. Length: {len(output)}")
        return output
    except FileNotFoundError:
        logger.error(f"YAML file not found: '{file_path}'", exc_info=True)
        raise RuntimeError(f"YAML file not found: {Path(file_path).name}")
    except yaml.YAMLError as yaml_e:
        logger.error(f"Invalid YAML format in file '{file_path}': {yaml_e}", exc_info=True)
        raise ValidationError(f"Invalid YAML format in {Path(file_path).name}: {str(yaml_e)}")
    except IOError as ioe:
        logger.error(f"IOError reading YAML file '{file_path}': {ioe}", exc_info=True)
        raise RuntimeError(f"YAML file reading failed for {Path(file_path).name}: {str(ioe)}")
    except Exception as e:
        logger.error(f"Unexpected error extracting data from YAML file '{file_path}': {e}", exc_info=True)
        raise RuntimeError(f"YAML extraction failed for {Path(file_path).name}: {str(e)}")

def extract_text_from_ini(file_path: str) -> str:
    logger.debug(f"Attempting to extract structured data from INI/CFG file: '{file_path}'")
    try:
        config = configparser.ConfigParser(interpolation=None) # Disable interpolation for safety
        # Detect encoding
        with open(file_path, 'rb') as f_enc:
            raw_bytes = f_enc.read(10000)
            encoding = chardet.detect(raw_bytes)['encoding'] or 'utf-8'
        logger.debug(f"Detected INI/CFG encoding for '{file_path}': {encoding}")

        with open(file_path, 'r', encoding=encoding) as f_ini:
            config.read_file(f_ini)

        # Convert to a dictionary and then to JSON string for consistent output
        data = {section: dict(config.items(section)) for section in config.sections()}
        output = json.dumps(data, indent=2)
        logger.info(f"Successfully extracted and formatted INI/CFG data from '{file_path}'. Length: {len(output)}")
        return output
    except FileNotFoundError:
        logger.error(f"INI/CFG file not found: '{file_path}'", exc_info=True)
        raise RuntimeError(f"INI/CFG file not found: {Path(file_path).name}")
    except configparser.Error as ini_e: # Catches parsing errors, missing sections, etc.
        logger.error(f"Invalid INI/CFG format in file '{file_path}': {ini_e}", exc_info=True)
        raise ValidationError(f"Invalid INI/CFG format in {Path(file_path).name}: {str(ini_e)}")
    except IOError as ioe:
        logger.error(f"IOError reading INI/CFG file '{file_path}': {ioe}", exc_info=True)
        raise RuntimeError(f"INI/CFG file reading failed for {Path(file_path).name}: {str(ioe)}")
    except Exception as e:
        logger.error(f"Unexpected error extracting data from INI/CFG file '{file_path}': {e}", exc_info=True)
        raise RuntimeError(f"INI/CFG extraction failed for {Path(file_path).name}: {str(e)}")

# --- Legacy Format Conversion & Extraction (.doc, .ppt) ---
def _convert_with_libreoffice(file_path: str, target_format: str, tmp_dir: str) -> str:
    logger.info(f"Attempting conversion of '{Path(file_path).name}' to {target_format} using LibreOffice in temp dir '{tmp_dir}'...")
    try:
        env = os.environ.copy()
        env['HOME'] = tmp_dir # Set HOME for LibreOffice profile to avoid user profile issues

        # Check if LibreOffice is installed and executable
        try:
            subprocess.run(['which', 'libreoffice'], capture_output=True, check=True)
            logger.debug("LibreOffice command found.")
        except (FileNotFoundError, subprocess.CalledProcessError):
            logger.error("LibreOffice command ('libreoffice') not found in PATH. Please ensure it is installed.")
            raise RuntimeError("LibreOffice not found. Cannot convert legacy office formats.")

        cmd = [
            "libreoffice", "--headless", "--invisible",
            "--convert-to", target_format,
            "--outdir", tmp_dir,
            file_path
        ]
        logger.debug(f"Executing LibreOffice command: {' '.join(cmd)}")

        process = subprocess.run(
            cmd,
            check=True, # Raises CalledProcessError for non-zero exit codes
            timeout=120, # 2-minute timeout for conversion
            env=env,
            capture_output=True # Capture stdout and stderr
        )

        converted_file_name = Path(file_path).stem + f".{target_format}"
        converted_file_path = os.path.join(tmp_dir, converted_file_name)

        if not os.path.exists(converted_file_path):
            logger.error(f"LibreOffice conversion failed: Output file '{converted_file_path}' not found. Stdout: {process.stdout.decode(errors='ignore')}, Stderr: {process.stderr.decode(errors='ignore')}")
            raise RuntimeError(f"Converted {target_format} file not found after LibreOffice execution.")

        logger.info(f"LibreOffice conversion successful: '{converted_file_path}'")
        return converted_file_path

    except subprocess.TimeoutExpired:
        logger.error(f"LibreOffice conversion timed out for '{file_path}' after 120 seconds.", exc_info=True)
        raise RuntimeError(f"LibreOffice conversion timed out for {Path(file_path).name}.")
    except subprocess.CalledProcessError as e:
        stderr_output = e.stderr.decode(errors='ignore') if e.stderr else "N/A"
        stdout_output = e.stdout.decode(errors='ignore') if e.stdout else "N/A"
        logger.error(f"LibreOffice conversion failed for '{file_path}'. Return code: {e.returncode}. Stderr: {stderr_output}. Stdout: {stdout_output}", exc_info=True)
        raise RuntimeError(f"LibreOffice conversion failed for {Path(file_path).name}: {stderr_output}")
    except Exception as e: # Catch any other unexpected errors
        logger.error(f"Unexpected error during LibreOffice conversion of '{file_path}': {e}", exc_info=True)
        raise RuntimeError(f"Unexpected error during LibreOffice conversion for {Path(file_path).name}: {str(e)}")


def extract_text_from_doc(file_path: str, user: Any = None) -> str:
    logger.debug(f"Attempting to extract text from DOC file: '{file_path}'")
    original_file_name = Path(file_path).name
    if platform.system() == "Windows":
        logger.warning("Windows .doc file conversion using LibreOffice or antiword is not directly supported here. Consider alternative methods for Windows.")
        # For now, returning an informative message. Could raise an error or try Windows-specific tools if available.
        return f"Windows .doc extraction for '{original_file_name}' is not currently supported via this method."

    with tempfile.TemporaryDirectory() as tmp_dir:
        try:
            logger.info(f"Attempting .doc to .docx conversion using LibreOffice for '{original_file_name}'...")
            converted_docx_path = _convert_with_libreoffice(file_path, "docx", tmp_dir)
            return extract_text_from_docx(converted_docx_path, user)
        except RuntimeError as libre_err: # Catch errors from _convert_with_libreoffice
            logger.warning(f"LibreOffice conversion from .doc to .docx failed for '{original_file_name}': {libre_err}. Attempting fallback with antiword.")
            try:
                # Check for antiword
                try:
                    subprocess.run(['which', 'antiword'], capture_output=True, check=True)
                    logger.debug("antiword command found.")
                except (FileNotFoundError, subprocess.CalledProcessError):
                    logger.error("antiword command not found. Cannot use antiword fallback for .doc extraction.")
                    raise RuntimeError(f"DOC extraction failed for '{original_file_name}': LibreOffice failed and antiword is not installed.")

                logger.info(f"Attempting .doc extraction using antiword for '{original_file_name}'...")
                result = subprocess.run(['antiword', file_path], capture_output=True, text=True, check=True, timeout=60)
                text = result.stdout.strip()
                if text:
                    logger.info(f"antiword successfully extracted text from '{original_file_name}'.")
                    return text
                else:
                    logger.warning(f"antiword extracted no text from '{original_file_name}'.")
                    return "" # Return empty if antiword yields nothing
            except subprocess.TimeoutExpired:
                logger.error(f"antiword processing timed out for '{original_file_name}'.", exc_info=True)
                raise RuntimeError(f"DOC extraction failed for '{original_file_name}': antiword timed out.")
            except subprocess.CalledProcessError as aw_e:
                logger.error(f"antiword failed for '{original_file_name}'. Return code: {aw_e.returncode}. Stderr: {aw_e.stderr}", exc_info=True)
                raise RuntimeError(f"DOC extraction failed for '{original_file_name}': antiword processing error ({aw_e.stderr}).")
            except Exception as fallback_e: # Catch any other error during fallback
                logger.error(f"Fallback antiword extraction failed for '{original_file_name}': {fallback_e}", exc_info=True)
                raise RuntimeError(f"DOC extraction failed for '{original_file_name}': Both LibreOffice and antiword fallback failed ({fallback_e}).")
        except Exception as e: # Catch all other unexpected errors
             logger.error(f"Unexpected error during .doc extraction for '{original_file_name}': {e}", exc_info=True)
             raise RuntimeError(f"Unexpected error during .doc extraction for '{original_file_name}'.")


def extract_text_from_ppt(file_path: str, user: Any = None) -> str:
    logger.debug(f"Attempting to extract text from PPT file: '{file_path}'")
    original_file_name = Path(file_path).name
    if platform.system() == "Windows":
        logger.warning("Windows .ppt file conversion using LibreOffice is not directly supported here. Consider alternative methods for Windows.")
        return f"Windows .ppt extraction for '{original_file_name}' is not currently supported via this method."

    with tempfile.TemporaryDirectory() as tmp_dir:
        try:
            logger.info(f"Attempting .ppt to .pptx conversion using LibreOffice for '{original_file_name}'...")
            converted_pptx_path = _convert_with_libreoffice(file_path, "pptx", tmp_dir)
            return extract_text_from_pptx(converted_pptx_path, user)
        except RuntimeError as libre_err: # Catch errors from _convert_with_libreoffice
            logger.error(f"LibreOffice conversion from .ppt to .pptx failed for '{original_file_name}': {libre_err}", exc_info=True)
            # No common CLI fallback for .ppt as readily available as antiword for .doc
            raise RuntimeError(f"PPT extraction failed for '{original_file_name}': LibreOffice conversion failed ({libre_err}).")
        except Exception as e: # Catch all other unexpected errors
             logger.error(f"Unexpected error during .ppt extraction for '{original_file_name}': {e}", exc_info=True)
             raise RuntimeError(f"Unexpected error during .ppt extraction for '{original_file_name}'.")

# --- Other Text Extraction Methods (.rtf, .odt, .epub, .tex, .msg) ---

def extract_text_from_rtf(file_path: str) -> str:
    logger.debug(f"Attempting to extract text from RTF file: '{file_path}'")
    original_file_name = Path(file_path).name
    try:
        # Try unrtf utility
        try:
            subprocess.run(['which', 'unrtf'], capture_output=True, check=True)
            logger.debug("unrtf command found. Using unrtf for RTF extraction.")
            result = subprocess.run(['unrtf', '--text', file_path], capture_output=True, text=True, check=True, timeout=60)
            text = result.stdout.strip()
            logger.info(f"Successfully extracted {len(text)} characters from RTF '{original_file_name}' using unrtf.")
            return text
        except (FileNotFoundError, subprocess.CalledProcessError) as unrtf_err: # unrtf not found or failed
            logger.warning(f"unrtf not available or failed for '{original_file_name}' ({unrtf_err}). Falling back to basic regex RTF extraction.")
            # Fallback to basic regex RTF extraction
            try:
                with open(file_path, 'r', encoding='latin-1', errors='ignore') as f: # Try latin-1 as common RTF encoding
                    data = f.read()
                # Basic regex cleanup (might not be perfect and can be slow on large files)
                text = re.sub(r'\\[a-zA-Z0-9]+\\*', '', data) # Remove control words
                text = re.sub(r'[{}]', '', text) # Remove braces
                text = re.sub(r'\\pard[\\sA-Za-z0-9]*', '\n', text) # Attempt to respect paragraphs
                text = re.sub(r'\s*\\[A-Za-z0-9]+\s*', ' ', text) # Remove other control sequences
                text = re.sub(r'\s+', ' ', text).strip() # Normalize whitespace
                logger.info(f"Successfully extracted {len(text)} characters from RTF '{original_file_name}' using regex fallback.")
                return text
            except Exception as regex_e:
                logger.error(f"Basic regex RTF extraction failed for '{original_file_name}': {regex_e}", exc_info=True)
                raise RuntimeError(f"RTF extraction failed for '{original_file_name}' (unrtf not found/failed, and regex fallback also failed): {str(regex_e)}")
    except FileNotFoundError: # For the initial open() in regex fallback
        logger.error(f"RTF file not found: '{original_file_name}'", exc_info=True)
        raise RuntimeError(f"RTF file not found: {original_file_name}")
    except Exception as e: # Catch-all for other unexpected errors
        logger.error(f"Unexpected error extracting text from RTF '{original_file_name}': {e}", exc_info=True)
        raise RuntimeError(f"RTF extraction failed for '{original_file_name}': {str(e)}")


def extract_text_from_odt(file_path: str) -> str:
    logger.debug(f"Attempting to extract text from ODT file: '{file_path}'")
    original_file_name = Path(file_path).name
    try:
        # Try odt2txt utility
        try:
            subprocess.run(['which', 'odt2txt'], capture_output=True, check=True)
            logger.debug("odt2txt command found. Using odt2txt for ODT extraction.")
            result = subprocess.run(['odt2txt', file_path], capture_output=True, text=True, check=True, timeout=60)
            text = result.stdout.strip()
            logger.info(f"Successfully extracted {len(text)} characters from ODT '{original_file_name}' using odt2txt.")
            return text
        except (FileNotFoundError, subprocess.CalledProcessError) as odt2txt_err: # odt2txt not found or failed
            logger.warning(f"odt2txt not available or failed for '{original_file_name}' ({odt2txt_err}). Falling back to manual XML parsing for ODT.")
            # Fallback to manual XML parsing
            try:
                text_parts = []
                with zipfile.ZipFile(file_path) as z:
                    if 'content.xml' not in z.namelist():
                        raise ValueError("content.xml not found in ODT file's archive.")
                    with z.open('content.xml') as f_xml:
                        xml_content = f_xml.read()

                    # Use BeautifulSoup for robust XML parsing, similar to HTML
                    soup = BeautifulSoup(xml_content, 'xml')
                    for tag_name in ['text:p', 'text:h', 'text:list-item']: # Common ODT text containers
                        for element in soup.find_all(tag_name):
                            text_parts.append(element.get_text(separator=' ', strip=True))

                text = '\n'.join(filter(None, text_parts)).strip()
                logger.info(f"Successfully extracted {len(text)} characters from ODT '{original_file_name}' using XML parsing fallback.")
                return text
            except Exception as xml_e:
                logger.error(f"Manual XML parsing for ODT '{original_file_name}' failed: {xml_e}", exc_info=True)
                raise RuntimeError(f"ODT extraction failed for '{original_file_name}' (odt2txt not found/failed, and XML parsing also failed): {str(xml_e)}")
    except FileNotFoundError: # For the initial zipfile.ZipFile or subprocess.run
        logger.error(f"ODT file not found: '{original_file_name}'", exc_info=True)
        raise RuntimeError(f"ODT file not found: {original_file_name}")
    except Exception as e: # Catch-all for other unexpected errors
        logger.error(f"Unexpected error extracting text from ODT '{original_file_name}': {e}", exc_info=True)
        raise RuntimeError(f"ODT extraction failed for '{original_file_name}': {str(e)}")


def extract_text_from_epub(file_path: str) -> str:
    logger.debug(f"Attempting to extract text from EPUB file: '{file_path}'")
    original_file_name = Path(file_path).name
    txt_path_for_ebook_convert: Optional[str] = None

    try:
        # Try ebook-convert from Calibre
        try:
            subprocess.run(['which', 'ebook-convert'], capture_output=True, check=True)
            logger.debug("ebook-convert (Calibre) command found. Using it for EPUB extraction.")

            with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp_txt:
                txt_path_for_ebook_convert = tmp_txt.name

            # ebook-convert <input_file> <output_file.txt>
            cmd = ['ebook-convert', file_path, txt_path_for_ebook_convert]
            logger.debug(f"Executing ebook-convert command: {' '.join(cmd)}")
            subprocess.run(cmd, check=True, timeout=120, capture_output=True) # 2 min timeout

            with open(txt_path_for_ebook_convert, 'r', encoding='utf-8') as f_text:
                text = f_text.read().strip()
            logger.info(f"Successfully extracted {len(text)} characters from EPUB '{original_file_name}' using ebook-convert.")
            return text
        except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired) as ebook_convert_err:
            logger.warning(f"ebook-convert (Calibre) not available or failed for '{original_file_name}' ({ebook_convert_err}). Falling back to manual HTML parsing for EPUB.")
            # Fallback to manual HTML parsing from EPUB (which is a ZIP archive)
            try:
                text_parts = []
                with zipfile.ZipFile(file_path) as z:
                    html_files = [name for name in z.namelist() if name.lower().endswith(('.html', '.xhtml', '.htm'))]
                    if not html_files:
                        logger.warning(f"No HTML/XHTML files found within EPUB archive '{original_file_name}'.")
                        return ""

                    logger.debug(f"Found {len(html_files)} HTML/XHTML files in EPUB '{original_file_name}'. Parsing them...")
                    for html_file_name in html_files:
                        try:
                            with z.open(html_file_name) as f_html_in_zip:
                                raw_bytes = f_html_in_zip.read()
                                encoding = chardet.detect(raw_bytes)['encoding'] or 'utf-8'
                                soup = BeautifulSoup(raw_bytes.decode(encoding, errors='replace'), 'html.parser')
                                for element_type in ['script', 'style', 'head', 'meta', 'link']: # Remove non-content tags
                                    for element in soup.find_all(element_type):
                                        element.decompose()
                                text_parts.append(soup.get_text(separator='\n', strip=True))
                        except Exception as file_parse_e:
                            logger.warning(f"Could not parse file '{html_file_name}' in EPUB '{original_file_name}': {file_parse_e}", exc_info=True)
                            continue # Skip problematic files within the EPUB

                text = '\n\n'.join(filter(None, text_parts)).strip()
                logger.info(f"Successfully extracted {len(text)} characters from EPUB '{original_file_name}' using manual HTML parsing fallback.")
                return text
            except Exception as zip_e: # Errors related to zipfile or parsing all HTMLs
                logger.error(f"Manual HTML parsing for EPUB '{original_file_name}' failed: {zip_e}", exc_info=True)
                raise RuntimeError(f"EPUB extraction failed for '{original_file_name}' (ebook-convert not found/failed, and manual HTML parsing also failed): {str(zip_e)}")
    except FileNotFoundError: # For the initial zipfile.ZipFile or subprocess.run
        logger.error(f"EPUB file not found: '{original_file_name}'", exc_info=True)
        raise RuntimeError(f"EPUB file not found: {original_file_name}")
    except Exception as e: # Catch-all for other unexpected errors
        logger.error(f"Unexpected error extracting text from EPUB '{original_file_name}': {e}", exc_info=True)
        raise RuntimeError(f"EPUB extraction failed for '{original_file_name}': {str(e)}")
    finally:
        if txt_path_for_ebook_convert and os.path.exists(txt_path_for_ebook_convert):
            try:
                os.unlink(txt_path_for_ebook_convert)
                logger.debug(f"Cleaned up temporary file from ebook-convert: {txt_path_for_ebook_convert}")
            except OSError as unlink_e:
                logger.warning(f"Could not clean up temporary file from ebook-convert '{txt_path_for_ebook_convert}': {unlink_e}")


def extract_text_from_tex(file_path: str) -> str:
    logger.debug(f"Attempting to extract text from TeX file: '{file_path}'")
    original_file_name = Path(file_path).name
    try:
        # Try detex utility
        try:
            subprocess.run(['which', 'detex'], capture_output=True, check=True)
            logger.debug("detex command found. Using detex for TeX extraction.")
            result = subprocess.run(['detex', file_path], capture_output=True, text=True, check=True, timeout=60)
            text = result.stdout.strip()
            logger.info(f"Successfully extracted {len(text)} characters from TeX '{original_file_name}' using detex.")
            return text
        except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired) as detex_err:
            logger.warning(f"detex not available or failed for '{original_file_name}' ({detex_err}). Falling back to basic text reading.")
            # Fallback to reading as plain text if detex fails or is not found
            return extract_text_from_txt(file_path)
    except FileNotFoundError: # For extract_text_from_txt fallback
        logger.error(f"TeX file not found (during fallback txt read): '{original_file_name}'", exc_info=True)
        raise RuntimeError(f"TeX file not found: {original_file_name}")
    except Exception as e: # Catch-all for other unexpected errors
        logger.error(f"Unexpected error extracting text from TeX '{original_file_name}': {e}", exc_info=True)
        raise RuntimeError(f"TeX extraction failed for '{original_file_name}': {str(e)}")


def extract_text_from_msg(file_path: str) -> str:
    logger.debug(f"Attempting to extract text from MSG file: '{file_path}'")
    original_file_name = Path(file_path).name
    try:
        import extract_msg # Import here to make it an optional dependency
        msg = extract_msg.Message(file_path)

        # Construct text from various parts of the email
        parts = [
            f"Subject: {msg.subject}",
            f"Date: {msg.date}",
            f"From: {msg.sender}",
            f"To: {msg.to}",
            f"CC: {msg.cc}" if msg.cc else "",
            "\n--- Body ---",
            msg.body
        ]
        # Include attachment names if any
        if msg.attachments:
            parts.append("\n--- Attachments ---")
            for i, attachment in enumerate(msg.attachments):
                # Depending on attachment type, you might try to extract content too
                # For now, just listing names.
                parts.append(f"Attachment {i+1}: {getattr(attachment, 'longFilename', getattr(attachment, 'shortFilename', 'UnknownAttachmentName'))}")

        full_text = "\n".join(filter(None, parts)).strip()
        logger.info(f"Successfully extracted content from MSG file '{original_file_name}'. Length: {len(full_text)}")
        return full_text
    except ImportError:
        logger.error("The 'extract-msg' library is not installed. Cannot process .msg files.")
        raise RuntimeError("MSG extraction failed: Required library 'extract-msg' not installed. Please install it (e.g., pip install extract-msg).")
    except FileNotFoundError:
        logger.error(f"MSG file not found: '{original_file_name}'", exc_info=True)
        raise RuntimeError(f"MSG file not found: {original_file_name}")
    except Exception as e: # Catch other errors from extract_msg library
        logger.error(f"Error extracting text from MSG file '{original_file_name}': {e}", exc_info=True)
        raise RuntimeError(f"MSG extraction failed for {original_file_name}: {str(e)}")

# --- Audio/Video Extraction ---
def extract_text_from_audio(file_path: str) -> str:
    logger.debug(f"Attempting to extract text from Audio file: '{file_path}'")
    original_file_name = Path(file_path).name
    wav_path: Optional[str] = None

    try:
        logger.info(f"Converting audio '{original_file_name}' to WAV format for transcription...")
        audio = AudioSegment.from_file(file_path)

        # Create a temporary WAV file
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp_wav:
            wav_path = tmp_wav.name
        # Export with parameters suitable for speech_recognition (mono, 16kHz)
        audio.export(wav_path, format="wav", parameters=["-ac", "1", "-ar", "16000"])
        logger.debug(f"Audio '{original_file_name}' converted to WAV: '{wav_path}'")

        recognizer = sr.Recognizer()
        with sr.AudioFile(wav_path) as source:
            # recognizer.adjust_for_ambient_noise(source) # Optional: adjust for noise
            audio_data = recognizer.record(source) # Read the entire audio file

        # Attempt transcription with Whisper (local or API via library)
        text = ""
        try:
            logger.info(f"Attempting transcription of '{original_file_name}' with Whisper...")
            import whisper # Optional import
            # Consider making model size configurable ("tiny", "base", "small", "medium", "large")
            # Smaller models are faster but less accurate.
            whisper_model = whisper.load_model("base")
            result = whisper_model.transcribe(wav_path, fp16=False) # fp16=False for CPU if GPU not available/configured
            text = result["text"].strip()
            if text:
                logger.info(f"Whisper transcription successful for '{original_file_name}'. Length: {len(text)}")
                return text
            logger.info(f"Whisper transcription yielded no text for '{original_file_name}'.")
        except ImportError:
            logger.warning("Whisper library not installed. Cannot use Whisper for audio transcription.")
            # Potentially raise an error here if Whisper is a hard requirement, or allow fallback.
        except Exception as whisper_e: # Catch specific whisper errors if known, otherwise general Exception
            logger.warning(f"Whisper transcription failed for '{original_file_name}': {whisper_e}. Trying other engines if available.", exc_info=True)

        # Fallback to CMU Sphinx if Whisper fails or is not available
        if not text:
            try:
                logger.info(f"Attempting transcription of '{original_file_name}' with CMU Sphinx (offline)...")
                import pocketsphinx # Optional import, though sr.recognize_sphinx implies it's needed
                text = recognizer.recognize_sphinx(audio_data).strip()
                if text:
                    logger.info(f"CMU Sphinx transcription successful for '{original_file_name}'. Length: {len(text)}")
                    return text
                logger.info(f"CMU Sphinx transcription yielded no text for '{original_file_name}'.")
            except ImportError:
                logger.warning("Pocketsphinx library not installed or CMU Sphinx models not found. Cannot use Sphinx for audio transcription.")
            except sr.UnknownValueError:
                logger.info(f"CMU Sphinx could not understand audio from '{original_file_name}'.")
            except sr.RequestError as e_sphinx: # Errors from Sphinx engine itself
                logger.error(f"CMU Sphinx error for '{original_file_name}': {e_sphinx}", exc_info=True)

        if not text:
            logger.warning(f"All transcription attempts failed for audio file: '{original_file_name}'. No text extracted.")
            return "" # Return empty if all methods fail

        return text # Should be unreachable if logic is correct, but as a safeguard

    except FileNotFoundError:
        logger.error(f"Audio file not found: '{original_file_name}'", exc_info=True)
        raise RuntimeError(f"Audio file not found: {original_file_name}")
    except Exception as e: # Catch other errors like pydub.exceptions.CouldntDecodeError
        logger.error(f"Error extracting text from Audio '{original_file_name}': {e}", exc_info=True)
        raise RuntimeError(f"Audio transcription failed for {original_file_name}: {str(e)}")
    finally:
        if wav_path and os.path.exists(wav_path):
            try:
                os.unlink(wav_path)
                logger.debug(f"Cleaned up temporary WAV file: {wav_path}")
            except OSError as unlink_e:
                logger.warning(f"Could not clean up temporary WAV file '{wav_path}': {unlink_e}")


def extract_text_from_video(file_path: str, user: Any = None) -> str:
    logger.debug(f"Attempting to extract text/description from Video: '{file_path}'")
    original_file_name = Path(file_path).name
    audio_text = ""
    tmp_audio_path: Optional[str] = None

    # 1. Extract audio and transcribe it
    try:
        logger.info(f"Extracting audio from video: '{original_file_name}'")
        with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmp_audio_file:
            tmp_audio_path = tmp_audio_file.name

        with VideoFileClip(file_path) as video_clip:
            if video_clip.audio is None:
                 logger.warning(f"Video '{original_file_name}' has no audio track.")
            else:
                 video_clip.audio.write_audiofile(tmp_audio_path, logger=None, codec='mp3') # Specify codec
                 logger.debug(f"Audio extracted from '{original_file_name}' to: '{tmp_audio_path}'")

        if tmp_audio_path and os.path.exists(tmp_audio_path) and os.path.getsize(tmp_audio_path) > 0:
            logger.info(f"Transcribing extracted audio from '{original_file_name}' (source: '{tmp_audio_path}')...")
            audio_text = extract_text_from_audio(tmp_audio_path) # Call the refactored audio extractor
            if audio_text:
                logger.info(f"Audio transcription successful for video '{original_file_name}'. Length: {len(audio_text)} chars.")
                # For now, if audio text is found, return it. LLM fallback for video frame is separate.
                return audio_text # Return immediately if audio transcription is primary
        else:
            logger.info(f"No audio extracted or audio file was empty for video '{original_file_name}'.")

    except FileNotFoundError: # VideoFileClip can raise this if ffmpeg is missing
        logger.error(f"Video file '{original_file_name}' not found or ffmpeg might be missing.", exc_info=True)
        # Depending on strictness, could return "" or raise RuntimeError
        # For now, let it proceed to LLM fallback if user is available
    except Exception as audio_e:
        logger.warning(f"Failed to extract or transcribe audio from video '{original_file_name}': {audio_e}", exc_info=True)
    finally:
        if tmp_audio_path and os.path.exists(tmp_audio_path):
            try:
                os.unlink(tmp_audio_path)
                logger.debug(f"Cleaned up temporary audio file: {tmp_audio_path}")
            except OSError as unlink_e:
                logger.warning(f"Could not clean up temporary audio file '{tmp_audio_path}': {unlink_e}")

    # 2. Fallback to LLM Description of a video frame if audio fails or yields no text, and user context is available
    if not user:
        logger.warning(f"No user context for video '{original_file_name}', and audio processing yielded no text. Returning empty.")
        return "" # Or placeholder like "Video content (Audio failed, LLM description unavailable: no user context)"

    logger.info(f"Using LLM to describe a frame from video: '{original_file_name}' (User: {user.id})")
    try:
        llm_config = get_llm_config(user.id)
        provider_instance = llm_config.get("provider_instance")
        chat_model = llm_config.get("chat_model")
        api_key = llm_config.get("api_key")

        if not provider_instance or not chat_model:
            logger.warning(f"LLM provider or model not configured for user {user.id}. Cannot describe video frame of '{original_file_name}'.")
            return "Video content (Audio processing failed/empty, LLM frame description unavailable due to configuration)"

        if not provider_instance.is_multimodal(model=chat_model, api_key=api_key):
            logger.warning(f"LLM provider '{llm_config.get('provider_name')}' model '{chat_model}' is not multimodal. Cannot describe video frame of '{original_file_name}'.")
            return "Video content (Audio processing failed/empty, LLM frame description unavailable as model is not multimodal)"

        logger.debug(f"Extracting middle frame from video: '{original_file_name}' for LLM description.")
        with VideoFileClip(file_path) as video_clip: # Re-open video file
            frame_time = video_clip.duration / 2 if video_clip.duration and video_clip.duration > 0 else 0
            frame = video_clip.get_frame(frame_time)

        img = Image.fromarray(frame)
        buffer = BytesIO()
        img.save(buffer, format="JPEG")
        image_data = base64.b64encode(buffer.getvalue()).decode('utf-8')
        logger.debug(f"Frame extracted and encoded for video '{original_file_name}'.")

        messages = [{"role": "user", "content": [
            {"type": "text", "text": "Describe this video frame in detail. This frame is from a video file."},
            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_data}"}}
        ]}]

        response = provider_instance.get_chat_completion(messages, model=chat_model, api_key=api_key)
        
        description = _extract_answer_text(response)
            
        if not description or not description.strip():
            logger.warning(f"LLM description for video frame of '{original_file_name}' was empty.")
            return "Video content (Audio processing failed/empty, LLM frame description was empty)"

        logger.info(f"LLM description of video frame for '{original_file_name}' successful. Length: {len(description)}")
        return description.strip()

    except FileNotFoundError: # VideoFileClip can raise this
        logger.error(f"Video file '{original_file_name}' not found during LLM frame description.", exc_info=True)
        # This implies the earlier audio extraction also likely failed due to this.
        raise RuntimeError(f"Video file not found: {original_file_name}")
    except Exception as llm_e:
        logger.error(f"Error getting LLM description for video frame of '{original_file_name}': {llm_e}", exc_info=True)
        # If audio_text was found earlier, it would have been returned.
        # This path means audio_text was empty or failed.
        return f"Video content (Audio processing failed/empty, LLM frame description failed: {str(llm_e)})"

# --- Archive Extraction ---
def _extract_archive_member(archive_obj, member_info, tmp_dir: str, archive_type: str) -> Optional[str]:
    member_name_attr = 'filename' if archive_type in ['zip', 'rar', '7z'] else 'name'
    member_name = getattr(member_info, member_name_attr, None)
    if not member_name:
        #logger.warning(f"Could not get filename for member in {archive_type} archive.")
        return None

    logger.debug(f"Attempting to extract archive member: '{member_name}' (Type: {archive_type}) into '{tmp_dir}'")
    target_path: Optional[str] = None # Define target_path here for cleanup scope

    try:
        # Security: Prevent path traversal attacks.
        # Normalize the member name and ensure it's relative.
        # os.path.basename is a good start, but for archives, member names can still be tricky.
        # A more robust check might involve ensuring the resolved path is within tmp_dir.
        safe_member_name = os.path.normpath(member_name)
        if safe_member_name.startswith(('..', '/', '\\')):
            logger.warning(f"Skipping potentially unsafe member path in archive: '{member_name}' (normalized: '{safe_member_name}')")
            return None

        # Use only the final component of the path for the target name to extract flat
        target_name = os.path.basename(safe_member_name)
        if not target_name: # Skip if basename is empty (e.g. member is '/')
            logger.debug(f"Skipping member with empty basename: '{member_name}'")
            return None

        target_path = os.path.join(tmp_dir, target_name)

        # Avoid overwriting if a file with the same (basename) already exists from this archive
        if os.path.exists(target_path):
             logger.warning(f"Skipping duplicate filename during archive member extraction: '{target_name}' in '{tmp_dir}'")
             return None

        logger.debug(f"Normalized target path for member '{member_name}': '{target_path}'")

        if archive_type == 'zip':
            if member_info.is_dir():
                logger.debug(f"Skipping directory member in ZIP: '{member_name}'")
                return None
            with archive_obj.open(member_info) as source, open(target_path, "wb") as target:
                shutil.copyfileobj(source, target)
        elif archive_type == 'tar':
            if not member_info.isfile():
                logger.debug(f"Skipping non-file member in TAR: '{member_name}' (Type: {member_info.type})")
                return None
            extracted_file_obj = archive_obj.extractfile(member_info)
            if extracted_file_obj:
                with open(target_path, "wb") as target_file:
                    shutil.copyfileobj(extracted_file_obj, target_file)
                extracted_file_obj.close()
            else: # Should not happen if member_info.isfile() is true
                logger.warning(f"Could not extract file data for TAR member: '{member_name}'")
                return None
        elif archive_type == 'rar':
             if member_info.isdir():
                logger.debug(f"Skipping directory member in RAR: '{member_name}'")
                return None
             # rarfile extracts with full path by default. We want it in tmp_dir.
             # We extract to tmp_dir, and then target_path will be os.path.join(tmp_dir, member_info.filename)
             # This might create subdirectories within tmp_dir if member_info.filename is nested.
             # For simplicity of recursive extraction, we'll use this path.
             # If a flat structure is strictly needed, files would have to be moved.
             archive_obj.extract(member_info, path=tmp_dir)
             target_path = os.path.join(tmp_dir, member_info.filename) # This is the actual path after rarfile extraction
             if not os.path.exists(target_path): # Should not happen if extract succeeded
                  logger.error(f"RAR extraction failed to produce expected file: '{target_path}' for member '{member_name}'")
                  return None
        elif archive_type == '7z':
            # For 7z, _extract_archive_member is called *after* extractall if that strategy is used.
            # This function would then just confirm the path.
            # If py7zr supports single member extraction to a specific path, that would be better here.
            # Assuming 'target_path' is the full path to the already extracted member (from extractall).
            if not os.path.exists(target_path) or not os.path.isfile(target_path):
                logger.warning(f"Expected pre-extracted 7z member not found or is not a file: '{target_path}' for member '{member_name}'")
                return None
            # No actual extraction operation here for 7z if extractall was used.
            logger.debug(f"Confirmed 7z member '{member_name}' is at '{target_path}' (from prior extractall).")

        logger.info(f"Successfully extracted archive member '{member_name}' to '{target_path}'.")
        return target_path
    except Exception as e:
        logger.error(f"Failed to extract archive member '{member_name}' (Type: {archive_type}): {e}", exc_info=True)
        # Clean up potentially partially extracted file if path was determined
        if target_path and os.path.exists(target_path):
            try:
                os.unlink(target_path)
                logger.debug(f"Cleaned up partially extracted file: '{target_path}'")
            except OSError as unlink_e:
                logger.error(f"Error cleaning up partially extracted file '{target_path}': {unlink_e}", exc_info=True)
        return None


def extract_text_from_archive(file_path: str) -> str:
    logger.debug(f"Attempting to extract text from Archive file: '{file_path}'")
    original_file_name = Path(file_path).name
    extracted_texts: List[str] = []
    archive_obj: Any = None # For typing, will be specific archive object

    try:
        mime_type = detect_mime_type(file_path)
        logger.info(f"Detected archive MIME type for '{original_file_name}': {mime_type}")
        file_suffix = Path(file_path).suffix.lower()
        archive_type: Optional[str] = None
        members: List[Any] = [] # To store member info objects

        if mime_type == 'application/zip' or file_suffix == '.zip':
            archive_type = 'zip'
            logger.debug(f"Opening ZIP archive: '{original_file_name}'")
            archive_obj = zipfile.ZipFile(file_path, 'r')
            members = archive_obj.infolist()
        elif mime_type in ['application/x-tar', 'application/x-gzip', 'application/x-bzip2'] or \
             file_suffix in ['.tar', '.gz', '.tgz', '.bz2', '.tbz2']:
            archive_type = 'tar'
            logger.debug(f"Opening TAR archive: '{original_file_name}' (auto-detecting compression)")
            archive_obj = tarfile.open(file_path, mode='r:*') # Handles various compressions
            members = archive_obj.getmembers()
        elif mime_type == 'application/vnd.rar' or file_suffix == '.rar':
            archive_type = 'rar'
            logger.debug(f"Attempting to open RAR archive: '{original_file_name}'")
            try:
                import rarfile
                # rarfile.UNRAR_TOOL = "/path/to/unrar" # Optional: configure if not in PATH
                if not rarfile.is_rarfile(file_path): # Pre-check
                    raise rarfile.NotRarFile(f"File '{original_file_name}' not recognized as RAR by rarfile.")
                archive_obj = rarfile.RarFile(file_path, 'r')
                members = archive_obj.infolist()
            except ImportError:
                logger.error("The 'rarfile' library is not installed. Cannot process .rar files.")
                raise RuntimeError("RAR extraction failed: Required library 'rarfile' not installed. Please install it (e.g., pip install rarfile).")
            except rarfile.NotRarFile as e_not_rar: # More specific error
                logger.error(f"File '{original_file_name}' is not a valid RAR archive: {e_not_rar}", exc_info=True)
                raise ValidationError(f"File '{original_file_name}' is not a valid RAR archive.")
        elif mime_type == 'application/x-7z-compressed' or file_suffix == '.7z':
            archive_type = '7z'
            logger.debug(f"Attempting to open 7z archive: '{original_file_name}'")
            try:
                import py7zr
                archive_obj = py7zr.SevenZipFile(file_path, mode='r')
                members = archive_obj.list() # This returns file info objects for py7zr
            except ImportError:
                logger.error("The 'py7zr' library is not installed. Cannot process .7z files.")
                raise RuntimeError("7z extraction failed: Required library 'py7zr' not installed. Please install it (e.g., pip install py7zr).")
            except py7zr.exceptions.Bad7zFile as e_bad_7z: # Specific error for bad 7z files
                logger.error(f"File '{original_file_name}' is not a valid 7z archive or is corrupted: {e_bad_7z}", exc_info=True)
                raise ValidationError(f"File '{original_file_name}' is not a valid 7z archive or is corrupted.")
        else:
            logger.warning(f"Unsupported archive type for '{original_file_name}'. MIME: '{mime_type}', Suffix: '{file_suffix}'")
            raise ValidationError(f"Unsupported archive type: {mime_type or file_suffix}")

        logger.info(f"Opened {archive_type.upper()} archive '{original_file_name}'. Contains approx. {len(members)} members.")

        with tempfile.TemporaryDirectory() as tmp_dir:
            logger.debug(f"Using temporary directory for archive extraction: {tmp_dir}")

            # Special handling for 7z: extractall is often more reliable or simpler with py7zr
            # if it doesn't easily support extracting single members to a stream/specific path without full paths.
            is_7z_extractall = False
            if archive_type == '7z' and isinstance(archive_obj, py7zr.SevenZipFile):
                logger.info(f"Using 'extractall' for 7z archive '{original_file_name}' into '{tmp_dir}'. This might be resource-intensive for very large archives.")
                try:
                    archive_obj.extractall(path=tmp_dir)
                    is_7z_extractall = True
                    logger.debug(f"Successfully extracted all members of 7z archive '{original_file_name}' to '{tmp_dir}'.")
                except Exception as extract_all_e:
                    logger.error(f"Failed to 'extractall' from 7z archive '{original_file_name}': {extract_all_e}", exc_info=True)
                    # If extractall fails, we likely can't proceed with member processing.
                    raise RuntimeError(f"Failed to extract 7z archive '{original_file_name}': {str(extract_all_e)}")

            processed_member_paths_for_7z = set() # Only used if is_7z_extractall is true

            for member_info in members:
                # Get member name consistently
                member_name = getattr(member_info, 'filename', getattr(member_info, 'name', None))
                if not member_name:
                    logger.warning(f"Skipping member with no name in {archive_type.upper()} archive '{original_file_name}'.")
                    continue

                logger.debug(f"Processing member '{member_name}' from archive '{original_file_name}'")
                extracted_member_path: Optional[str] = None

                if is_7z_extractall: # For 7z, members are already extracted
                    # Construct the expected path of the member within tmp_dir
                    # Note: member_name from archive_obj.list() should be the relative path within the archive.
                    potential_path = os.path.join(tmp_dir, member_name)
                    if os.path.exists(potential_path) and os.path.isfile(potential_path):
                         extracted_member_path = potential_path
                         # Avoid reprocessing if somehow listed multiple times or symlinked etc.
                         if extracted_member_path in processed_member_paths_for_7z:
                              logger.debug(f"Skipping already processed path (7z extractall): '{extracted_member_path}'")
                              continue
                         processed_member_paths_for_7z.add(extracted_member_path)
                    else:
                         logger.warning(f"Expected extracted 7z member not found or is not a file: '{potential_path}' for member '{member_name}'. Skipping.")
                         continue
                else: # For other archive types, extract member individually
                    extracted_member_path = _extract_archive_member(archive_obj, member_info, tmp_dir, archive_type)

                if extracted_member_path and os.path.exists(extracted_member_path):
                    try:
                        logger.debug(f"Extracting text from successfully extracted member file: '{extracted_member_path}' (Original name: '{member_name}')")
                        # Pass user=None as this is a sub-extraction, LLM context might not apply or be too complex here.
                        # The main extract_text_from_file gets the user context.
                        member_text = extract_text_from_file(extracted_member_path, Path(member_name).name, user=None)
                        if member_text:
                            extracted_texts.append(f"--- File: {member_name} ---\n{member_text}")
                            logger.debug(f"Successfully extracted text from member '{member_name}'.")
                        else:
                            logger.debug(f"No text extracted from member '{member_name}' (empty or unsupported sub-format).")
                    except (ValidationError, RuntimeError) as e_member_extract:
                        logger.warning(f"Could not extract text from archive member '{member_name}' (Path: '{extracted_member_path}'): {e_member_extract}", exc_info=True)
                    except Exception as e_unexpected_member: # Catch any other unexpected error
                         logger.error(f"Unexpected error processing archive member '{member_name}' (Path: '{extracted_member_path}'): {e_unexpected_member}", exc_info=True)
                    finally:
                        # Clean up individual member file *only if using 7z extractall*
                        if not is_7z_extractall and os.path.exists(extracted_member_path):
                            try:
                                os.unlink(extracted_member_path)
                                logger.debug(f"Cleaned up extracted member file: '{extracted_member_path}'")
                            except OSError as e_unlink:
                                logger.warning(f"Could not clean up extracted member file '{extracted_member_path}': {e_unlink}")
                elif not is_7z_extractall : # If individual extraction failed for non-7z
                    logger.warning(f"Skipped member '{member_name}' as it was not extracted or is a directory.")

            # Temporary directory tmp_dir and its contents (like 7z extractall) are auto-cleaned by TemporaryDirectory context manager

        logger.info(f"Finished processing archive '{original_file_name}'. Extracted text from {len(extracted_texts)} members.")
        return "\n\n".join(extracted_texts).strip()

    except (zipfile.BadZipFile, tarfile.TarError, rarfile.Error, py7zr.exceptions.Bad7zFile, magic.MagicException) as archive_specific_err:
        logger.error(f"Invalid, corrupted, or specific archive processing error for '{original_file_name}': {archive_specific_err}", exc_info=True)
        raise ValidationError(f"Invalid or corrupted archive file '{original_file_name}': {str(archive_specific_err)}")
    except ImportError as import_err: # Should be caught by specific archive type blocks, but as a fallback
         logger.error(f"Missing library for archive processing of '{original_file_name}': {import_err}", exc_info=True)
         raise RuntimeError(f"Archive extraction failed for '{original_file_name}': Missing required library ({import_err.name}).")
    except FileNotFoundError:
        logger.error(f"Archive file not found: '{original_file_name}'", exc_info=True)
        raise RuntimeError(f"Archive file not found: {original_file_name}")
    except Exception as e: # Catch-all for other unexpected errors
        logger.error(f"Unexpected error extracting text from Archive '{original_file_name}': {e}", exc_info=True)
        if isinstance(e, (ValidationError, RuntimeError)): # Re-raise if already one of our types
            raise e
        raise RuntimeError(f"Archive extraction failed for '{original_file_name}': {str(e)}")
    finally:
        if archive_obj:
            try:
                archive_obj.close()
                logger.debug(f"Closed archive object for '{original_file_name}'.")
            except Exception as close_e:
                logger.warning(f"Error closing archive file '{original_file_name}': {close_e}", exc_info=True)

# --- LLM Fallback Extraction (Updated) ---
def extract_text_with_chatgpt(file_path: str, user: Any = None) -> str:
    #Logger.info(f"Using LLM fallback analysis for file: {file_path}")
    try:
        llm_config = get_llm_config(user.id if user else None)
        provider_instance = llm_config["provider_instance"]
        chat_model = llm_config["chat_model"]
        api_key = llm_config["api_key"]
        provider_name = llm_config["provider_name"]

        if not provider_instance:
             #logger.error("LLM Provider instance is None in fallback analysis.")
             return f"File: {Path(file_path).name} (LLM provider configuration error)"

        #logger.debug(f"LLM Fallback using: Provider={provider_name}, Model={chat_model}")

        ext = Path(file_path).suffix.lower()
        file_name = Path(file_path).name
        mime_type = detect_mime_type(file_path)

        # Base prompt for analysis
        prompt_text = (
            f"Analyze the following file and extract its text content or provide a detailed description.\n"
            f"File Name: {file_name}\n"
            f"File Extension: {ext}\n"
            f"MIME Type: {mime_type}\n"
        )
        messages: List[Dict[str, Any]] = []
        can_process_directly = False

        # Handle images (requires multimodal model)
        if mime_type.startswith("image/"):
            #logger.debug("File is an image. Checking for multimodal LLM.")
            if provider_instance.is_multimodal(model=chat_model, api_key=api_key):
                #Logger.info(f"Using multimodal model {chat_model} to process image.")
                with open(file_path, "rb") as img_file:
                    encoded = base64.b64encode(img_file.read()).decode("utf-8")
                messages = [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt_text + "Describe the image content."},
                            {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{encoded}"}}
                        ]
                    }
                ]
                can_process_directly = True
            else:
                #logger.warning(f"Provider {provider_name} model {chat_model} is not multimodal. Cannot process image directly.")
                return f"Image file: {file_name} (LLM cannot process images)"

        # Handle text-based files (prefer returning the raw content for chunking)
        elif mime_type.startswith("text/") or mime_type in ["application/json", "application/xml", "application/x-yaml", "application/javascript", "application/x-python"] or ext in {'.py', '.js', '.java', '.c', '.cpp', '.sh', '.md', '.log', '.sql'}:
            # Directly return the raw text (bounded) so downstream chunking keeps the code structure
            try:
                raw_text = _read_text_fallback(file_path, file_name)
                if raw_text:
                    return raw_text
                # If we somehow cannot read raw text, fall back to a minimal LLM description
                with open(file_path, "rb") as f:
                    raw_bytes = f.read(20000)  # Read first 20KB for context
                    encoding = chardet.detect(raw_bytes)['encoding'] or "utf-8"
                    content_snippet = raw_bytes.decode(encoding, errors="ignore")
                prompt_text += (
                    "\n\nContent Snippet (first 20KB):\n```\n"
                    f"{content_snippet}\n```\n\nExtract the full text content if possible, or summarize."
                )
                messages = [{"role": "user", "content": prompt_text}]
                can_process_directly = True
            except Exception as read_e:
                #logger.error(f"Failed to read text-based file {file_path} for LLM fallback: {read_e}")
                return f"Text file: {file_name} (Error reading content for LLM analysis)"

        # Handle other known types (PDF, DOCX etc.) - provide metadata
        elif mime_type in ["application/pdf", "application/msword", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
             #logger.debug("File is a document type (PDF/DOCX/PPTX). Providing metadata.")
             prompt_text += "\nThis is a document file. Describe its likely content based on the metadata."
             messages = [{"role": "user", "content": prompt_text}]
             can_process_directly = True

        # Handle Audio/Video - provide metadata
        elif mime_type.startswith("audio/") or mime_type.startswith("video/"):
             #logger.debug("File is audio/video. Providing metadata.")
             media_type = "audio" if mime_type.startswith("audio/") else "video"
             prompt_text += f"\nThis is an {media_type} file. Describe its likely purpose or content based on the metadata."
             messages = [{"role": "user", "content": prompt_text}]
             can_process_directly = True

        # Handle Archives - provide metadata (actual content extracted elsewhere)
        elif mime_type in ["application/zip", "application/x-tar", "application/gzip", "application/vnd.rar", "application/x-7z-compressed"]:
             #logger.debug("File is an archive. Providing metadata.")
             prompt_text += "\nThis is an archive file. Describe its likely contents based on the metadata."
             messages = [{"role": "user", "content": prompt_text}]
             can_process_directly = True

        # Default fallback for unknown binary types
        else:
            #logger.debug("File is an unknown binary type. Providing metadata.")
            prompt_text += "\nThis is an unknown binary file type. Describe its likely purpose based on the metadata."
            messages = [{"role": "user", "content": prompt_text}]
            can_process_directly = True

        # Make the LLM call if possible
        if can_process_directly and messages:
            #Logger.info(f"Sending request to LLM ({provider_name}, {chat_model}) for fallback analysis.")
            response = provider_instance.get_chat_completion(messages, model=chat_model, api_key=api_key)
            
            result = _extract_answer_text(response)
                
            #Logger.info(f"LLM fallback analysis completed for {file_name}.")
            return result.strip()
        else:
            # Should not happen with current logic, but as a safeguard
            #logger.error(f"Could not prepare LLM request for file {file_path}")
            return f"File: {file_name} (Could not be analyzed by LLM)"

    except Exception as e:
        #logger.error(f"Error in LLM fallback extraction for {file_path}: {e}")
        # Don't raise ValidationError here, return an error message
        return f"File analysis failed: {str(e)}"

# --- Master Extraction Function (Updated) ---
def extract_text_from_file(file_path: str, original_file_name: str, user: Any = None) -> str:
    logger.info(f"Starting text extraction for file: '{original_file_name}' (Path: '{file_path}', User: {user.id if user else 'N/A'})")
    ext = Path(original_file_name).suffix.lower()
    extracted_data: str = "" # Initialize as empty string

    # Centralized mapping of extensions to their primary extraction functions
    # Functions requiring 'user' context (for LLM calls) are noted implicitly by their signature if needed
    extraction_methods: Dict[str, Callable[..., str]] = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".png": extract_text_from_image, ".jpg": extract_text_from_image, ".jpeg": extract_text_from_image,
        ".gif": extract_text_from_image, ".bmp": extract_text_from_image, ".tiff": extract_text_from_image, ".tif": extract_text_from_image,
        ".pptx": extract_text_from_pptx,
        ".xls": extract_structured_from_excel_or_csv, ".xlsx": extract_structured_from_excel_or_csv, ".csv": extract_structured_from_excel_or_csv,
        ".txt": extract_text_from_txt, ".log": extract_text_from_txt,
        ".json": extract_structured_from_json,
        ".xml": extract_text_from_xml,
        ".html": extract_text_from_html, ".htm": extract_text_from_html,
        ".md": extract_text_from_md, # Delegates to txt
        ".yaml": extract_text_from_yaml, ".yml": extract_text_from_yaml,
        ".ini": extract_text_from_ini, ".cfg": extract_text_from_ini,
        ".doc": extract_text_from_doc, # Handles conversion, may need user for LLM fallback in sub-functions
        ".ppt": extract_text_from_ppt, # Handles conversion, may need user for LLM fallback in sub-functions
        ".rtf": extract_text_from_rtf,
        ".odt": extract_text_from_odt,
        ".epub": extract_text_from_epub,
        ".tex": extract_text_from_tex,
        ".msg": extract_text_from_msg,
        ".mp4": extract_text_from_video, ".avi": extract_text_from_video, ".mov": extract_text_from_video,
        ".wmv": extract_text_from_video, ".mkv": extract_text_from_video, ".flv": extract_text_from_video,
        ".mp3": extract_text_from_audio, ".wav": extract_text_from_audio, ".m4a": extract_text_from_audio,
        ".ogg": extract_text_from_audio, ".flac": extract_text_from_audio, ".aac": extract_text_from_audio,
        ".zip": extract_text_from_archive, ".tar": extract_text_from_archive, ".gz": extract_text_from_archive,
        ".tgz": extract_text_from_archive, ".bz2": extract_text_from_archive, ".tbz2": extract_text_from_archive,
        ".rar": extract_text_from_archive, ".7z": extract_text_from_archive
    }

    method_to_try = extraction_methods.get(ext)
    extraction_attempted = False

    if method_to_try:
        logger.info(f"Attempting specific extraction method '{method_to_try.__name__}' for file extension '{ext}' on '{original_file_name}'.")
        extraction_attempted = True
        try:
            # Determine if the method needs the 'user' argument
            # This is a bit simplistic; inspect.signature would be more robust but adds complexity.
            # For now, we assume methods that might use LLMs (image, video, doc, ppt, pdf (if OCR uses LLM)) accept 'user'.
            # Simpler extractors (txt, json, etc.) might not, but will ignore extra kwargs if defined with **kwargs.
            # To be safer, only pass user if it's in a predefined list of user-aware methods.
            user_aware_methods = [
                extract_text_from_image, extract_text_from_video,
                extract_text_from_doc, extract_text_from_ppt,
                extract_text_from_pdf # PDF might use LLM for OCR fallback in some setups
            ]
            if method_to_try in user_aware_methods:
                extracted_data = method_to_try(file_path, user=user)
            elif method_to_try is extract_structured_from_excel_or_csv:
                extracted_data = method_to_try(file_path, original_file_name=original_file_name)
            else:
                extracted_data = method_to_try(file_path)

            if extracted_data:
                 logger.info(f"Successfully extracted data using '{method_to_try.__name__}' for '{original_file_name}'. Length: {len(extracted_data)}")
            else:
                 logger.warning(f"Method '{method_to_try.__name__}' for '{original_file_name}' completed but returned no data. Will attempt LLM fallback if applicable.")

        except (ValidationError, RuntimeError) as e_specific_extractor:
            logger.warning(f"Specific extraction method '{method_to_try.__name__}' failed for '{original_file_name}': {e_specific_extractor}. Attempting LLM fallback.", exc_info=True)
            extracted_data = "" # Ensure data is empty for LLM fallback
        except Exception as e_unexpected: # Catch truly unexpected errors from specific extractors
            logger.error(f"Unexpected error in specific extraction method '{method_to_try.__name__}' for '{original_file_name}': {e_unexpected}", exc_info=True)
            extracted_data = "" # Ensure data is empty for LLM fallback
    else:
        logger.warning(f"No specific extraction method found for file extension '{ext}' of file '{original_file_name}'. Proceeding directly to LLM fallback analysis.")

    # Use LLM fallback if no specific method was found, or if it failed/returned no data
    if not extracted_data.strip(): # Check if string is empty or only whitespace
        # Try a raw text read before invoking LLMs so we retain code/content structure
        raw_text_fallback = _read_text_fallback(file_path, original_file_name)
        if raw_text_fallback:
            extracted_data = raw_text_fallback
        elif not user:
            logger.warning(f"LLM fallback for '{original_file_name}' skipped: No user context provided and primary extraction failed or yielded no data.")
            # If no specific method was even attempted, and no user for LLM, then it's a failure.
            if not extraction_attempted:
                 raise RuntimeError(f"No specific extractor for '{ext}' and LLM fallback unavailable without user context for '{original_file_name}'.")
        else:
            logger.info(f"Attempting LLM fallback analysis for '{original_file_name}' (User: {user.id}).")
            try:
                extracted_data = extract_text_with_chatgpt(file_path, user)
                if not extracted_data.strip():
                    logger.warning(f"LLM fallback analysis for '{original_file_name}' also resulted in empty or whitespace-only data.")
                else:
                    logger.info(f"LLM fallback analysis successful for '{original_file_name}'. Length: {len(extracted_data)}")
            except Exception as llm_fallback_e:
                logger.error(f"LLM fallback analysis failed for '{original_file_name}': {llm_fallback_e}", exc_info=True)
                # If even LLM fallback fails, this is the point of ultimate failure for text extraction.
                raise RuntimeError(f"All text extraction methods (specific and LLM fallback) failed for '{original_file_name}': {str(llm_fallback_e)}")

    # Final check: Ensure we have some usable text (even if it's just a description from LLM)
    if not isinstance(extracted_data, str) or not extracted_data.strip():
        final_error_msg = f"No text or structured data could be extracted from the file: {original_file_name} (Extension: {ext}). All methods, including LLM fallback (if applicable), yielded no content."
        logger.error(final_error_msg)
        raise ValidationError(final_error_msg) # Use ValidationError as it often relates to file content/type

    logger.info(f"Successfully completed text/data extraction for '{original_file_name}'. Final data length: {len(extracted_data)} chars.")

    # Ensure consistent return type (string) - mainly for structured data extractors that might return dict/list internally
    if isinstance(extracted_data, (dict, list)): # Should ideally be handled by individual structured extractors returning string
        logger.warning(f"Structured data extractor for '{original_file_name}' returned non-string type. Converting to JSON string.")
        return json.dumps(extracted_data, indent=2)

    return extracted_data.strip() # Ensure leading/trailing whitespace is removed

# --- Text Splitting --- 
def split_text_into_chunks(text: str, file_ext: str = None) -> List[LangChainDocument]:
    #logger.debug(f"Splitting text ({len(text)} chars) into chunks. Chunk size: {CHUNK_SIZE}, Overlap: {CHUNK_OVERLAP}")
    if not text:
        logger.warning("Input text for splitting is empty.")
        return []
    
    # Use different strategies for different file types
    if file_ext in ['.xlsx', '.xls', '.csv']:
        # For tabular data, use larger chunks
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            length_function=len,
            is_separator_regex=False, # Use default separators
        )
    else:
        # For regular text
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            length_function=len,
            is_separator_regex=False,
        )
    split_docs = text_splitter.create_documents([text])
    logger.info(f"Text split into {len(split_docs)} chunks using {'table-aware' if file_ext in ['.xlsx', '.xls', '.csv'] else 'general'} strategy.")
    return split_docs

# --- Vector Store Insertion (Updated) ---
def insert_document_to_vectorstore(
    text: str,
    file_type: str,
    file_ext: str,
    document_id: str, # This should be the DB Document object ID
    document_name: str,
    user: Any,
    collection_name: str,
    vector_store_id: str, # This should be the DB VectorStore object ID
):
    #Logger.info(f"Starting vector store insertion for doc_id: {document_id}, user: {user.id}, collection: {collection_name}")
    try:
        logger.info(f"Starting vector store insertion for doc_id: {document_id}, user: {user.id}, collection: {collection_name}")

        # COMPREHENSIVE LOGGING - START
        log_document_processing_status(document_id, "START_INSERTION", {
            "text_length": len(text),
            "file_type": file_type,
            "file_ext": file_ext,
            "collection": collection_name,
            "vector_store_id": vector_store_id,
            "user_id": user.id
        })

        # Monitor memory at start
        monitor_memory_usage("VECTOR_INSERT_START", document_id)
        # 1. Get the correctly initialized vector store instance
        # This now uses the refactored get_llm_config implicitly
        vector_store = get_qdrant_vector_store(user, collection_name)

        # 2. Split text into chunks
        docs_to_embed = split_text_into_chunks(text, file_ext)

        if not docs_to_embed:
            logger.warning(f"Text splitting resulted in zero chunks for doc_id: {document_id}. Skipping insertion.")
            log_document_processing_status(document_id, "NO_CHUNKS_CREATED", {
                "reason": "Text splitting resulted in zero chunks"
            })
            return
        
        # COMPREHENSIVE LOGGING - CHUNKS CREATED
        log_document_processing_status(document_id, "CHUNKS_CREATED", {
            "chunk_count": len(docs_to_embed),
            "first_chunk_preview": docs_to_embed[0].page_content[:200] if docs_to_embed else "None",
            "last_chunk_preview": docs_to_embed[-1].page_content[:200] if len(docs_to_embed) > 1 else "Same as first"
        })
        
        logger.info(f"Created {len(docs_to_embed)} chunks for document {document_id}")
        
        # Monitor memory after text splitting
        monitor_memory_usage("AFTER_TEXT_SPLIT", document_id)

        # 3. Prepare metadata for each chunk and add it to the documents
        point_ids = []  # Initialize empty list
        llm_config = get_llm_config(user.id if user else None)
        embedding_model = llm_config.get("embedding_model", "unknown")
        embed_ts = datetime.utcnow().isoformat()

        for i, doc in enumerate(docs_to_embed):
            # Generate proper UUID for point ID instead of string
            chunk_uuid = str(uuid.uuid4())  # This creates a proper UUID string
            point_ids.append(chunk_uuid)

            content = doc.page_content or ""

            metadata = {
                "document_id": str(document_id),  # Link to the Document model ID
                "vector_store_id": str(vector_store_id),  # Link to the VectorStore model ID
                "tenant_id": str(user.tenant.id),
                "user_id": str(user.id),
                "document_name": document_name,
                "file_type": file_type,
                "file_ext": file_ext,
                "chunk_index": i,
                "chunk_id": f"{document_id}chunk{i}",  # Keep this in metadata for reference
                "source": f"doc_{document_id}chunk{i}",  # Example source identifier
                # --- Extended metadata ---
                "page_number": doc.metadata.get("page_number", 1),
                "section_path": doc.metadata.get("section_path", ""),
                # "language": detect(content) if content.strip() else "unknown",
                "content_sha256": hashlib.sha256(content.encode("utf-8")).hexdigest(),
                "num_tokens": len(tiktoken.get_encoding("cl100k_base").encode(content)),
                "embedding_model": embedding_model,
                "embedding_ts": embed_ts,
                "ingest_ts": embed_ts,
            }
            doc.metadata = metadata


        # Monitor memory before vector insertion
        monitor_memory_usage("BEFORE_QDRANT_INSERT", document_id)
            
        # 4. Add documents (text chunks with metadata) to the vector store
        # Langchain's QdrantVectorStore.add_documents handles embedding and batching internally.
        #Logger.info(f"Adding {len(docs_to_embed)} chunks to Qdrant collection '{collection_name}' for doc_id: {document_id}")
        batch_size = getattr(settings, 'QDRANT_BATCH_SIZE', 50)
        close_connections_before_io("Qdrant document insertion")

        # For very large documents, use smaller batches
        if len(docs_to_embed) > 100:
            batch_size = min(batch_size, 25)  # Smaller batches for large files
            logger.info(f"Processing large document with {len(docs_to_embed)} chunks in batches of {batch_size}")

        # Process in batches
        total_inserted = 0
        failed_batches = []

        for i in range(0, len(docs_to_embed), batch_size):
            batch_docs = docs_to_embed[i:i + batch_size]
            batch_ids = point_ids[i:i + batch_size]
            batch_num = i // batch_size + 1
            total_batches = (len(docs_to_embed) - 1) // batch_size + 1

            # COMPREHENSIVE LOGGING - BATCH START
            log_document_processing_status(document_id, "BATCH_START", {
                "batch_num": batch_num,
                "total_batches": total_batches,
                "batch_size": len(batch_docs),
                "start_index": i,
                "end_index": i + len(batch_docs),
                "batch_ids_sample": batch_ids[:3] if len(batch_ids) > 3 else batch_ids,
                "first_chunk_in_batch": batch_docs[0].page_content[:100] if batch_docs else "None"
            })

            retry_count = 0
            max_retries = 3
            batch_success = False

            while retry_count < max_retries and not batch_success:
                try:
                    logger.info(f"Inserting batch {batch_num}/{total_batches} (attempt {retry_count + 1}) - {len(batch_docs)} chunks")

                    # COMPREHENSIVE LOGGING - INSERTION ATTEMPT
                    log_document_processing_status(document_id, "BATCH_INSERTION_ATTEMPT", {
                        "batch_num": batch_num,
                        "attempt": retry_count + 1,
                        "chunk_count": len(batch_docs),
                        "point_ids_count": len(batch_ids)
                    })
                    
                    call_with_resilience(
                        lambda: vector_store.add_documents(batch_docs, ids=batch_ids),
                        service="qdrant_add_documents",
                        exceptions=(Exception,),
                    )
                    total_inserted += len(batch_docs)
                    batch_success = True

                    # COMPREHENSIVE LOGGING - BATCH SUCCESS
                    log_document_processing_status(document_id, "BATCH_SUCCESS", {
                        "batch_num": batch_num,
                        "inserted_count": len(batch_docs),
                        "total_inserted_so_far": total_inserted,
                        "attempt": retry_count + 1
                    })

                    logger.info(f"Successfully inserted batch {batch_num}/{total_batches} ({len(batch_docs)} chunks)")

                    # Monitor memory after each batch for large files
                    if len(docs_to_embed) > 50:
                        monitor_memory_usage(f"AFTER_BATCH_{batch_num}", document_id)
                
                except Exception as batch_e:
                    retry_count += 1

                    # COMPREHENSIVE LOGGING - BATCH FAILURE
                    log_document_processing_status(document_id, "BATCH_FAILURE", {
                        "batch_num": batch_num,
                        "attempt": retry_count,
                        "error": str(batch_e),
                        "error_type": type(batch_e).__name__,
                        "chunk_count": len(batch_docs),
                        "will_retry": retry_count < max_retries
                    })
                    
                    logger.error(f"Failed to insert batch {batch_num} (attempt {retry_count}): {batch_e}")

                    if retry_count < max_retries:
                        wait_time = 2 ** retry_count  # Exponential backoff
                        logger.info(f"Retrying batch {batch_num} in {wait_time} seconds...")
                        time.sleep(wait_time)
                    else:
                        logger.error(f"Failed to insert batch {batch_num} after {max_retries} attempts")
                        failed_batches.append({
                            'batch_num': batch_num,
                            'start_index': i,
                            'end_index': i + len(batch_docs),
                            'error': str(batch_e),
                            'chunk_count': len(batch_docs)
                        })

                        # COMPREHENSIVE LOGGING - BATCH PERMANENTLY FAILED
                        log_document_processing_status(document_id, "BATCH_PERMANENTLY_FAILED", {
                            "batch_num": batch_num,
                            "final_error": str(batch_e),
                            "chunks_lost": len(batch_docs),
                            "start_index": i,
                            "end_index": i + len(batch_docs)
                        })

            # Small delay between batches to prevent overwhelming Qdrant
            if batch_num < total_batches:
                time.sleep(0.1)


        # Monitor memory after all insertions
        monitor_memory_usage("AFTER_QDRANT_INSERT", document_id)

        # COMPREHENSIVE LOGGING - INSERTION COMPLETE
        log_document_processing_status(document_id, "INSERTION_COMPLETE", {
            "total_chunks_attempted": len(docs_to_embed),
            "total_inserted": total_inserted,
            "failed_batches_count": len(failed_batches),
            "failed_batches": failed_batches,
            "success_rate": (total_inserted / len(docs_to_embed)) * 100 if docs_to_embed else 0,
            "missing_chunks": len(docs_to_embed) - total_inserted
        })

        # Report results
        if failed_batches:
            logger.error(f"Document {document_id}: {len(failed_batches)} batches failed to insert. Failed batches: {failed_batches}")
            raise RuntimeError(f"Vector store insertion partially failed: {len(failed_batches)} of {total_batches} batches failed")
        
        logger.info(f"Successfully inserted {total_inserted} chunks for document {document_id}")
    
    except Exception as e:
        # COMPREHENSIVE LOGGING - INSERTION FAILED
        log_document_processing_status(document_id, "INSERTION_FAILED", {
            "error": str(e),
            "error_type": type(e).__name__,
            "total_inserted": total_inserted if 'total_inserted' in locals() else 0,
            "attempted_chunks": len(docs_to_embed) if 'docs_to_embed' in locals() else 0
        })
        logger.error(f"Vector store insertion failed for document {document_id}: {e}")
        raise RuntimeError(f"Vector store insertion failed: {str(e)}")

    #     # Generate unique IDs for each Qdrant point
    #     point_ids = [str(uuid.uuid4()) for _ in docs_to_embed]

    #     # Use add_documents which handles batching
    #     vector_store.add_documents(
    #         docs_to_embed,
    #         ids=point_ids,
    #         batch_size=QDRANT_BATCH_SIZE
    #     )

    #     #Logger.info(f"Successfully inserted document chunks for doc_id: {document_id} into collection '{collection_name}'")

    #     # Optional: Force flush Qdrant client if needed (usually not necessary)
    #     # vector_store.client.conn.sync() # Example if direct client access was used

    # except Exception as e:
    #     #logger.error(f"Error inserting document (doc_id: {document_id}) into vector store '{collection_name}': {e}", exc_info=True)
    #     # Propagate the error to be handled by the calling task (_async_ingest)
    #     raise RuntimeError(f"Vector store insertion failed: {str(e)}")

# --- Document Retrieval ---
def scroll_all_points_by_vector_store_db_id(vector_store_db_id: str, user: Any, collection_name: str) -> List[Dict]:
    """
    Retrieves all Qdrant points/chunks associated with a specific VectorStore DB ID.
    The `vector_store_db_id` parameter is the ID of the Django VectorStore model instance.
    This function scrolls through all points matching `metadata.vector_store_id`.
    """
    logger.debug(f"Scrolling all points for vector_store_db_id: {vector_store_db_id} from collection: {collection_name}")
    try:
        close_connections_before_io("Qdrant scroll")
        # Get LLM config to determine the expected dimension for initializing the client,
        # even though embeddings aren't directly used for retrieval by ID here.
        # This ensures the collection is accessed/initialized correctly.
        llm_config = get_llm_config(user.id if user else None)
        if not llm_config['embedding_dimension']: # Should be caught by get_llm_config
            raise ValueError("Embedding dimension is unknown, cannot safely access Qdrant collection.")

        # Initialize Qdrant client via initialize_qdrant_collection to ensure collection exists
        # and its dimension is validated before attempting to retrieve.
        client = initialize_qdrant_collection(collection_name, llm_config['embedding_dimension'])

        # Determine accessible document IDs for this user and vector store
        accessible_doc_ids = get_accessible_document_ids(user, vector_store_db_id)
        if not accessible_doc_ids:
            logger.info(f"No accessible documents found for user {user.id} in vector store {vector_store_db_id}.")
            return []

        # Filter by tenant and accessible document IDs
        accessible_doc_ids_str = [str(doc_id) for doc_id in accessible_doc_ids]
        filter_condition = Filter(must=[
            FieldCondition(key="metadata.tenant_id", match=MatchValue(value=str(user.tenant.id))),
            FieldCondition(key="metadata.document_id", match=MatchAny(any=accessible_doc_ids_str))
        ])

        logger.debug(f"Qdrant scroll request filter for vector_store_db_id {vector_store_db_id}: {filter_condition.model_dump_json()}")

        # Use scroll API for potentially large results
        results = []
        offset = None
        limit = 100 # Scroll in batches
        while True:
            scroll_result, next_offset = call_with_resilience(
                lambda: client.scroll(
                    collection_name=collection_name,
                    scroll_filter=filter_condition,
                    limit=limit,
                    offset=offset,
                    with_payload=True,
                    with_vectors=False, # Don't need vectors for retrieval by ID
                ),
                service="qdrant_scroll",
                exceptions=(Exception,),
            )
            #logger.debug(f"Scroll result for vector_store_id {vector_id} (first page): {scroll_result}")
            results.extend(scroll_result)
            offset = next_offset
            if offset is None:
                break # No more results

        # Format results
        documents = [
            {
                "id": point.id, # Qdrant point ID
                "payload": point.payload,
            }
            for point in results
        ]
        logger.info(f"Retrieved {len(documents)} points (chunks) for vector_store_db_id {vector_store_db_id} from collection '{collection_name}'.")
        return documents
    except Exception as e:
        logger.error(f"Error retrieving document chunks by vector_store_db_id {vector_store_db_id} from '{collection_name}': {e}", exc_info=True)
        raise RuntimeError(f"Document chunk retrieval by vector_store_db_id failed: {str(e)}")

# --- Document Deletion --- 
def delete_points_by_document_db_id(document_db_id: str, user: Any, collection_name: str):
    """
    Deletes Qdrant points/chunks associated with a specific Document DB ID.
    The `document_db_id` parameter is the ID of the Django Document model instance.
    This function deletes points matching `metadata.document_id`.
    """
    logger.info(f"Attempting to delete document chunks by Document DB ID: {document_db_id} from collection: {collection_name}")
    try:
        close_connections_before_io("Qdrant delete")
        # Get LLM config to determine the expected dimension for initializing the client.
        llm_config = get_llm_config(user.id if user else None)
        if not llm_config['embedding_dimension']: # Should be caught by get_llm_config
            raise ValueError("Embedding dimension is unknown, cannot safely access Qdrant collection for deletion.")

        # Initialize Qdrant client via initialize_qdrant_collection to ensure collection exists
        # and its dimension is validated before attempting deletion.
        client = initialize_qdrant_collection(collection_name, llm_config['embedding_dimension'])

        # Delete points matching the Document DB ID in metadata
        points_selector = Filter(must=[
            FieldCondition(key="metadata.document_id", match=MatchValue(value=str(document_db_id))),
            FieldCondition(key="metadata.tenant_id", match=MatchValue(value=str(user.tenant.id))) # Important for tenancy
        ])
        logger.debug(f"Qdrant delete request for Document DB ID {document_db_id} using filter: {points_selector.model_dump_json()}")

        delete_result = call_with_resilience(
            lambda: client.delete(
                collection_name=collection_name,
                points_selector=points_selector,
                wait=True, # Wait for operation to complete
            ),
            service="qdrant_delete_points",
            exceptions=(Exception,),
        )
        logger.info(f"Qdrant deletion result for Document DB ID {document_db_id}: {delete_result}")
        logger.info(f"Successfully submitted deletion request for points (chunks) related to Document DB ID {document_db_id} from collection {collection_name}")
    except Exception as e:
        logger.error(f"Error deleting document chunks by Document DB ID {document_db_id} from collection '{collection_name}': {e}", exc_info=True)
        raise RuntimeError(f"Document chunk deletion by Document DB ID failed: {str(e)}")

# --- Hybrid Search (BM25 + Vector) ---
from .search import hybrid_search_all_files

# --- External Web Search Utility ---
def perform_web_search(query: str, max_results: int = 5) -> Tuple[str, List[str]]:
    """Return summarized text and source URLs using a DuckDuckGo search.

    Utilizes the ``ddgs`` package's ``text`` method to retrieve web results. The
    method yields dictionaries with ``title``, ``href``, and ``body`` keys. The
    collected snippets are concatenated until ``max_results`` items are
    gathered. If the library is unavailable, an empty result is returned and a
    warning is logged.
    """

    if DDGS is None:
        logger.warning("DDGS package is not installed; skipping web search")
        return "", []

    snippets: List[str] = []
    sources: List[str] = []

    try:
        with DDGS() as ddgs:
            for result in ddgs.text(
                query,
                region="us-en",
                safesearch="moderate",
                timelimit=None,
                max_results=max_results,
                page=1,
                backend="auto",
            ):
                title = result.get("title")
                body = result.get("body")
                if title or body:
                    snippets.append(" ".join(filter(None, [title, body])))
                href = result.get("href")
                if href:
                    sources.append(href)
                if len(snippets) >= max_results:
                    break

        return "\n".join(snippets[:max_results]), sources[:max_results]
    except Exception as e:
        logger.warning(f"Web search failed for query '{query}': {e}")
        return "", []

# --- Tracing Utilities ---
def trace_retrieval(question: str, docs: List[LangChainDocument]) -> None:
    """Basic tracing for retrieval stage."""
    try:
        ids = [d.metadata.get("document_id") for d in docs]
        logger.info(f"[TRACE] question='{question}' | retrieved={len(docs)} | doc_ids={ids}")
    except Exception as e:
        logger.debug(f"Trace logging failed: {e}")


# --- RAG (Retrieval-Augmented Generation) (Updated) ---
def _extract_answer_text(response: Any) -> str:
    """
    Normalize provider responses (NormalizedChatResponse, Chat Completions, Responses API, or raw strings)
    into a text answer.
    """
    if response is None:
        return ""

    # New standardized format
    from .llm_providers import NormalizedChatResponse
    if isinstance(response, NormalizedChatResponse):
        return response.content or ""

    # OpenAI Chat Completions style (Legacy support)
    if hasattr(response, "choices") and getattr(response, "choices"):
        first_choice = response.choices[0]
        message = getattr(first_choice, "message", None)
        if message and hasattr(message, "content"):
            return message.content
        content = getattr(first_choice, "content", None)
        if content:
            return content

    # OpenAI Responses API convenience property
    if hasattr(response, "output_text"):
        return getattr(response, "output_text") or ""

    # OpenAI Responses API structured output
    if hasattr(response, "output"):
        output = getattr(response, "output") or []
        if isinstance(output, list) and output:
            first = output[0]
            content = getattr(first, "content", None)
            if isinstance(content, list) and content:
                text_part = getattr(content[0], "text", None)
                if text_part and hasattr(text_part, "value"):
                    return text_part.value
                if text_part:
                    return str(text_part)

    # Dict-like fallbacks
    if isinstance(response, dict):
        if "answer" in response:
            return str(response.get("answer") or "")
        if "content" in response:
            return str(response.get("content") or "")

    # String or other objects
    return str(response)


def _get_json_format_for_provider(provider_name: str) -> Dict[str, Any]:
    """
    Get the appropriate JSON format parameter for the given provider.
    
    Args:
        provider_name: Name of the LLM provider (e.g., "OpenAI", "Ollama")
        
    Returns:
        Dictionary with format parameter key and value for the provider
    """
    provider_name_lower = provider_name.lower() if provider_name else ""
    
    if provider_name_lower == "openai":
        return {"response_format": {"type": "json_object"}}
    elif provider_name_lower == "ollama":
        # Ollama JSON schema for structured output
        json_schema = {
            "type": "object",
            "properties": {
                "answer": {
                    "type": "string",
                    "description": "The answer text"
                },
                "used_document_ids": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "List of document IDs actually used"
                }
            },
            "required": ["answer", "used_document_ids"]
        }
        return {"format": json_schema}
    else:
        # For other providers, return empty dict (no structured output support)
        return {}


def _extract_used_document_ids(response: Any, answer_text: str, tenant: Optional[Any] = None, user: Optional[Any] = None) -> List[str]:
    """
    Extract used_document_ids from LLM response.
    
    Tries multiple strategies:
    1. Parse JSON from response object or answer_text
    2. Extract used_document_ids from parsed JSON
    3. Fallback: Regex scan for [document_id=...] patterns in answer_text
    4. Deduplicate and return list
    5. Validate that extracted IDs exist in database (if tenant provided)
    
    Args:
        response: LLM response object (may contain JSON)
        answer_text: Extracted answer text string
        tenant: Optional Tenant instance to validate document IDs exist
        user: Optional User instance to validate document IDs exist
        
    Returns:
        List of document IDs that were actually used by the LLM and exist in the database
    """
    import json
    import re
    
    used_doc_ids = []
    
    # Strategy 1: Try to parse JSON from response or answer_text
    json_data = None
    
    # Try parsing answer_text first as it might already be extracted/cleaned
    try:
        json_data = json.loads(answer_text)
    except (json.JSONDecodeError, TypeError):
        # Try to extract JSON from text if it's embedded
        json_match = re.search(r'\{[^{}]*"used_document_ids"[^{}]*\}', answer_text, re.DOTALL)
        if json_match:
            try:
                json_data = json.loads(json_match.group(0))
            except (json.JSONDecodeError, TypeError):
                pass
    
    # If not found in answer_text, try the raw response content (via _extract_answer_text logic)
    if json_data is None:
        raw_text = _extract_answer_text(response)
        if raw_text and raw_text != answer_text:
            try:
                json_data = json.loads(raw_text)
            except (json.JSONDecodeError, TypeError):
                pass
    
    # Strategy 2: Already covered by Strategy 1 unified logic
    
    # Strategy 3: Extract used_document_ids from parsed JSON
    if json_data and isinstance(json_data, dict):
        doc_ids = json_data.get("used_document_ids", [])
        if isinstance(doc_ids, list):
            used_doc_ids.extend([str(doc_id) for doc_id in doc_ids if doc_id])
    
    # Strategy 4: Fallback - Regex scan for [document_id=...] patterns in answer_text
    if not used_doc_ids:
        doc_id_pattern = r'\[document_id=([^\]]+)\]'
        matches = re.findall(doc_id_pattern, answer_text)
        if matches:
            used_doc_ids.extend([str(match.strip()) for match in matches if match.strip()])
    
    # Deduplicate while preserving order
    seen = set()
    unique_doc_ids = []
    for doc_id in used_doc_ids:
        if doc_id and doc_id not in seen:
            seen.add(doc_id)
            unique_doc_ids.append(doc_id)
    
    # Validate document IDs exist in database if tenant is provided
    if unique_doc_ids and tenant:
        unique_doc_ids = validate_document_ids(unique_doc_ids, tenant, user)
        logger.info(f"[_extract_used_document_ids] Validated {len(unique_doc_ids)} document IDs (filtered out deleted ones)")
    
    return unique_doc_ids


def validate_document_ids(document_ids: List[str], tenant: Any, user: Optional[Any] = None) -> List[str]:
    """
    Validate that document IDs exist in the database and filter out deleted ones.
    
    Args:
        document_ids: List of document IDs to validate
        tenant: Tenant instance to filter documents by tenant
        user: Optional User instance to filter documents by user (if None, checks all tenant documents)
    
    Returns:
        List of valid document IDs that exist in the database
    """
    if not document_ids:
        return []
    
    # Build query to check if documents exist
    query = Document.objects.filter(id__in=document_ids, tenant=tenant)
    if user:
        query = query.filter(user=user)
    
    # Get IDs of documents that actually exist
    existing_ids = set(query.values_list('id', flat=True))
    
    # Filter to only return IDs that exist
    valid_ids = [doc_id for doc_id in document_ids if doc_id in existing_ids]
    
    # Log if any IDs were filtered out
    invalid_ids = set(document_ids) - existing_ids
    if invalid_ids:
        logger.warning(f"[validate_document_ids] Filtered out {len(invalid_ids)} non-existent document IDs: {invalid_ids}")
    
    return valid_ids


def extract_document_ids_from_response(response_text: str, available_doc_ids: Optional[List[str]] = None, tenant: Optional[Any] = None, user: Optional[Any] = None) -> List[str]:
    """
    Extract document IDs that the LLM referenced in its response text.
    
    Parses the response text for document ID patterns like [document_id=...] or document_id=...
    and returns only those IDs that were actually mentioned by the LLM.
    If tenant is provided, validates that extracted IDs exist in the database.
    
    Args:
        response_text: The LLM response text to parse
        available_doc_ids: Optional list (for logging only, not used for filtering)
        tenant: Optional Tenant instance to validate document IDs exist
        user: Optional User instance to validate document IDs exist
    
    Returns:
        Sorted list of unique document ID strings that were referenced in the response and exist in the database
    """
    import re
    
    logger.debug(f"[extract_document_ids_from_response] Starting extraction. Response length: {len(response_text) if response_text else 0}")
    if available_doc_ids:
        logger.debug(f"[extract_document_ids_from_response] Available doc IDs (for reference only): {available_doc_ids}")
    
    if not response_text:
        logger.warning("[extract_document_ids_from_response] Response text is empty, returning empty list")
        return []
    
    # Log first 500 chars of response for debugging
    response_preview = response_text[:500] if len(response_text) > 500 else response_text
    logger.debug(f"[extract_document_ids_from_response] Response preview: {response_preview!r}")
    
    # Pattern to match [document_id=...] or document_id=... in various formats
    patterns = [
        r'\[document_id=([^\]]+)\]',  # [document_id=doc_123]
        r'document_id[=:]\s*([^\s,\]\n\)]+)',  # document_id=doc_123 or document_id: doc_123
        r'\(document_id[=:]\s*([^\s,\]\n\)]+)\)',  # (document_id=doc_123)
    ]
    
    found_ids = []
    for i, pattern in enumerate(patterns):
        matches = re.findall(pattern, response_text, re.IGNORECASE)
        if matches:
            logger.debug(f"[extract_document_ids_from_response] Pattern {i} matched {len(matches)} times: {matches}")
        for match in matches:
            # Clean up the match (remove quotes, whitespace, etc.)
            doc_id = match.strip().strip('"').strip("'").strip()
            if doc_id:
                found_ids.append(doc_id)
                logger.debug(f"[extract_document_ids_from_response] Found document ID: {doc_id}")
    
    # Remove duplicates and sort
    unique_ids = sorted(list(set(found_ids)))
    logger.info(f"[extract_document_ids_from_response] Extracted {len(unique_ids)} unique document IDs from LLM response: {unique_ids}")
    if not unique_ids:
        logger.warning(f"[extract_document_ids_from_response] No document IDs found in response. Response may not contain citations.")
        return []
    
    # Validate extracted IDs against available_doc_ids first
    # This filters out garbage extractions like "151", "152" etc. from numbered references
    validated_ids = []
    if available_doc_ids:
        # Create lookup sets for both prefixed and non-prefixed versions
        available_set = set(available_doc_ids)
        available_uuids = set()
        for aid in available_doc_ids:
            if aid.startswith('doc_'):
                available_uuids.add(aid[4:])  # Strip doc_ prefix to get UUID
            else:
                available_uuids.add(aid)
        
        for doc_id in unique_ids:
            # Check if it matches directly (with prefix)
            if doc_id in available_set:
                validated_ids.append(doc_id)
                logger.debug(f"[extract_document_ids_from_response] Matched '{doc_id}' directly to available IDs")
            # Check if it matches with doc_ prefix added
            elif f"doc_{doc_id}" in available_set:
                validated_ids.append(f"doc_{doc_id}")
                logger.debug(f"[extract_document_ids_from_response] Matched '{doc_id}' -> 'doc_{doc_id}' to available IDs")
            # Check if it matches by UUID (without prefix)
            elif doc_id.startswith('doc_') and doc_id[4:] in available_uuids:
                validated_ids.append(doc_id)
                logger.debug(f"[extract_document_ids_from_response] Matched '{doc_id}' via UUID portion")
            elif doc_id in available_uuids:
                validated_ids.append(f"doc_{doc_id}")
                logger.debug(f"[extract_document_ids_from_response] Matched UUID '{doc_id}' -> 'doc_{doc_id}'")
            else:
                logger.debug(f"[extract_document_ids_from_response] Filtered out '{doc_id}' - not in available_doc_ids")
        
        if validated_ids:
            logger.info(f"[extract_document_ids_from_response] Validated {len(validated_ids)} against available_doc_ids: {validated_ids}")
        else:
            logger.warning(f"[extract_document_ids_from_response] None of the extracted IDs matched available_doc_ids")
    else:
        # No available_doc_ids provided - do basic normalization only
        for doc_id in unique_ids:
            if not doc_id.startswith('doc_'):
                validated_ids.append(f"doc_{doc_id}")
            else:
                validated_ids.append(doc_id)
        logger.info(f"[extract_document_ids_from_response] No available_doc_ids provided, normalized to: {validated_ids}")
    
    # Final DB validation if tenant is provided
    if validated_ids and tenant:
        validated_ids = validate_document_ids(validated_ids, tenant, user)
        logger.info(f"[extract_document_ids_from_response] DB validated {len(validated_ids)} document IDs")
    
    return validated_ids


def remove_document_id_citations(response_text: str) -> str:
    """
    Remove document ID citations from response text while preserving the rest of the content.
    
    Removes patterns like [document_id=...], document_id=..., (document_id=...) etc.
    This cleans the content for display while keeping citations in metadata.
    
    Args:
        response_text: The response text that may contain document ID citations
    
    Returns:
        Cleaned response text with document ID citations removed
    """
    import re
    
    if not response_text:
        return response_text
    
    # Patterns to remove (same as extraction patterns)
    patterns_to_remove = [
        r'\[document_id=[^\]]+\]',  # [document_id=doc_123]
        r'\(document_id[=:]\s*[^\s,\]\n\)]+\)',  # (document_id=doc_123)
        r'document_id[=:]\s*[^\s,\]\n\)]+',  # document_id=doc_123 or document_id: doc_123
    ]
    
    cleaned_text = response_text
    for pattern in patterns_to_remove:
        cleaned_text = re.sub(pattern, '', cleaned_text, flags=re.IGNORECASE)
    
    # Clean up ONLY extra spaces that might be left behind where citations were removed
    # We do NOT want to collapse newlines as that destroys Markdown formatting
    
    # Replace multiple spaces with a single space, but preserve newlines
    # This regex looks for 2 or more spaces in a row and replaces them with 1 space
    cleaned_text = re.sub(r'[ \t]{2,}', ' ', cleaned_text)
    
    # Ensure no trailing/leading whitespace
    cleaned_text = cleaned_text.strip()
    
    logger.debug(f"[remove_document_id_citations] Removed citations. Original length: {len(response_text)}, Cleaned length: {len(cleaned_text)}")
    
    return cleaned_text


def extract_used_document_ids(documents: List[Any]) -> List[str]:
    """
    Extract unique document IDs from documents/chunks that were included in LLM context.
    NOTE: This function extracts IDs from documents, not from LLM response.
    Use extract_document_ids_from_response() to get IDs from LLM response text.
    
    Args:
        documents: List of document objects (LangChain Document, dict, or objects with metadata)
    
    Returns:
        Sorted list of unique document ID strings
    """
    doc_ids = []
    for doc in documents:
        if not doc:
            continue
        
        doc_id = None
        try:
            # Handle dict format (from _convert_docs_for_custom_context or similar)
            if isinstance(doc, dict):
                # Check if it has payload structure
                if "payload" in doc:
                    metadata = doc.get("payload", {}).get("metadata", {})
                    if isinstance(metadata, dict):
                        doc_id = metadata.get("document_id") or metadata.get("doc_id") or metadata.get("documentId")
                # Check direct metadata
                elif "metadata" in doc:
                    metadata = doc.get("metadata", {})
                    if isinstance(metadata, dict):
                        doc_id = metadata.get("document_id") or metadata.get("doc_id") or metadata.get("documentId")
            # Handle LangChain Document objects
            elif hasattr(doc, "metadata"):
                metadata = getattr(doc, "metadata", {}) or {}
                if isinstance(metadata, dict):
                    doc_id = metadata.get("document_id") or metadata.get("doc_id") or metadata.get("documentId")
                else:
                    doc_id = getattr(metadata, "document_id", None) or getattr(metadata, "doc_id", None)
        except Exception:
            # Skip documents with unparseable metadata
            continue
        
        if doc_id:
            doc_ids.append(str(doc_id))
    
    # Return sorted unique list
    return sorted(list(set(doc_ids)))


def ask_question(
    question: str,
    vector_store_id: str,
    user: Any,
    collection_name: str,
    assistant_instructions: Optional[str] = None,
    thread_history: Optional[List[Dict[str, str]]] = None,
    documents: Optional[List[Dict]] = None,
    mode: str = "document",
    metadata_filters: Optional[Dict[str, Any]] = None,
    rerank_model: Optional[str] = None,
    use_mmr: bool = False,
) -> Dict[str, Any]:
    try:
        logger.info(f"[ASK_QUESTION] Processing question for user {user.id}, vector_store_id: {vector_store_id}, mode: {mode}")

        # Get LLM configuration
        llm_config = get_llm_config(user.id)
        provider_instance = llm_config["provider_instance"]
        provider_name = llm_config["provider_name"]
        chat_model = llm_config["chat_model"]
        api_key = llm_config["api_key"]

        if not provider_instance:
            raise ValueError("Could not initialize LLM provider instance.")

        thread_context = ""
        if thread_history:
            thread_context = "\n\nPrevious conversation:\n" + "\n".join(
                f"{msg['role']}: {msg['content']}" for msg in thread_history
            )

        if mode == "normal":
            messages = [
                {"role": "system", "content": assistant_instructions or "You are a helpful assistant."}
            ]
            if thread_history:
                messages.extend(thread_history)
            messages.append({"role": "user", "content": question})
            response = provider_instance.get_chat_completion(
                messages=messages,
                model=chat_model,
                api_key=api_key,
            )
            answer = _extract_answer_text(response)
            return {
                "answer": answer,
                "sources": [],
                "used_document_ids": [],
                "search_method": "normal_chat",
                "retrieval_time_ms": 0,
                "total_chunks_retrieved": 0,
            }

        if mode == "web":
            search_text, web_sources = perform_web_search(question)
            base_prompt = (
                f"Web search results:\n{search_text}\n{thread_context}\n\nQuestion: {question}\n"
                "Provide the most relevant and up-to-date answer using the search results above."
            )
            messages = [
                {"role": "system", "content": assistant_instructions or "You are a helpful assistant leveraging web search results."},
                {"role": "user", "content": base_prompt},
            ]
            response = provider_instance.get_chat_completion(
                messages=messages,
                model=chat_model,
                api_key=api_key,
            )
            answer = _extract_answer_text(response)
            return {
                "answer": answer,
                "sources": web_sources,
                "search_method": "web_search",
                "retrieval_time_ms": 0,
                "total_chunks_retrieved": 0,
            }

        vector_store = get_qdrant_vector_store(user, collection_name)

        context = ""
        sources = []
        retrieval_time_ms = 0
        total_chunks_retrieved = 0
        relevant_chunks_found = False

        # Handle custom documents if provided
        available_doc_ids = []  # Track available doc IDs for validation
        if documents is not None:
            # Extract available document IDs from documents that will be included in context
            available_doc_ids = extract_used_document_ids(documents)
            # Combine page contents from documents
            context = "\n\n---\n\n".join([
                doc.get("payload", {}).get("page_content", "")
                for doc in documents
                if doc.get("payload", {}).get("page_content")
            ])
            sources = ["custom_documents"]
            relevant_chunks_found = bool(context.strip())
        else:
            # Get accessible document IDs for the user
            accessible_doc_ids = get_accessible_document_ids(user, vector_store_id)
            if not accessible_doc_ids:
                context = "No accessible documents found."
                sources = []
            else:
                # Single retrieval using MatchAny for all accessible document IDs
                accessible_doc_ids_str = [str(doc_id) for doc_id in accessible_doc_ids]
                must_conditions = [
                    FieldCondition(key="metadata.tenant_id", match=MatchValue(value=str(user.tenant.id))),
                    FieldCondition(key="metadata.document_id", match=MatchAny(any=accessible_doc_ids_str)),
                ]
                if metadata_filters:
                    for key, value in metadata_filters.items():
                        must_conditions.append(FieldCondition(key=f"metadata.{key}", match=MatchValue(value=value)))
                search_filter = Filter(must=must_conditions)
                logger.info(f"Starting retrieval for question: '{question}'")
                # Use the new HybridSearchRetriever
                retriever = HybridSearchRetriever(
                    vector_store=vector_store,
                    search_filter=search_filter,
                    k=HYBRID_SEARCH_TOP_K,
                    rerank_model=rerank_model or os.getenv("RERANKER_MODEL"),
                    use_mmr=use_mmr or os.getenv("USE_MMR", "false").lower() == "true"
                )
                
                logger.info(f"Starting retrieval for question: '{question}' using HybridSearchRetriever")
                start_time = time.time()
                retrieved_docs = retriever.invoke(question)
                trace_retrieval(question, retrieved_docs)
                end_time = time.time()

                retrieval_time_ms += (end_time - start_time) * 1000
                logger.info(f"-----Retrieval completed in {retrieval_time_ms:.2f} ms | Found {len(retrieved_docs)} documents using retriever-----")

                total_chunks_retrieved = len(retrieved_docs)

                if not retrieved_docs:
                    context = ""
                    sources = []
                    available_doc_ids = []
                else:
                    context = "\n\n---\n\n".join([doc.page_content for doc in retrieved_docs])
                    # Extract available document IDs from documents that were actually included in context
                    available_doc_ids = extract_used_document_ids(retrieved_docs)
                    # Sources: list of unique document_ids we searched across
                    sources = accessible_doc_ids_str
                    relevant_chunks_found = True

        # Handle case when no relevant context is found
        if not relevant_chunks_found or not context.strip():
            missing_info_prompt = (
                f"The provided documents do not contain enough information to answer the following question: '{question}'.\n"
                "Please list the specific facts, data, or details that would be required to answer this question. "
                "If possible, explain what is missing or what type of document/source would be needed."
            )
            messages = [
                {
                    "role": "system",
                    "content": assistant_instructions
                    or (
                        "You are a helpful assistant. If you cannot answer, list what information is missing."
                    ),
                },
                {"role": "user", "content": missing_info_prompt},
            ]
            response = provider_instance.get_chat_completion(
                messages=messages,
                model=chat_model,
                api_key=api_key,
            )
            
            # Extract the answer content based on provider response type
            answer = _extract_answer_text(response)
            logger.info(f"[ASK_QUESTION] Extracted answer text (length: {len(answer) if answer else 0})")
            # Extract document IDs that the LLM actually referenced in its response
            logger.info(f"[ASK_QUESTION] Available document IDs for validation: {available_doc_ids}")
            used_document_ids = extract_document_ids_from_response(answer, available_doc_ids, tenant=user.tenant, user=user)
            logger.info(f"[ASK_QUESTION] Extracted {len(used_document_ids)} used document IDs: {used_document_ids}")
            return {
                "answer": answer,
                "sources": sources,
                "used_document_ids": used_document_ids,
                "search_method": "hybrid_bm25_vector",
                "retrieval_time_ms": round(retrieval_time_ms, 2),
                "total_chunks_retrieved": total_chunks_retrieved
            }

        # Construct the main prompt with context and thread history
        base_prompt = f"""Context from documents:
        {context}

        {thread_context}

        Question: {question}

        Please provide a comprehensive answer based on the context above. When you use information from the context, cite the document ID by including [document_id=<id>] in your response where <id> is the document_id from the context headers. If the context does not contain enough information to fully answer the question, clearly list the specific facts, data, or details that are missing and would be required to answer it completely. Do not make up information not present in the context."""

        # Prepare messages for the language model
        messages = [
            {
                "role": "system",
                "content": assistant_instructions or (
                    "You are a helpful assistant answering questions based on provided document context. "
                    "When you reference information from the documents, you must cite the document ID using the format [document_id=<id>] where <id> matches the document_id from the context headers. "
                    "If the context does not contain the answer, state what is missing. "
                    "Do not make up information not present in the context. Keep responses concise and grounded in the source material."
                )
            },
            {
                "role": "user",
                "content": base_prompt
            }
        ]
        # Log message preview
        logger.debug(f"[ASK_QUESTION] Messages sent to {provider_name}: {json.dumps(messages, indent=2)[:500]}")

        # Get the answer from the language model
        response = provider_instance.get_chat_completion(
            messages=messages,
            model=chat_model,
            api_key=api_key,
        )
        
        # Extract the answer content based on provider response type
        answer = _extract_answer_text(response)

        # Post-processing: enhance answer if generic
        generic_responses = [
            "I don't know",
            "not enough information",
            "not available",
            "insufficient information",
            "cannot answer"
        ]
        if isinstance(answer, str) and any(generic in answer.lower() for generic in generic_responses):
            answer += "\n\nIf you need more details, please provide or upload documents containing the required information."

        logger.info(f"[ASK_QUESTION] Successfully generated answer for user {user.id}")
        logger.info(f"[ASK_QUESTION] Answer text length: {len(answer) if answer else 0}")

        # Extract document IDs that the LLM actually referenced in its response
        available_for_validation = available_doc_ids if 'available_doc_ids' in locals() else None
        logger.info(f"[ASK_QUESTION] Available document IDs for validation: {available_for_validation}")
        used_document_ids = extract_document_ids_from_response(answer, available_for_validation, tenant=user.tenant, user=user)
        logger.info(f"[ASK_QUESTION] Final used_document_ids: {used_document_ids}")
        
        # Remove document ID citations from the answer text for clean display
        cleaned_answer = remove_document_id_citations(answer)
        logger.info(f"[ASK_QUESTION] Cleaned answer text (removed citations). Original length: {len(answer)}, Cleaned length: {len(cleaned_answer)}")

        # Return the final response
        return {
            "answer": cleaned_answer,
            "sources": sources,
            "used_document_ids": used_document_ids,
            "search_method": "hybrid_bm25_vector",
            "retrieval_time_ms": round(retrieval_time_ms, 2),
            "total_chunks_retrieved": total_chunks_retrieved
        }

    except TimeoutException:
        logger.error(f"[ASK_QUESTION] Timeout for user {user.id}, question: {question[:100]}")
        raise LLMServiceUnavailableError("Request timed out. Please try a shorter query.")
    except LLMProviderError:
        # Re-raise custom provider errors to be caught by views
        raise


def insert_large_file_to_vectorstore_from_path(
    file_path: str,
    file_ext: str,
    document_id: str,
    document_name: str,
    user: Any,
    collection_name: str,
    vector_store_id: str,
) -> None:
    """
    Stream large CSV/Excel files into Qdrant without loading the full file into memory.
    """
    logger.info("Starting streaming vector ingestion for large file doc_id: %s", document_id)
    close_connections_before_io("large file vector ingestion")

    vector_store = get_qdrant_vector_store(user, collection_name)
    llm_config = get_llm_config(user.id if user else None)
    embedding_model = llm_config.get("embedding_model", "unknown")
    embed_ts = datetime.utcnow().isoformat()
    encoding = tiktoken.get_encoding("cl100k_base")
    batch_size = getattr(settings, 'QDRANT_BATCH_SIZE', 50)
    batch_docs: List[LangChainDocument] = []
    batch_ids: List[str] = []
    chunk_index = 0

    def flush_batch() -> None:
        if not batch_docs:
            return
        call_with_resilience(
            lambda: vector_store.add_documents(batch_docs, ids=batch_ids),
            service="qdrant_add_documents_streaming",
            exceptions=(Exception,),
        )
        batch_docs.clear()
        batch_ids.clear()

    def add_docs_from_chunk(chunk_text: str) -> None:
        nonlocal chunk_index
        docs_to_embed = split_text_into_chunks(chunk_text, file_ext)
        for doc in docs_to_embed:
            content = doc.page_content or ""
            doc.metadata = {
                "document_id": str(document_id),
                "vector_store_id": str(vector_store_id),
                "tenant_id": str(user.tenant.id),
                "user_id": str(user.id),
                "document_name": document_name,
                "file_type": "file",
                "file_ext": file_ext,
                "chunk_index": chunk_index,
                "chunk_id": f"{document_id}chunk{chunk_index}",
                "source": f"doc_{document_id}chunk{chunk_index}",
                "page_number": doc.metadata.get("page_number", 1),
                "section_path": doc.metadata.get("section_path", ""),
                "content_sha256": hashlib.sha256(content.encode("utf-8")).hexdigest(),
                "num_tokens": len(encoding.encode(content)),
                "embedding_model": embedding_model,
                "embedding_ts": embed_ts,
                "ingest_ts": embed_ts,
            }
            batch_docs.append(doc)
            batch_ids.append(str(uuid.uuid4()))
            chunk_index += 1
            if len(batch_docs) >= batch_size:
                flush_batch()

    try:
        if file_ext == ".csv":
            for chunk in pd.read_csv(file_path, chunksize=CHUNK_SIZE_LARGE):
                chunk_text = chunk.to_string(index=False, max_colwidth=50)
                add_docs_from_chunk(chunk_text)
                del chunk
        elif file_ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                chunk_rows = []
                for row in sheet.iter_rows(values_only=True):
                    chunk_rows.append(row)
                    if len(chunk_rows) >= CHUNK_SIZE_LARGE:
                        add_docs_from_chunk(_format_rows_to_text(sheet_name, chunk_rows))
                        chunk_rows.clear()
                if chunk_rows:
                    add_docs_from_chunk(_format_rows_to_text(sheet_name, chunk_rows))
                    chunk_rows.clear()
            wb.close()
        elif file_ext == ".xls":
            import xlrd
            wb = xlrd.open_workbook(file_path, on_demand=True)
            for sheet in wb.sheets():
                for start_idx in range(0, sheet.nrows, CHUNK_SIZE_LARGE):
                    end_idx = min(start_idx + CHUNK_SIZE_LARGE, sheet.nrows)
                    chunk_rows = [sheet.row_values(i) for i in range(start_idx, end_idx)]
                    add_docs_from_chunk(_format_rows_to_text(sheet.name, chunk_rows))
            wb.release_resources()
        else:
            raise ValueError(f"Unsupported file extension for streaming ingestion: {file_ext}")

        flush_batch()
        logger.info("Completed streaming vector ingestion for doc_id: %s", document_id)
    except CircuitBreakerOpenError as exc:
        logger.error("Circuit breaker open during streaming ingestion: %s", exc)
        raise
    except Exception as exc:
        logger.error("Streaming vector ingestion failed for doc_id %s: %s", document_id, exc, exc_info=True)
        raise


# --- Document Enrichment & Alerting (Updated) --- 
def enrich_document(doc: Document, text: str, file_ext: str, user: Any):
    #Logger.info(f"Starting enrichment for document: {doc.id} ({doc.title})")
    try:
        llm_config = get_llm_config(user.id)
        provider_instance = llm_config["provider_instance"]
        chat_model = llm_config["chat_model"]
        api_key = llm_config["api_key"]

        if not provider_instance:
            #logger.error("Cannot enrich document: LLM provider instance is None.")
            return # Cannot proceed without LLM

        # Use a flag to track if any enrichment succeeded
        enrichment_succeeded = False

        # 1. Summarization
        try:
            #logger.debug("Generating summary...")
            summary_prompt = f"Provide a concise summary (2-3 sentences) of the following document content:\n\n{text[:MAX_TOKEN_THRESHOLD]}" # Limit context size
            response = provider_instance.get_chat_completion([{'role': 'user', 'content': summary_prompt}], model=chat_model, api_key=api_key)
            
            summary = _extract_answer_text(response)
                
            doc.summary = summary.strip()
            #logger.debug(f"Summary generated: {doc.summary}")
            enrichment_succeeded = True
        except Exception as summary_e:
            #logger.error(f"Summary generation failed for doc {doc.id}: {summary_e}")
            doc.summary = "(Summary generation failed)"

        # 2. Keyword Extraction
        try:
            #logger.debug("Extracting keywords...")
            keyword_prompt = f"Extract the 5-10 most important keywords or keyphrases from the following document content. List them separated by commas:\n\n{text[:MAX_TOKEN_THRESHOLD]}"
            response = provider_instance.get_chat_completion([{'role': 'user', 'content': keyword_prompt}], model=chat_model, api_key=api_key)
            
            keywords_str = _extract_answer_text(response)
                
            doc.keywords = [k.strip() for k in keywords_str.split(',') if k.strip() and len(k.strip()) < 100] # Add length limit
            #logger.debug(f"Keywords extracted: {doc.keywords}")
            enrichment_succeeded = True
        except Exception as keyword_e:
            #logger.error(f"Keyword extraction failed for doc {doc.id}: {keyword_e}")
            doc.keywords = []

        # 3. Sentiment Analysis (using NLTK VADER)
        analyzer = get_sentiment_analyzer()
        if analyzer is not None:
            try:
                #logger.debug("Performing sentiment analysis...")
                # Limit text size for performance
                sentiment_scores = analyzer.polarity_scores(text[:50000])
                doc.sentiment_score = sentiment_scores['compound'] # Compound score: -1 (neg) to +1 (pos)
                # Determine label based on score
                if doc.sentiment_score >= 0.05:
                    doc.sentiment_label = 'positive'
                elif doc.sentiment_score <= -0.05:
                    doc.sentiment_label = 'negative'
                else:
                    doc.sentiment_label = 'neutral'
                #logger.debug(f"Sentiment: Score={doc.sentiment_score}, Label={doc.sentiment_label}")
                enrichment_succeeded = True
            except Exception as sentiment_e:
                #logger.error(f"Sentiment analysis failed for doc {doc.id}: {sentiment_e}")
                doc.sentiment_score = None
                doc.sentiment_label = 'unknown'
        else:
            doc.sentiment_score = None
            doc.sentiment_label = 'unknown'

        # 4. Entity Recognition (using spaCy)
        try:
            #logger.debug("Performing entity recognition...")
            # Limit text size for spaCy performance
            nlp_model = get_nlp()
            spacy_doc = nlp_model(text[:50000])
            entities = defaultdict(list)
            for ent in spacy_doc.ents:
                # Filter entities (optional)
                # if ent.label_ in ['PERSON', 'ORG', 'GPE', 'DATE']:
                if len(ent.text) < 100: # Limit entity text length
                     entities[ent.label_].append(ent.text)
            # Store as JSON, ensuring unique entities per type
            doc.entities = {k: list(set(v)) for k, v in entities.items()}
            #logger.debug(f"Entities recognized: {json.dumps(doc.entities)[:200]}...") # Log snippet
            enrichment_succeeded = True
        except Exception as entity_e:
            #logger.error(f"Entity recognition failed for doc {doc.id}: {entity_e}")
            doc.entities = {}

        # Save the document only if some enrichment was attempted
        if enrichment_succeeded:
            doc.save()
            #Logger.info(f"Enrichment completed for document: {doc.id}")
        else:
            logger.warning(f"No enrichment steps succeeded for document: {doc.id}")

    except Exception as e:
        #logger.error(f"Error during document enrichment setup for doc_id {doc.id}: {e}", exc_info=True)
        # Log error but don't necessarily fail the whole process
        # Optionally update the document status to indicate enrichment failure
        pass

def detect_alerts(doc: Document, text: str):
    #Logger.info(f"Starting alert detection for document: {doc.id}")
    # Example: Detect PII (using regex or more advanced methods)
    # Consider using libraries like presidio-analyzer for better PII detection
    pii_patterns = {
        'EMAIL': r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}',
        'PHONE_US': r'\b\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}\b',
        'SSN': r'\b\d{3}-\d{2}-\d{4}\b',
        'CREDIT_CARD': r'\b(?:\d[ -]*?){13,16}\b',
        'IP_ADDRESS': r'\b(?:[0-9]{1,3}\.){3}[0-9]{1,3}\b',
        'DOB': r'\b(0[1-9]|1[0-2])[\/.-](0[1-9]|[12][0-9]|3[01])[\/.-](19|20)\d\d\b',
        'ADDRESS': r'\d+\s+[A-Za-z0-9.\s]+\s+(Street|St|Avenue|Ave|Boulevard|Blvd|Road|Rd|Lane|Ln|Drive|Dr)\b',
        'PASSPORT': r'\b[A-PR-WYa-pr-wy][1-9]\d{6}\b',
    }
    alerts_found = []
    try:
        for pii_type, pattern in pii_patterns.items():
            # Limit search text size for performance
            matches = re.findall(pattern, text[:100000])
            if matches:
                unique_matches = list(set(matches))
                #logger.warning(f"Potential PII ({pii_type}) detected in document {doc.id}: {len(unique_matches)} unique instance(s).")
                alert = DocumentAlert(
                    document=doc,
                    user=getattr(doc, 'user', None),
                    keyword=pii_type,
                    snippet=f"Potential {pii_type} found: {len(unique_matches)} unique instance(s). Matches: {', '.join(unique_matches[:3])}"
                )
                alerts_found.append(alert)

        # Example: Detect negative sentiment (using previously enriched data)
        if hasattr(doc, 'sentiment_label') and doc.sentiment_label == 'negative' and hasattr(doc, 'sentiment_score') and doc.sentiment_score is not None and doc.sentiment_score < -0.5: # Threshold
            #logger.warning(f"Strong negative sentiment detected in document {doc.id}")
            alert = DocumentAlert(
                document=doc,
                user=getattr(doc, 'user', None),
                keyword="NEGATIVE_SENTIMENT",
                snippet=f"Strong negative sentiment detected (Score: {doc.sentiment_score:.2f}). Summary: {(doc.summary or '')[:200]}"
            )
            alerts_found.append(alert)

        # Bulk create alerts
        if alerts_found:
            try:
                DocumentAlert.objects.bulk_create(alerts_found)
                #Logger.info(f"Created {len(alerts_found)} alerts for document {doc.id}")
            except Exception as e:
                logger.error(f"Failed to save alerts for document {doc.id}: {e}")

        logger.info(f"Alert detection completed for document: {doc.id}")

    except Exception as alert_e:
        logger.error(f"Error during alert detection for doc {doc.id}: {alert_e}", exc_info=True)

# --- Document Access Control (Example) ---
def check_document_access(user: Any, document_id: str) -> bool:
    #logger.debug(f"Checking access for user {user.id} to document {document_id}")
    try:
        # Ensure user is authenticated and has a tenant
        if not user or not user.is_authenticated or not hasattr(user, 'tenant') or not user.tenant:
            #logger.warning(f"Access check failed: Invalid user object for user ID {getattr(user, 'id', 'N/A')}.")
            return False

        # Check if document exists and belongs to the user's tenant first
        doc = Document.objects.filter(id=document_id, tenant=user.tenant).first()
        if not doc:
            #logger.warning(f"Access denied: Document {document_id} not found or does not belong to tenant {user.tenant.id}.")
            return False

        # 1. Check if user owns the document
        if doc.user == user:
            #Logger.info(f"User {user.id} owns document {document_id}. Access granted.")
            return True

        # 2. Check if user has explicit access via DocumentAccess model
        # Ensure the DocumentAccess entry also matches the tenant
        if DocumentAccess.objects.filter(document=doc, user=user, tenant=user.tenant).exists():
            #Logger.info(f"User {user.id} has explicit access to document {document_id}. Access granted.")
            return True

        # 3. Add other rules (e.g., group access, tenant-wide access)
        # Example: Allow access if user is an admin in the same tenant
        # if user.is_staff or user.is_superuser: # Or a custom role
        #     #Logger.info(f"Admin user {user.id} accessing document {document_id} in tenant {user.tenant.id}. Access granted.")
        #     return True

        # Example: Allow access if document is marked tenant-sharable (add is_tenant_shared field to Document model)
        # if getattr(doc, 'is_tenant_shared', False):
        #      #Logger.info(f"Document {document_id} is shared within tenant {user.tenant.id}. Access granted to user {user.id}.")
        #      return True

        logger.warning(f"Access denied for user {user.id} to document {document_id}.")
        return False
    except Exception as e:
        logger.error(f"Error checking document access for user {user.id}, doc {document_id}: {e}", exc_info=True)
        return False # Deny access on error

def get_accessible_document_ids(user: Any, vector_store_id: str) -> list:
    """
    Returns a list of document IDs the user has access to for the given vector store.
    Includes documents owned by the user and those granted via DocumentAccess.
    """
    # Documents owned by the user in this vector store
    owned_docs = set(Document.objects.filter(user=user, vector_store=vector_store_id).values_list('id', flat=True))
    # Documents granted via DocumentAccess
    access_docs = set(DocumentAccess.objects.filter(vector_store=vector_store_id, granted_by=user).values_list('document_id', flat=True))
    # Optionally, include documents granted to the user (if DocumentAccess has a 'user' field for grantee)
    # access_docs = set(DocumentAccess.objects.filter(vector_store=vector_store_id, user=user).values_list('document_id', flat=True))
    return list(owned_docs | access_docs)

def get_all_accessible_document_ids(user: Any, vector_store_ids: list) -> list:
    """
    Consolidated version of get_accessible_document_ids for multiple vector stores.
    Reduces database round-trips by querying for all IDs at once.
    """
    if not vector_store_ids:
        return []
    
    # Documents owned by the user across all specified vector stores
    owned_docs = set(Document.objects.filter(
        user=user, 
        vector_store__id__in=vector_store_ids
    ).values_list('id', flat=True))
    
    # Documents granted via DocumentAccess across all specified vector stores
    access_docs = set(DocumentAccess.objects.filter(
        vector_store__id__in=vector_store_ids, 
        granted_by=user
    ).values_list('document_id', flat=True))
    
    return [str(doc_id) for doc_id in (owned_docs | access_docs)]

def generate_thread_title(first_message_content: str, user: Any = None) -> str:
    """
    Generate a concise title for a thread based on the first user message.
    Uses LLM to create a short, descriptive title.
    """
    try:
        # Limit the input to avoid token limits
        content_preview = first_message_content[:500]
        
        # Create a prompt for title generation
        title_prompt = f"""Generate a concise, descriptive title (maximum 6-8 words) for a conversation thread based on this first message:

"{content_preview}"

The title should:
- Be clear and descriptive
- Capture the main topic or question
- Be professional and concise
- Not include quotes or special characters

Title:"""

        llm_config = get_llm_config(user.id if user else None)
        provider_instance = llm_config["provider_instance"]
        chat_model = llm_config.get("chat_model")
        api_key = llm_config.get("api_key")

        response = provider_instance.get_chat_completion(
            [{"role": "user", "content": title_prompt}],
            model=chat_model,
            api_key=api_key
        )
        
        # Normalize answer extraction to support Chat Completions and Responses API outputs.
        title = _extract_answer_text(response).strip() or _generate_fallback_title(content_preview)
        title = title.strip('"\'')
        if len(title) > 255:
            title = title[:252] + "..."
        return title
            
    except Exception as e:
        logger.warning(f"Failed to generate thread title for user {user.id if user else 'N/A'}: {e}")
        return _generate_fallback_title(first_message_content)

def _generate_fallback_title(content: str) -> str:
    """
    Generate a simple fallback title when LLM generation fails.
    """
    # Take first few words and clean them up
    words = content.split()[:6]  # Take first 6 words
    title = " ".join(words)
    
    # Remove common punctuation and clean up
    title = title.strip('.,!?;:"')
    
    # Ensure it's not too long
    if len(title) > 50:
        title = title[:47] + "..."
    
    # If empty or too short, use a default
    if len(title.strip()) < 3:
        title = "New Conversation"
    
    return title
